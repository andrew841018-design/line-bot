"""
LINE → FastAPI webhook → Gemini → LINE

回應觸發條件：
1. Explicit：@mention bot，或用 /ai /問 /ask 前綴 → 立刻回
2. Implicit：任何文字訊息都進 burst_filter 佇列 → 30 秒 debounce →
   規則 + 啟發式 + Gemini classifier 判定「值得回」才回（絕大多數情況不回）

訊息類型：
- 文字：走上面兩條路徑之一；URL 會由 burst 啟發式觸發，Gemini 自帶 url_context 會讀網頁
- 圖片 / 影片 / 音訊：預設不主動處理。唯一例外：使用者 @mention 並引用該媒體，
  就下載 bytes 丟給 Gemini multimodal 分析。
- 文字檔：自動分析

指令列表請用 /help 查看。
"""

from __future__ import annotations

import json as _json
import fcntl
import hashlib
import logging
import mimetypes
import os
import re
import threading
import time
import uuid as _uuid
from contextlib import asynccontextmanager, contextmanager
from datetime import datetime, timedelta
from zoneinfo import ZoneInfo

import requests as _requests
from bs4 import BeautifulSoup

try:
    import yt_dlp as _yt_dlp  # type: ignore[import-untyped]

    _YTDLP_AVAILABLE = True
except ImportError:
    _yt_dlp = None
    _YTDLP_AVAILABLE = False

from fastapi import FastAPI, Header, HTTPException, Request
from google.genai import types
from linebot.v3 import WebhookParser  # type: ignore[import-untyped]
from linebot.v3.exceptions import InvalidSignatureError  # type: ignore[import-untyped]
from linebot.v3.messaging import (  # type: ignore[import-untyped]
    ApiClient,
    Configuration,
    ImageMessage,
    MessagingApi,
    MessagingApiBlob,
    PushMessageRequest,
    ReplyMessageRequest,
    TextMessage,
)
from linebot.v3.webhooks import (  # type: ignore[import-untyped]
    AudioMessageContent,
    FileMessageContent,
    GroupSource,
    ImageMessageContent,
    MemberJoinedEvent,
    MemberLeftEvent,
    MessageEvent,
    TextMessageContent,
    VideoMessageContent,
)

import burst_filter
import feedback_collector
import gemini_client
import memory
import output_validator
from pending_reply_policy import PENDING_REPLY_ENABLED as _DEFAULT_PENDING_REPLY_ENABLED
import review
from config import settings

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(name)s | %(message)s",
)


# log 時間戳用 PT（Gemini quota 以這個為準）+ TW（使用者看這個）雙時區，
# 避免「現在才 0800 為什麼 quota 就爆了」的誤判 — Gemini 的「一天」是 PT 的一天。
class _DualTZFormatter(logging.Formatter):
    _PT = ZoneInfo("America/Los_Angeles")
    _TW = ZoneInfo("Asia/Taipei")

    def formatTime(self, record, datefmt=None):
        pt = datetime.fromtimestamp(record.created, tz=self._PT)
        tw = datetime.fromtimestamp(record.created, tz=self._TW)
        return f"{pt.strftime('%m-%d %H:%M:%S')} PT ({tw.strftime('%H:%M')} TW)"


for _h in logging.getLogger().handlers:
    _h.setFormatter(
        _DualTZFormatter("%(asctime)s %(levelname)s %(name)s | %(message)s")
    )

logger = logging.getLogger("line_bot")


@asynccontextmanager
async def _app_lifespan(_app: FastAPI):
    """Run the same startup hooks previously registered via app.on_event."""
    if os.getenv("JOBS_ROUTES_ENABLED") == "1":
        from jobs_router import startup_sweep as _ss
        _ss()
    _process_pending_on_startup()
    _init_on_startup()
    yield


app = FastAPI(lifespan=_app_lifespan)

# Optional jobs router (n8n / external scheduler trigger surface).
# Feature flag: JOBS_ROUTES_ENABLED=1 enables /jobs/* endpoints.
# Load .env explicitly because pydantic-settings doesn't push to os.environ.
try:
    from dotenv import load_dotenv as _load_dotenv
    _load_dotenv(dotenv_path=os.path.join(os.path.dirname(__file__), ".env"), override=False)
except ImportError:
    pass
if os.getenv("JOBS_ROUTES_ENABLED") == "1":
    from jobs_router import router as jobs_router
    app.include_router(jobs_router)

_parser = WebhookParser(settings.line_channel_secret)


def _get_line_config() -> Configuration:
    """每次 LINE API 呼叫前用這個拿最新 token：
    - 有 LINE_CHANNEL_ID → v3 stateless short-lived，每 15 分自動 refresh
    - 沒設 → fallback 到 .env 的 long-lived token（向後相容）
    """
    from line_push_client import line_configuration

    return line_configuration()


# ── LINE 訊息配額 ─────────────────────────────────────────────────────────────


def _get_quota_footer() -> str:
    """配額 footer — 用戶反饋（2026-05-05）「拿掉 📊 Gemini 今日用量 99.0%」，
    平常 80%+ 警示移除（破壞對話自然度）。2026-05-31 再移除「今日用量已用完」
    footer，避免低價值 fallback 變成干擾訊息。
    """
    return ""


_OUTBOUND_SYSTEM_STATUS_MARKERS = (
    "咪寶聽到了但這個話題不太接得上",
    "Gemini 今日用量已用完",
    "Gemini 免費層今日請求額度已用完",
    "今日請求額度已用完",
    "可以再使用的時間",
    "aistudio.google.com",
    "Gemini 暫時忙不過來",
    "Gemini 那邊好像塞車",
    "Gemini 那邊暫時斷線",
    "Gemini API key",
    "Gemini 說這個輸入有問題",
    "分析失敗:",
    "model 沒載成功",
    "傳 LINE 失敗",
)


_OUTBOUND_INTERNAL_TRACE_MARKERS = (
    "The user is asking",
    "My response should",
    "Let's ensure it adheres to the rules",
    "As per rule",
    "I must respond",
    "I need to clarify",
    "First sentence needs",
    "This directly answers",
    "Rule 0:",
    "規則 0:",
)


def _is_internal_trace_outbound(text: str) -> bool:
    """True when an LLM exposed hidden reasoning/checklist text."""
    s = (text or "").lstrip()
    if not s:
        return False
    head = s[:2000]
    if re.match(r"(?is)^(THOUGHT|ANALYSIS|REASONING)\b", head):
        return True
    hits = sum(marker in head for marker in _OUTBOUND_INTERNAL_TRACE_MARKERS)
    return hits >= 2


def _is_system_status_outbound(text: str) -> bool:
    """True when text is an internal/quota/status message that must not go to LINE."""
    s = text or ""
    return (
        any(marker in s for marker in _OUTBOUND_SYSTEM_STATUS_MARKERS)
        or _is_internal_trace_outbound(s)
    )


def _strip_user_visible_mode_labels(text: str) -> str:
    """Remove internal fallback/model labels from text before it reaches LINE."""
    s = text or ""
    replacements = {
        "🔍 lite mode（Google 首頁 snippet）：": "🔍 Google 搜尋片段：",
        "🔍 lite mode（DuckDuckGo 搜尋片段）：": "🔍 DuckDuckGo 搜尋片段：",
    }
    for old, new in replacements.items():
        s = s.replace(old, new)
    patterns = (
        r"\n{0,3}（\s*lite mode\s*[—-]\s*local LLM\s*）",
        r"\n{0,3}（\s*local LLM fallback\s*）",
        r"\n{0,3}（\s*lite mode[^）]*）",
        r"\n{0,3}\(\s*lite mode[^)]*\)",
        r"\n{0,3}\(\s*local LLM fallback\s*\)",
    )
    for pattern in patterns:
        s = re.sub(pattern, "", s, flags=re.IGNORECASE)
    return s.strip()


def _prepare_outbound_text(text: str, *, source: str = "reply") -> str:
    """Normalize and validate text before it reaches any LINE text API."""
    prepared = _strip_user_visible_mode_labels(_md_to_line(text))
    if not prepared.strip():
        return ""
    result = output_validator.validate_outbound_text(prepared)
    if not result.ok:
        logger.warning(
            "outbound validation blocked source=%s reason=%s preview=%r",
            source,
            result.reason,
            prepared[:160],
        )
        return result.text
    return result.text


def _append_bot_turn(group_id: str, text: str, *, source: str = "bot_memory") -> None:
    """Store the validated user-visible bot text, not a blocked raw draft."""
    safe_text = _prepare_outbound_text(text, source=source)
    if safe_text:
        memory.append_turn(group_id, "bot", safe_text)


def _is_market_quote_outbound(text: str) -> bool:
    """Market quotes are reply-only; an expired reply token must not create a push."""
    return (text or "").lstrip().startswith(("【市場報價", "【即時股價"))


def _is_market_quote_request(text: str, context: list | None = None) -> bool:
    """Cheaply detect quote-seeking turns before Gemini may paraphrase the answer."""
    try:
        import stock_quote
        return stock_quote.should_try_contextual_quote(text or "", context=context)
    except Exception as e:
        logger.debug("market quote request detection skipped: %s", e)
        return False


def _get_explicit_market_quote_reply(
    text: str,
    context: list | None = None,
) -> str | None:
    """Return deterministic market quote text for explicit asks, before any LLM.

    Market data is factual/time-sensitive, so Gemini should not be the primary
    answerer. This also keeps Gemini quota exhaustion from blocking stock/gold
    quote requests that can be answered through public quote APIs.
    """
    if not text:
        return None
    try:
        import stock_quote
        technical = stock_quote.get_taiex_month_line_text(text)
        if technical:
            return technical
        return stock_quote.get_contextual_quotes_text(text, context=context) or None
    except Exception as e:
        logger.info("explicit market quote deterministic path failed: %s", e)
        return None


_LOCAL_TEXT_FALLBACK_SYSTEM_PROMPT = """你是 LINE 群組助理咪寶。現在雲端 Gemini 額度暫時用完，只能用本機 LLM 回覆。
請用繁體中文，第一句直接給判斷或回答，不要重述使用者問題。
只能輸出要發到 LINE 的正式回覆；嚴禁輸出 THOUGHT、ANALYSIS、REASONING、英文推理、自我檢查或規則清單。
若題目需要即時查證、最新價格、外部網頁或醫療法律投資專業判斷，而使用者沒有提供足夠資料，請明確說「本機模式無法即時查證」，但仍給可用的大方向、條件與下一步。
控制在 80 到 260 個中文字，避免空泛安慰。
如果只能輸出「有一定道理」「需要進一步查證」「需要多方面考量」「保持健康生活方式」這類模板句，請輸出空字串，不要回。"""


def _runtime_time_context_for_local_llm() -> str:
    now_tw = datetime.now(ZoneInfo("Asia/Taipei"))
    weekday = "一二三四五六日"[now_tw.weekday()]
    return (
        "即時時間基準："
        f"目前台灣時間 {now_tw.strftime('%Y-%m-%d %H:%M:%S')}（週{weekday}，Asia/Taipei）。"
        f"今天日期是 {now_tw.year}年{now_tw.month}月{now_tw.day}日。"
        "回答現在、目前、今年、去年、最新、完整數據時，以此為準，不要沿用舊年份。"
    )


def _local_text_fallback_system_prompt() -> str:
    return (
        _LOCAL_TEXT_FALLBACK_SYSTEM_PROMPT.rstrip()
        + "\n"
        + _runtime_time_context_for_local_llm()
    )

_GEMINI_QUOTA_RECHECK_INTERVAL_SEC = int(
    os.environ.get("GEMINI_QUOTA_RECHECK_INTERVAL_SEC", "1800")
)


def _local_text_llm_fallback(
    user_text: str,
    context: list[tuple[str, str]] | None = None,
) -> str:
    """Direct local text LLM fallback for quota outage after lite_reply misses."""
    text = (user_text or "").strip()
    if not text:
        return ""
    try:
        from local_llm import chat as local_chat

        out = local_chat(
            text,
            context=context,
            system_prompt=_local_text_fallback_system_prompt(),
            max_tokens=360,
        )
    except Exception as e:
        logger.warning("local text fallback failed: %s", e)
        return ""
    if out and isinstance(out, str) and len(out.strip()) > 5:
        return out.strip()
    return ""


def _llm_chat(
    user_input,
    context: list[tuple[str, str]],
    facts: list[str],
    pnotes: list[dict] | None = None,
) -> str:
    """Gemini chat; when quota is exhausted, fall back to local text generation."""
    if _quota_exhausted():
        if _quota_recheck_allowed():
            _record_quota_recheck_attempt()
            try:
                reply = gemini_client.chat(user_input, context, facts, pnotes)
            except Exception as e:
                if _is_quota_error(e):
                    _mark_quota_exhausted()
                    logger.warning("gemini quota recheck still exhausted")
                else:
                    logger.warning("gemini quota recheck failed: %s", e)
            else:
                _clear_quota_exhausted_after_recheck()
                return reply
        try:
            import lite_reply
            from gemini_client import _extract_text
            user_text = _extract_text(user_input)
            out = lite_reply.lite_reply(user_text, context=context)
            if out:
                logger.info(
                    "quota exhausted → lite_reply hit (text=%r)", user_text[:50]
                )
                return out
        except Exception as e:
            logger.warning("lite_reply fallback failed: %s", e)
            try:
                from gemini_client import _extract_text
                user_text = _extract_text(user_input)
            except Exception:
                user_text = str(user_input)
        out = _local_text_llm_fallback(user_text, context=context)
        if out:
            logger.info("quota exhausted → local text fallback hit")
            return out
        return ""
    return gemini_client.chat(user_input, context, facts, pnotes)


# ── URL 預抓取（繞過 Gemini url_context 的限制）─────────────────────────────

_URL_RE = re.compile(r"https?://\S+")
_YOUTUBE_BARE_URL_RE = re.compile(
    r"(?<![A-Za-z0-9./:-])"
    r"(?:www\.|m\.)?"
    r"(?:youtube\.com|youtube-nocookie\.com|youtu\.be)/\S+",
    re.IGNORECASE,
)
_PREFETCH_TIMEOUT = 5  # 秒，避免拖太久讓 reply_token 過期
_PREFETCH_MAX_CHARS = 5000  # 截斷上限，避免塞爆 prompt
_PREFETCH_MAX_URLS = 2  # 一次最多抓幾個連結
_PREFETCH_MIN_CHARS = 80  # 低於此長度視為垃圾（JS 渲染空殼），不塞進 prompt

_YTDLP_TIMEOUT = 12  # yt-dlp 單次提取上限（秒）
_YTDLP_SUBTITLE_MAX_CHARS = 3000
_YTDLP_SUBTITLE_LANGS = ["zh-TW", "zh-Hant", "zh", "zh-Hans", "en"]
_YTDLP_DESCRIPTION_MAX_CHARS = 1500

# Gemini Video Understanding fallback 參數
_GEMINI_VIDEO_TIMEOUT = 60  # 整個 download + upload + analyze 的總上限（秒）
_GEMINI_VIDEO_MAX_FILESIZE = "50M"  # yt-dlp --max-filesize
_GEMINI_VIDEO_MIN_REMAINING_QUOTA = 5  # 今日剩餘 < 此值就不啟動（影片呼叫 token 較重）
_GEMINI_VIDEO_THIN_THRESHOLD = 300  # 前面 prefetch 結果 < 此字數才 fallback
# 哪些 domain 才允許走 Gemini video fallback（短影音 / JS 渲染）
# 注意：YouTube 因為 yt-dlp 已能抓字幕，不啟動 fallback
_GEMINI_VIDEO_DOMAINS = re.compile(
    r"https?://(?:[a-z0-9-]+\.)*("
    r"tiktok\.com|instagram\.com|threads\.net|facebook\.com|fb\.watch|x\.com|twitter\.com"
    r")/",
    re.IGNORECASE,
)

# JS 渲染 / Cloudflare 保護的網站，requests.get() 抓不到有效內容
# 這些網站一律不 prefetch，直接讓 Gemini 用 Google Search 處理
_JS_RENDERED_DOMAINS = re.compile(
    r"https?://(?:[a-z0-9-]+\.)*("
    r"tiktok\.com|instagram\.com|threads\.net|facebook\.com|fb\.watch|"
    r"dcard\.tw|x\.com|twitter\.com|reddit\.com|"
    r"youtube\.com/shorts|youtu\.be"
    r")/",
    re.IGNORECASE,
)

# TikTok 短網址 pattern（vt.tiktok.com / vm.tiktok.com），需先 redirect 才能丟 oEmbed
_TIKTOK_SHORT_DOMAIN = re.compile(r"https?://(?:vt|vm)\.tiktok\.com/", re.IGNORECASE)
# 從 oEmbed html 欄位抽背景音樂
# html 結構：<a title="♬ xxx" href="..."> ♬ xxx</a>，title 裡也有 ♬ 會誤匹配，
# 所以要求 ♬ 前面必須是 `>`（真正的 anchor content，不是屬性值）
_TIKTOK_MUSIC_RE = re.compile(r">\s*♬\s*([^<]+?)\s*</a>", re.UNICODE)


def _parse_vtt(vtt_text: str) -> str:
    """WebVTT → 純文字，去掉時間碼、HTML tag、相鄰重複行。"""
    lines = []
    for line in vtt_text.splitlines():
        line = line.strip()
        if not line or line.startswith("WEBVTT") or "-->" in line or line.isdigit():
            continue
        line = re.sub(r"<[^>]+>", "", line)
        if line:
            lines.append(line)
    deduped: list[str] = []
    for ln in lines:
        if not deduped or ln != deduped[-1]:
            deduped.append(ln)
    return "\n".join(deduped)


def _clean_prefetch_url(url: str) -> str:
    """Trim punctuation often included when users paste links in chat."""
    return (url or "").strip().rstrip("。．，,、；;：:！!？?）)]}>\"'")


def _with_scheme_for_youtube(url: str) -> str:
    cleaned = _clean_prefetch_url(url)
    if cleaned and not re.match(r"https?://", cleaned, re.IGNORECASE):
        if _YOUTUBE_BARE_URL_RE.fullmatch(cleaned):
            return f"https://{cleaned}"
    return cleaned


def _is_youtube_url(url: str) -> bool:
    from urllib.parse import urlparse

    lower = _with_scheme_for_youtube(url).lower()
    host = urlparse(lower).netloc
    return (
        host == "youtu.be"
        or host.endswith(".youtu.be")
        or host == "youtube.com"
        or host.endswith(".youtube.com")
        or host == "youtube-nocookie.com"
        or host.endswith(".youtube-nocookie.com")
    )


def _extract_prefetch_urls(text: str) -> list[str]:
    """Extract HTTP URLs plus bare YouTube URLs, prioritizing YouTube later."""
    raw_urls = list(_URL_RE.findall(text or ""))
    raw_urls.extend(
        match.group(0) for match in _YOUTUBE_BARE_URL_RE.finditer(text or "")
    )

    deduped: list[str] = []
    seen: set[str] = set()
    for raw_url in raw_urls:
        cleaned = _with_scheme_for_youtube(raw_url)
        if not cleaned or cleaned in seen:
            continue
        deduped.append(cleaned)
        seen.add(cleaned)
    return deduped


def _extract_youtube_video_id(url: str) -> str | None:
    """Return a YouTube video id from common watch/shorts/live/embed forms."""
    from urllib.parse import parse_qs, urlparse

    cleaned = _with_scheme_for_youtube(url)
    if not cleaned:
        return None
    parsed = urlparse(cleaned)
    host = parsed.netloc.lower()
    path = parsed.path or ""

    def _valid(candidate: str | None) -> str | None:
        candidate = (candidate or "").strip()
        if re.fullmatch(r"[A-Za-z0-9_-]{11}", candidate):
            return candidate
        return None

    if host == "youtu.be" or host.endswith(".youtu.be"):
        return _valid(path.lstrip("/").split("/", 1)[0])

    is_youtube_host = (
        host == "youtube.com"
        or host.endswith(".youtube.com")
        or host == "youtube-nocookie.com"
        or host.endswith(".youtube-nocookie.com")
    )
    if not is_youtube_host:
        return None

    if path == "/watch":
        return _valid((parse_qs(parsed.query).get("v") or [""])[0])

    for prefix in ("/shorts/", "/live/", "/embed/", "/v/"):
        if path.startswith(prefix):
            return _valid(path[len(prefix) :].split("/", 1)[0])

    return None


def _canonical_youtube_url(url: str) -> str:
    video_id = _extract_youtube_video_id(url)
    if video_id:
        return f"https://www.youtube.com/watch?v={video_id}"
    return _clean_prefetch_url(url)


def _format_duration_seconds(duration: object) -> str | None:
    try:
        total = int(duration)  # type: ignore[arg-type]
    except (TypeError, ValueError):
        return None
    if total <= 0:
        return None
    h, rem = divmod(total, 3600)
    m, s = divmod(rem, 60)
    if h:
        return f"{h}:{m:02d}:{s:02d}"
    return f"{m}:{s:02d}"


def _youtube_live_status_label(info: dict) -> str | None:
    live_status = str(info.get("live_status") or "").strip().lower()
    if info.get("is_live") or live_status == "is_live":
        return "正在直播"
    if live_status in {"is_upcoming", "is_upcoming_event"}:
        return "預定直播"
    if info.get("was_live") or live_status in {"was_live", "post_live"}:
        return "直播已結束 / 重播"
    if live_status:
        return live_status
    return None


def _append_youtube_metadata_lines(lines: list[str], info: dict) -> None:
    title = (info.get("title") or "").strip()
    uploader = (info.get("uploader") or info.get("channel") or "").strip()
    channel_url = (info.get("channel_url") or info.get("uploader_url") or "").strip()
    webpage_url = (info.get("webpage_url") or "").strip()
    duration_text = (
        info.get("duration_string") or _format_duration_seconds(info.get("duration"))
    )
    live_label = _youtube_live_status_label(info)
    view_count = info.get("view_count")
    description = (info.get("description") or "").strip()

    if title:
        lines.append(f"標題：{title}")
    if uploader:
        lines.append(f"頻道：{uploader}")
    if channel_url:
        lines.append(f"頻道網址：{channel_url}")
    if webpage_url:
        lines.append(f"影片網址：{webpage_url}")
    if live_label:
        lines.append(f"直播狀態：{live_label}")
    if duration_text:
        lines.append(f"長度：{duration_text}")
    if isinstance(view_count, int):
        lines.append(f"觀看次數：{view_count:,}")
    if description:
        desc = (
            description[:_YTDLP_DESCRIPTION_MAX_CHARS] + "…（描述截斷）"
            if len(description) > _YTDLP_DESCRIPTION_MAX_CHARS
            else description
        )
        lines.append(f"描述：{desc}")


def _extract_subtitles_from_info(info: dict) -> str | None:
    """從 yt-dlp info dict 拿字幕文字（優先人工字幕 → 自動生成，語言優先順序見常數）。"""
    for subs_dict in (
        info.get("subtitles") or {},
        info.get("automatic_captions") or {},
    ):
        for lang in _YTDLP_SUBTITLE_LANGS:
            entries = subs_dict.get(lang)
            if not entries:
                continue
            entry = next((e for e in entries if e.get("ext") == "vtt"), entries[0])
            sub_url = entry.get("url") if entry else None
            if not sub_url:
                continue
            try:
                resp = _requests.get(sub_url, timeout=8)
                resp.raise_for_status()
                text = _parse_vtt(resp.text)
                if text and len(text) > 50:
                    if len(text) > _YTDLP_SUBTITLE_MAX_CHARS:
                        text = text[:_YTDLP_SUBTITLE_MAX_CHARS] + "…（字幕截斷）"
                    return text
            except Exception as e:
                logger.debug("subtitle download failed lang=%s: %s", lang, e)
    return None


def _fetch_video_ytdlp(url: str) -> str | None:
    """
    用 yt-dlp 抓影片 metadata + 字幕，支援 YouTube、TikTok、IG、FB、X 等 1000+ 網站。

    優先抓字幕（中文 > 英文）；沒字幕就用 title + description。
    任何錯誤都回 None（讓 caller fallback）。
    """
    if not _YTDLP_AVAILABLE:
        return None
    try:
        ydl_opts = {
            "quiet": True,
            "no_warnings": True,
            "skip_download": True,
            "socket_timeout": _YTDLP_TIMEOUT,
        }
        with _yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=False)
        if not info:
            return None

        if info.get("_type") == "playlist" and info.get("entries"):
            first = next((e for e in info.get("entries") or [] if e), None)
            if isinstance(first, dict):
                info = first

        title = (info.get("title") or "").strip()
        uploader = (info.get("uploader") or info.get("channel") or "").strip()
        description = (info.get("description") or "").strip()
        subtitle_text = _extract_subtitles_from_info(info)

        if not any((title, uploader, description, subtitle_text)):
            return None

        lines = [f"（以下是影片連結 {url} 的內容，透過 yt-dlp 擷取）"]
        lines.append("--- 影片資訊開始 ---")
        _append_youtube_metadata_lines(lines, info)
        if subtitle_text:
            lines.append(f"\n字幕內容：\n{subtitle_text}")
        else:
            lines.append(
                "\n字幕內容：未取得可用字幕或逐字稿；請根據上方標題、頻道、"
                "直播狀態與描述判斷主題，不要聲稱已看完整影片。"
            )

        lines.append("--- 影片資訊結束 ---")
        block = "\n".join(lines)
        logger.info(
            "ytdlp OK url=%s chars=%d has_subs=%s", url, len(block), bool(subtitle_text)
        )
        return block
    except Exception as e:
        logger.info("ytdlp failed url=%s: %s", url, e)
        return None


def _fetch_tiktok_meta(url: str) -> str | None:
    """
    TikTok 專用 prefetch：走官方 oEmbed API（公開 endpoint，免 token）取 caption / 作者 / 音樂。

    為什麼要這層：TikTok 是 JS 渲染，requests.get() 只抓到空殼；而 Gemini url_context
    對 TikTok 實測 100% 回空字串（連三次 empty reply 後 raise RuntimeError）。
    oEmbed endpoint 直接吐 JSON，能拿到 title（caption + hashtags）/ author_name /
    author_unique_id / html（內含音樂資訊）。

    失敗時回 None，由 caller fallback 回原本 skip 行為，不會退步。
    """
    try:
        # 短網址（vt.tiktok.com / vm.tiktok.com）先 HEAD follow redirect 拿完整 URL
        target_url = url
        if _TIKTOK_SHORT_DOMAIN.search(url):
            try:
                r = _requests.head(
                    url,
                    timeout=_PREFETCH_TIMEOUT,
                    allow_redirects=True,
                    headers={
                        "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7)"
                    },
                )
                target_url = r.url
                logger.info("tiktok short url resolved: %s → %s", url, target_url)
            except Exception as e:
                logger.info("tiktok short url resolve failed url=%s: %s", url, e)
                return None

        # 呼叫 oEmbed API
        resp = _requests.get(
            "https://www.tiktok.com/oembed",
            params={"url": target_url},
            timeout=_PREFETCH_TIMEOUT,
            headers={"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7)"},
        )
        if resp.status_code != 200:
            logger.info("tiktok oembed HTTP %d url=%s", resp.status_code, target_url)
            return None

        data = resp.json()
        # oEmbed error response 會是 {"message": "...", "code": 4xx}
        if data.get("code") and int(data.get("code", 0)) >= 400:
            logger.info(
                "tiktok oembed error code=%s url=%s", data.get("code"), target_url
            )
            return None

        title = (data.get("title") or "").strip()
        author_name = (data.get("author_name") or "").strip()
        author_id = (data.get("author_unique_id") or "").strip()

        # 從 html 欄位抽出背景音樂資訊
        html_field = data.get("html") or ""
        music_match = _TIKTOK_MUSIC_RE.search(html_field)
        music = music_match.group(1).strip() if music_match else ""

        if not title and not author_name:
            logger.info("tiktok oembed empty content url=%s", target_url)
            return None

        lines = [f"（以下是 TikTok 連結 {url} 的影片資訊，透過 oEmbed API 擷取）"]
        lines.append("--- TikTok 影片資訊開始 ---")
        if author_name:
            author_line = f"作者：{author_name}"
            if author_id:
                author_line += f" (@{author_id})"
            lines.append(author_line)
        if title:
            lines.append(f"影片描述：{title}")
        if music:
            lines.append(f"背景音樂：{music}")
        lines.append("--- TikTok 影片資訊結束 ---")

        block = "\n".join(lines)
        logger.info(
            "tiktok oembed OK url=%s author=%s chars=%d",
            url,
            author_id or author_name,
            len(block),
        )
        return block
    except Exception as e:
        logger.info("tiktok oembed failed url=%s: %s", url, e)
        return None


def _fetch_youtube_meta(url: str) -> str | None:
    """
    YouTube（含 shorts、youtu.be）走官方 oEmbed API 拿 title + 頻道。免 token、免 auth。

    為什麼要這層：
      - youtube.com/shorts / youtu.be 在 JS 白名單裡（目前 skip，讓 Gemini 處理，但 url_context 對 shorts 吐 metadata 不穩定）
      - youtube.com/watch 走 generic HTML prefetch 只抓到 ~280 chars boilerplate
      - oEmbed endpoint 穩定吐 title + author_name，比前兩條路都好

    限制：oEmbed 不提供 description，想拿內容描述還是只能靠 Gemini；但至少 title 有了。
    """
    try:
        resp = _requests.get(
            "https://www.youtube.com/oembed",
            params={"url": url, "format": "json"},
            timeout=_PREFETCH_TIMEOUT,
            headers={"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7)"},
        )
        if resp.status_code != 200:
            logger.info("youtube oembed HTTP %d url=%s", resp.status_code, url)
            return None

        data = resp.json()
        title = (data.get("title") or "").strip()
        author = (data.get("author_name") or "").strip()
        if not title and not author:
            return None

        lines = [f"（以下是 YouTube 連結 {url} 的影片資訊，透過 oEmbed API 擷取）"]
        lines.append("--- YouTube 影片資訊開始 ---")
        if title:
            lines.append(f"標題：{title}")
        if author:
            lines.append(f"頻道：{author}")
        lines.append(
            "資料限制：oEmbed 只提供標題與頻道，未取得字幕或逐字稿；"
            "請仍根據可取得的 metadata 判斷直播/影片主題，並標示資訊限制。"
        )
        lines.append("--- YouTube 影片資訊結束 ---")
        block = "\n".join(lines)
        logger.info("youtube oembed OK url=%s chars=%d", url, len(block))
        return block
    except Exception as e:
        logger.info("youtube oembed failed url=%s: %s", url, e)
        return None


def _meta_content(soup: BeautifulSoup, key: str) -> str:
    for attr in ("property", "name", "itemprop"):
        tag = soup.find("meta", attrs={attr: key})
        if tag and tag.get("content"):
            return str(tag.get("content")).strip()
    return ""


def _extract_youtube_player_response(html_text: str) -> dict | None:
    marker = "ytInitialPlayerResponse"
    start_at = html_text.find(marker)
    while start_at >= 0:
        equals_at = html_text.find("=", start_at)
        json_start = html_text.find("{", equals_at if equals_at >= 0 else start_at)
        if json_start < 0:
            return None
        try:
            data, _ = _json.JSONDecoder().raw_decode(html_text[json_start:])
        except ValueError:
            start_at = html_text.find(marker, start_at + len(marker))
            continue
        if isinstance(data, dict):
            return data
        return None
    return None


def _format_count(value: object) -> str | None:
    try:
        return f"{int(str(value)):,}"
    except (TypeError, ValueError):
        return None


def _fetch_youtube_html_meta(url: str) -> str | None:
    """Fetch YouTube watch HTML and parse og/meta + ytInitialPlayerResponse."""
    try:
        canonical_url = _canonical_youtube_url(url)
        resp = _requests.get(
            canonical_url,
            timeout=_PREFETCH_TIMEOUT,
            headers={
                "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7)",
                "Accept-Language": "zh-TW,zh;q=0.9,en;q=0.6",
            },
        )
        if resp.status_code != 200:
            logger.info("youtube html HTTP %d url=%s", resp.status_code, canonical_url)
            return None

        soup = BeautifulSoup(resp.text, "html.parser")
        player = _extract_youtube_player_response(resp.text) or {}
        details = player.get("videoDetails") if isinstance(player, dict) else {}
        if not isinstance(details, dict):
            details = {}

        title = (
            (details.get("title") or "").strip()
            or _meta_content(soup, "og:title")
            or (soup.title.get_text(strip=True) if soup.title else "")
        )
        author = (details.get("author") or "").strip() or _meta_content(
            soup, "author"
        )
        description = (
            (details.get("shortDescription") or "").strip()
            or _meta_content(soup, "og:description")
            or _meta_content(soup, "description")
        )
        duration_text = _format_duration_seconds(details.get("lengthSeconds"))
        view_count = _format_count(details.get("viewCount"))
        is_live = bool(details.get("isLiveContent"))
        video_url = _meta_content(soup, "og:url") or canonical_url

        if not any((title, author, description)):
            logger.info("youtube html empty metadata url=%s", canonical_url)
            return None

        lines = [
            f"（以下是 YouTube 連結 {url} 的影片資訊，透過 YouTube HTML metadata 擷取）"
        ]
        lines.append("--- YouTube 影片資訊開始 ---")
        if title:
            lines.append(f"標題：{title}")
        if author:
            lines.append(f"頻道：{author}")
        if video_url:
            lines.append(f"影片網址：{video_url}")
        if is_live:
            lines.append("直播狀態：直播內容或直播重播")
        if duration_text:
            lines.append(f"長度：{duration_text}")
        if view_count:
            lines.append(f"觀看次數：{view_count}")
        if description:
            desc = (
                description[:_YTDLP_DESCRIPTION_MAX_CHARS] + "…（描述截斷）"
                if len(description) > _YTDLP_DESCRIPTION_MAX_CHARS
                else description
            )
            lines.append(f"描述：{desc}")
        lines.append(
            "字幕內容：未取得可用字幕或逐字稿；請根據上方標題、頻道、"
            "直播狀態與描述判斷主題，不要要求使用者自行點擊觀看。"
        )
        lines.append("--- YouTube 影片資訊結束 ---")
        block = "\n".join(lines)
        logger.info("youtube html metadata OK url=%s chars=%d", canonical_url, len(block))
        return block
    except Exception as e:
        logger.info("youtube html metadata failed url=%s: %s", url, e)
        return None


def _youtube_unavailable_block(url: str) -> str:
    cleaned = _clean_prefetch_url(url)
    canonical_url = _canonical_youtube_url(cleaned)
    video_id = _extract_youtube_video_id(cleaned)
    lines = [
        f"（以下是 YouTube 連結 {cleaned} 的抓取狀態；系統已辨識為 YouTube 影片或直播）",
        "--- YouTube 影片資訊開始 ---",
    ]
    if video_id:
        lines.append(f"影片 ID：{video_id}")
    if canonical_url:
        lines.append(f"影片網址：{canonical_url}")
    lines.append(
        "抓取結果：yt-dlp、oEmbed、YouTube HTML metadata 這次都沒有取得標題、描述或字幕。"
    )
    lines.append(
        "回覆要求：請明確說這次抓取未取得可用 metadata，建議稍後重試或補貼標題/截圖；"
        "禁止使用把操作交回使用者、只描述平台網域、或要求補關鍵字才處理的退讓話術。"
    )
    lines.append("--- YouTube 影片資訊結束 ---")
    return "\n".join(lines)


def _fetch_youtube_context(url: str) -> str:
    canonical_url = _canonical_youtube_url(url)
    return (
        _fetch_video_ytdlp(canonical_url)
        or _fetch_youtube_meta(canonical_url)
        or _fetch_youtube_html_meta(canonical_url)
        or _youtube_unavailable_block(url)
    )


# Reddit 短網址 pattern：
#   舊 redd.it/xxx（短 domain）
#   新 reddit.com/r/sub/s/xxx（share link）
_REDDIT_SHORT_DOMAIN = re.compile(r"https?://redd\.it/", re.IGNORECASE)
_REDDIT_SHARE_PATH = re.compile(r"reddit\.com/r/[^/]+/s/", re.IGNORECASE)


def _fetch_reddit_meta(url: str) -> str | None:
    """
    Reddit 走公開 .json endpoint 拿 post + 前三條 top comments。免 token，但需 User-Agent。

    為什麼要這層：
      - reddit.com 在 JS 白名單，目前 skip；而 Gemini url_context 對 reddit 常常只拿到 meta tag
      - .json endpoint 是 reddit 官方認可的 public API，吐結構化 JSON（title / selftext / comments）
      - 拿到 selftext + 熱門留言，資訊量遠大於 Gemini 原本拿到的 meta
    """
    try:
        # 短網址 resolve：redd.it/xxx 和 reddit.com/r/.../s/xxx 都要先 follow redirect
        target = url
        if _REDDIT_SHORT_DOMAIN.search(url) or _REDDIT_SHARE_PATH.search(url):
            try:
                r = _requests.head(
                    url,
                    timeout=_PREFETCH_TIMEOUT,
                    allow_redirects=True,
                    headers={"User-Agent": "andrew-line-bot/1.0"},
                )
                target = r.url
                logger.info("reddit short url resolved: %s → %s", url, target)
            except Exception as e:
                logger.info("reddit short url resolve failed url=%s: %s", url, e)
                return None

        # 非貼文 URL（例如 subreddit 首頁、使用者頁面）沒 .json 可抓
        if "/comments/" not in target:
            logger.info("reddit url 非貼文格式 (no /comments/) url=%s", target)
            return None

        # 砍 query/fragment，path 結尾加 .json
        from urllib.parse import urlsplit, urlunsplit

        parts = urlsplit(target)
        json_path = parts.path.rstrip("/") + ".json"
        json_url = urlunsplit((parts.scheme, parts.netloc, json_path, "", ""))

        resp = _requests.get(
            json_url,
            timeout=_PREFETCH_TIMEOUT,
            headers={"User-Agent": "andrew-line-bot/1.0 (LINE chatbot prefetcher)"},
        )
        if resp.status_code != 200:
            logger.info("reddit .json HTTP %d url=%s", resp.status_code, json_url)
            return None

        data = resp.json()
        # 正常 response：[post_listing, comments_listing]
        if not isinstance(data, list) or len(data) < 1:
            return None

        post_children = data[0].get("data", {}).get("children", [])
        if not post_children:
            return None
        post = post_children[0].get("data", {}) or {}

        title = (post.get("title") or "").strip()
        if not title:
            return None

        selftext = (post.get("selftext") or "").strip()
        subreddit = (post.get("subreddit") or "").strip()
        author = (post.get("author") or "").strip()
        score = post.get("score", 0)
        num_comments = post.get("num_comments", 0)

        # 前三條 top-level 留言（跳過 deleted / removed）
        top_comments: list[str] = []
        if len(data) > 1:
            for child in data[1].get("data", {}).get("children", []):
                if len(top_comments) >= 3:
                    break
                c = child.get("data", {}) or {}
                body = (c.get("body") or "").strip()
                if not body or body in ("[deleted]", "[removed]"):
                    continue
                if len(body) > 300:
                    body = body[:300] + "…"
                top_comments.append(
                    f"  - u/{c.get('author', '?')} ({c.get('score', 0)} 分): {body}"
                )

        # 內文截斷（避免塞爆 prompt；留空間給 comments）
        if len(selftext) > _PREFETCH_MAX_CHARS - 500:
            selftext = selftext[: _PREFETCH_MAX_CHARS - 500] + "…（內文截斷）"

        lines = [f"（以下是 Reddit 貼文 {url} 的內容，透過 .json endpoint 擷取）"]
        lines.append("--- Reddit 貼文開始 ---")
        lines.append(f"版：r/{subreddit}")
        lines.append(f"作者：u/{author}")
        lines.append(f"標題：{title}")
        lines.append(f"分數：{score} / 留言數：{num_comments}")
        if selftext:
            lines.append(f"內文：\n{selftext}")
        if top_comments:
            lines.append("熱門留言：")
            lines.extend(top_comments)
        lines.append("--- Reddit 貼文結束 ---")

        block = "\n".join(lines)
        logger.info(
            "reddit .json OK url=%s subreddit=%s comments=%d chars=%d",
            url,
            subreddit,
            len(top_comments),
            len(block),
        )
        return block
    except Exception as e:
        logger.info("reddit .json failed url=%s: %s", url, e)
        return None


_IG_REEL_RE = re.compile(r"instagram\.com/(reel|p)/([A-Za-z0-9_-]+)", re.IGNORECASE)


def _fetch_instagram_embed(url: str) -> str | None:
    """
    Instagram Reels / Posts 的 embed 頁面 fallback。

    yt-dlp 對 IG 失敗率高，這層直接抓 /embed/ 公開頁面，
    用 BeautifulSoup 解出 caption（不需要 token / 登入）。
    """
    m = _IG_REEL_RE.search(url)
    if not m:
        return None
    shortcode = m.group(2)
    kind = m.group(1).lower()  # reel 或 p
    embed_url = f"https://www.instagram.com/{kind}/{shortcode}/embed/"
    try:
        resp = _requests.get(
            embed_url,
            timeout=_PREFETCH_TIMEOUT,
            headers={
                "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/124.0.0.0 Safari/537.36",
                "Accept-Language": "zh-TW,zh;q=0.9,en;q=0.8",
            },
        )
        if resp.status_code != 200:
            logger.info("ig embed HTTP %d url=%s", resp.status_code, embed_url)
            return None

        soup = BeautifulSoup(resp.text, "html.parser")

        # caption 通常在 <div class="Caption"> 或 meta description
        caption = ""
        caption_div = soup.find("div", class_=re.compile(r"Caption", re.I))
        if caption_div:
            caption = caption_div.get_text(separator=" ", strip=True)

        if not caption:
            meta = soup.find("meta", attrs={"name": "description"}) or soup.find(
                "meta", attrs={"property": "og:description"}
            )
            if meta:
                caption = str(meta.get("content") or "").strip()

        if not caption or len(caption) < 10:
            logger.info("ig embed: no caption found url=%s", url)
            return None

        if len(caption) > 800:
            caption = caption[:800] + "…"

        block = (
            f"（以下是 Instagram 連結 {url} 的內容，透過 embed 頁面擷取）\n"
            f"--- Instagram 內容開始 ---\n"
            f"Caption：{caption}\n"
            f"--- Instagram 內容結束 ---"
        )
        logger.info("ig embed OK url=%s chars=%d", url, len(block))
        return block
    except Exception as e:
        logger.info("ig embed failed url=%s: %s", url, e)
        return None


def _gemini_video_quota_ok() -> bool:
    """今日 Gemini quota 剩餘 >= _GEMINI_VIDEO_MIN_REMAINING_QUOTA 才允許走 video fallback。

    影片呼叫的 token 量比文字大很多（一段 30 秒影片 ~3-10k tokens），需要保留餘裕給
    一般對話。失敗時 conservative，回 False（寧可不啟動也不要爆 quota）。
    """
    try:
        info = gemini_client.get_gemini_quota_info()
        if info is None:
            return False
        remaining_req = info["limit_requests"] - info["used_requests"]
        return remaining_req >= _GEMINI_VIDEO_MIN_REMAINING_QUOTA
    except Exception:
        return False


_GEMINI_SIDE_TASK_MIN_REMAINING_REQUESTS = int(
    os.environ.get("GEMINI_SIDE_TASK_MIN_REMAINING_REQUESTS", "16")
)
_GEMINI_SIDE_TASK_MIN_REMAINING_TOKENS = int(
    os.environ.get("GEMINI_SIDE_TASK_MIN_REMAINING_TOKENS", "50000")
)


def _gemini_side_task_allowed(reason: str = "side_task") -> bool:
    """Return True only when optional Gemini work can run without eating reply budget."""
    try:
        if _quota_exhausted():
            logger.info("skip Gemini %s: quota exhausted", reason)
            return False
        info = gemini_client.get_gemini_quota_info()
        if info is None:
            logger.info("skip Gemini %s: usage unavailable", reason)
            return False
        remaining_req = int(info["limit_requests"]) - int(info["used_requests"])
        remaining_tokens = int(info["limit_tokens"]) - int(info["used_tokens"])
        if remaining_req <= _GEMINI_SIDE_TASK_MIN_REMAINING_REQUESTS:
            logger.info(
                "skip Gemini %s: remaining_requests=%d reserve=%d",
                reason,
                remaining_req,
                _GEMINI_SIDE_TASK_MIN_REMAINING_REQUESTS,
            )
            return False
        if remaining_tokens <= _GEMINI_SIDE_TASK_MIN_REMAINING_TOKENS:
            logger.info(
                "skip Gemini %s: remaining_tokens=%d reserve=%d",
                reason,
                remaining_tokens,
                _GEMINI_SIDE_TASK_MIN_REMAINING_TOKENS,
            )
            return False
        return True
    except Exception as e:
        logger.info("skip Gemini %s: budget check failed: %s", reason, e)
        return False


def _fetch_video_gemini(url: str) -> str | None:
    """
    終極 fallback：用 yt-dlp 把短影片抓下來，丟給 Gemini 2.5 Flash 用 Files API 分析。

    觸發條件（caller 負責判斷）：
      - prefetch chain 拿到的內容 < _GEMINI_VIDEO_THIN_THRESHOLD 字
      - URL 屬於 _GEMINI_VIDEO_DOMAINS（TikTok / IG / FB / Threads / X，不含 YouTube）
      - 今日 Gemini quota 剩餘 >= _GEMINI_VIDEO_MIN_REMAINING_QUOTA

    流程：
      1. yt-dlp 下載 MP4 到 /tmp/linebot_video_<hash>.mp4（--max-filesize 50M）
      2. 上傳到 Gemini Files API
      3. 輪詢直到 file.state == ACTIVE
      4. generate_content("用繁體中文簡短描述...")
      5. 清理：刪除 Files API 上的檔案 + 本地 mp4

    任何錯誤回 None，由 caller fallback。整個流程有 60 秒總 timeout 保護。
    """
    if not _YTDLP_AVAILABLE:
        return None

    import hashlib
    from pathlib import Path

    start_ts = time.time()
    url_hash = hashlib.md5(url.encode()).hexdigest()[:12]
    local_path = f"/tmp/linebot_video_{url_hash}.mp4"
    uploaded_file = None

    def _elapsed_ok() -> bool:
        return (time.time() - start_ts) < _GEMINI_VIDEO_TIMEOUT

    try:
        # 1) 下載影片：優先低畫質（短影音分析不需要 1080p，能省 5-10x bandwidth）
        ydl_opts = {
            "quiet": True,
            "no_warnings": True,
            "noplaylist": True,
            "max_filesize": 50 * 1024 * 1024,
            "outtmpl": local_path,
            # 偏好最差但有畫質的 mp4：降低 bandwidth、加快 download
            "format": "worst[ext=mp4]/worst/mp4",
            # 短影音 CDN 偶爾很慢，整個下載階段給 35 秒（剩 25 秒給 upload + analyze）
            "socket_timeout": 35,
            "overwrites": True,
        }
        with _yt_dlp.YoutubeDL(ydl_opts) as ydl:
            ydl.download([url])

        if not os.path.exists(local_path) or os.path.getsize(local_path) == 0:
            logger.info("gemini video fallback: download empty url=%s", url)
            return None

        if not _elapsed_ok():
            logger.info("gemini video fallback: timeout after download url=%s", url)
            return None

        # 2) 上傳到 Files API
        uploaded_file = gemini_client._client.files.upload(file=Path(local_path))

        # 3) 輪詢直到 ACTIVE
        while True:
            if not _elapsed_ok():
                logger.info(
                    "gemini video fallback: timeout while waiting active url=%s", url
                )
                return None
            assert uploaded_file.name is not None  # narrow for mypy
            f = gemini_client._client.files.get(name=uploaded_file.name)
            state = getattr(f.state, "name", str(f.state))
            if state == "ACTIVE":
                uploaded_file = f
                break
            if state == "FAILED":
                logger.info("gemini video fallback: file FAILED url=%s", url)
                return None
            time.sleep(1)

        # 4) generate_content — flash 滿了自動 fallback flash-lite
        prompt = (
            "用繁體中文簡短描述這部影片的主要內容（含口白、畫面、訊息），"
            "200 字內。如果是廣告/業配/沒重點請直說。"
        )
        models_to_try = [settings.gemini_model]
        if settings.gemini_light_model and settings.gemini_light_model != settings.gemini_model:
            models_to_try.append(settings.gemini_light_model)

        response = None
        last_error = None
        contents_list = [uploaded_file, prompt]
        for model_name in models_to_try:
            try:
                response = gemini_client._client.models.generate_content(
                    model=model_name,
                    contents=contents_list,  # type: ignore[arg-type]
                )
                logger.info(
                    "gemini video fallback: model=%s succeeded url=%s",
                    model_name, url,
                )
                break
            except Exception as e:
                last_error = e
                err_str = str(e)
                is_quota = "429" in err_str or "RESOURCE_EXHAUSTED" in err_str
                if is_quota and model_name != models_to_try[-1]:
                    logger.info(
                        "gemini video fallback: model=%s 429 quota, trying lite url=%s",
                        model_name, url,
                    )
                    continue
                # 非 quota 錯誤、或已經試完所有模型 → 拋給外層 except
                raise

        if response is None:
            logger.info("gemini video fallback: all models exhausted url=%s err=%s", url, last_error)
            return None

        # 計入 quota counter
        try:
            gemini_client._track_usage(response)
        except Exception:
            pass

        text = (response.text or "").strip()
        if not text:
            logger.info("gemini video fallback: empty response url=%s", url)
            return None

        block = (
            f"（以下是影片連結 {url} 的內容，由 Gemini Video Understanding 分析）\n"
            f"--- 影片分析開始 ---\n"
            f"{text}\n"
            f"--- 影片分析結束 ---"
        )
        logger.info(
            "gemini video fallback used url=%s chars=%d", url, len(block)
        )
        return block
    except Exception as e:
        logger.info("gemini video fallback failed url=%s: %s", url, e)
        return None
    finally:
        # 清理 Files API 上傳的檔案
        if uploaded_file is not None and uploaded_file.name is not None:
            try:
                gemini_client._client.files.delete(name=uploaded_file.name)
            except Exception as e:
                logger.debug("gemini video fallback: delete remote failed: %s", e)
        # 清理本地 mp4
        try:
            if os.path.exists(local_path):
                os.unlink(local_path)
        except Exception as e:
            logger.debug("gemini video fallback: unlink local failed: %s", e)


def _maybe_video_fallback(url: str, current_block: str | None) -> str | None:
    """若 current_block 太薄且 URL 是視訊平台，嘗試 Gemini video fallback。

    回傳：fallback 結果 OR 原 current_block（保持原行為）
    """
    if not _GEMINI_VIDEO_DOMAINS.search(url):
        return current_block
    chars = len(current_block) if current_block else 0
    if current_block is not None and chars >= _GEMINI_VIDEO_THIN_THRESHOLD:
        return current_block
    # 不做本地 quota pre-block——讓 Google server-side 決定。429 由 _fetch_video_gemini
    # 內部 catch 後回 None，pipeline 自然 fallback 不影響流程。能送就送、不要白白浪費機會。
    fallback = _fetch_video_gemini(url)
    if fallback:
        # 如果有原本的 block（薄但非空），把兩者拼起來給 model 更多 context
        if current_block:
            return current_block + "\n\n" + fallback
        return fallback
    return current_block


def _prefetch_urls(text: str) -> str:
    """
    從文字中抽出 URL，用 Python requests 預先抓取網頁內容，
    轉成純文字後塞進 prompt。

    特殊平台優先走公開 API（oEmbed / .json），比 HTML prefetch 或 Gemini url_context 穩定：
      - TikTok  → www.tiktok.com/oembed（caption + author + music）
      - YouTube → www.youtube.com/oembed（title + channel）
      - Reddit  → <permalink>.json（title + selftext + top 3 comments）

    其他 JS 渲染網站（IG/threads/FB/X/dcard）仍 skip，交給 Gemini Google Search。
    一般靜態網頁走 HTML prefetch + BeautifulSoup 文字萃取。
    """
    urls = _extract_prefetch_urls(text)
    if not urls:
        return text

    blocks = []
    youtube_urls = [url for url in urls if _is_youtube_url(url)]
    other_urls = [url for url in urls if not _is_youtube_url(url)]
    ordered_urls = (youtube_urls + other_urls)[:_PREFETCH_MAX_URLS]
    for raw_url in ordered_urls:
        try:
            url = _clean_prefetch_url(raw_url)
            if not url:
                continue
            u_lower = url.lower()

            # 1) 影片平台：yt-dlp 優先（支援字幕），失敗才 fallback oEmbed
            if "tiktok.com" in u_lower:
                block = _fetch_video_ytdlp(url) or _fetch_tiktok_meta(url)
                # 內容太薄（< 300 chars）→ Gemini Video Understanding fallback
                block = _maybe_video_fallback(url, block)
                if block:
                    blocks.append(block)
                else:
                    logger.info("tiktok: ytdlp + oembed + gemini all failed, skip url=%s", url)
                continue
            if _is_youtube_url(url):
                # YouTube 不走 Gemini video fallback：yt-dlp/oEmbed/HTML metadata
                # 已能提供可用 context；全失敗時也保留辨識狀態避免 LLM 只看裸網址。
                blocks.append(_fetch_youtube_context(url))
                continue
            if "reddit.com" in u_lower or "redd.it" in u_lower:
                block = _fetch_reddit_meta(url)
                if block:
                    blocks.append(block)
                else:
                    logger.info("reddit .json failed, skip url=%s", url)
                continue

            # 2) Instagram Reels / Posts：yt-dlp → embed 頁面 → Gemini video → Google Search
            if "instagram.com" in u_lower:
                block = _fetch_video_ytdlp(url) or _fetch_instagram_embed(url)
                block = _maybe_video_fallback(url, block)
                if block:
                    blocks.append(block)
                else:
                    logger.info(
                        "instagram: all methods failed url=%s → Gemini Google Search",
                        url,
                    )
                continue

            # 3) 其他 JS 渲染網站（FB / X / Threads / dcard）：試 yt-dlp，再試 Gemini video，失敗才 Google Search
            if _JS_RENDERED_DOMAINS.search(url):
                block = _fetch_video_ytdlp(url)
                block = _maybe_video_fallback(url, block)
                if block:
                    blocks.append(block)
                else:
                    logger.info(
                        "prefetch skip (JS/CF site, ytdlp failed) url=%s → Gemini Google Search",
                        url,
                    )
                continue

            # 一般網頁：直接抓取 HTML
            resp = _requests.get(
                url,
                timeout=_PREFETCH_TIMEOUT,
                headers={
                    "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7)"
                },
            )
            resp.raise_for_status()
            soup = BeautifulSoup(resp.text, "html.parser")
            for tag in soup(["script", "style", "nav", "footer", "header"]):
                tag.decompose()
            content = soup.get_text(separator="\n", strip=True)

            # 內容太短 = JS 渲染空殼或 Cloudflare 頁面，不塞垃圾進 prompt
            if len(content) < _PREFETCH_MIN_CHARS:
                logger.info(
                    "prefetch skip (too short %d chars) url=%s", len(content), url
                )
                continue

            if len(content) > _PREFETCH_MAX_CHARS:
                content = content[:_PREFETCH_MAX_CHARS] + "\n…（內容截斷）"
            blocks.append(
                f"（以下是連結 {url} 的網頁內容，已預先擷取）\n"
                f"--- 網頁內容開始 ---\n{content}\n--- 網頁內容結束 ---"
            )
            logger.info("prefetch OK url=%s chars=%d", url, len(content))
        except Exception as e:
            logger.info("prefetch failed url=%s: %s", url, e)

    if blocks:
        return "\n\n".join(blocks) + "\n\n" + text
    return text


# ── Health ────────────────────────────────────────────────────────────────────


@app.get("/health")
def health():
    return {
        "status": "ok",
        "gemini_model": settings.gemini_model,
        "gemini_light_model": settings.gemini_light_model,
        "group_locked": bool(settings.allowed_group_ids),
        "allowed_group_count": len(settings.allowed_group_ids),
    }


_GENIMG_DIR = "/tmp/line_bot_genimg"
_GENIMG_FILENAME_RE = re.compile(r"^[a-f0-9]{32}\.png$")


@app.get("/static/img/{filename}")
def serve_genimg(filename: str):
    """Serve generated images for LINE ImageMessage（public via cloudflared）。

    路徑驗證：filename 必須是 UUID hex（32 字）+ .png，避免 path traversal。
    """
    from fastapi.responses import FileResponse
    from fastapi import HTTPException
    if not _GENIMG_FILENAME_RE.match(filename):
        raise HTTPException(status_code=400, detail="invalid filename")
    path = os.path.join(_GENIMG_DIR, filename)
    if not os.path.isfile(path):
        raise HTTPException(status_code=404, detail="not found")
    return FileResponse(path, media_type="image/png")


# ── Webhook ───────────────────────────────────────────────────────────────────


def _hash_text(value: str) -> str:
    import hashlib as _hl

    return _hl.sha256((value or "").encode("utf-8")).hexdigest()[:12]


def _safe_event_debug_dump(event) -> str:
    src = getattr(event, "source", None)
    msg = getattr(event, "message", None)
    payload = {
        "event_type": type(event).__name__,
        "source_type": type(src).__name__ if src else None,
        "group_id": getattr(src, "group_id", None) if src else None,
    }
    user_id = getattr(src, "user_id", None) if src else None
    if user_id:
        payload["user_id_sha256"] = _hash_text(user_id)
    if msg is not None:
        payload["message_type"] = getattr(msg, "type", type(msg).__name__)
        payload["message_id"] = getattr(msg, "id", None)
        text = getattr(msg, "text", None)
        if text is not None:
            payload["text_len"] = len(text)
            payload["text_sha256"] = _hash_text(text)
    return _json.dumps(payload, ensure_ascii=False, sort_keys=True)


@app.post("/callback")
async def callback(request: Request, x_line_signature: str = Header(None)):
    # N3 fix (2026-05-30): errors="replace" 避免非 UTF-8 垃圾請求在驗章前就
    # UnicodeDecodeError → 未處理的 500。簽章仍會擋掉偽造的合法訊息。
    body = (await request.body()).decode("utf-8", errors="replace")
    if not x_line_signature:
        client = request.client.host if request.client else "?"
        logger.warning("missing x-line-signature header from %s body_len=%d", client, len(body))
        raise HTTPException(status_code=400, detail="missing signature")
    # PII protection: never print raw signature or body; hashes are enough to correlate retries.
    if os.getenv("DEBUG_RAW_BODY") == "1":
        print(f"[RAW] debug=1 sig_sha256={_hash_text(x_line_signature)} len={len(body)} body_sha256={_hash_text(body)}", flush=True)
    else:
        print(f"[RAW] sig_sha256={_hash_text(x_line_signature)} len={len(body)} body_sha256={_hash_text(body)}", flush=True)
    try:
        events = _parser.parse(body, x_line_signature)
    except InvalidSignatureError:
        logger.warning("invalid signature")
        # N5 fix (2026-05-30): from None 斷開例外鏈，log 不再夾雜 InvalidSignatureError traceback
        raise HTTPException(status_code=400, detail="invalid signature") from None

    print(f"[PARSED] event_count={len(events)}", flush=True)
    for event in events:
        src = getattr(event, "source", None)
        gid = getattr(src, "group_id", None) if src else None
        print(
            f"[EVENT] type={type(event).__name__} source={type(src).__name__ if src else None} group_id={gid}",
            flush=True,
        )
        # Debug dump is sanitized: no raw message text or raw user_id.
        if os.getenv("DEBUG_EVENT_DUMP") == "1":
            try:
                print(f"[EVENT_DUMP] {_safe_event_debug_dump(event)}", flush=True)
            except Exception:
                print("[EVENT_DUMP] (could not dump sanitized event)", flush=True)
        try:
            _handle_event(event)
        except Exception as e:
            logger.exception("handle_event failed: %s", e)
    return {"ok": True}


def _handle_event(event) -> None:
    # MemberJoinedEvent / MemberLeftEvent: 其他成員進出群組
    if isinstance(event, (MemberJoinedEvent, MemberLeftEvent)):
        if os.getenv("DEBUG_EVENT_DUMP") == "1":
            try:
                print(
                    f"[MEMBER_EVT] {_safe_event_debug_dump(event)}",
                    flush=True,
                )
            except Exception:
                pass
        return

    # 只處理群組裡的 message
    if not isinstance(event, MessageEvent):
        return
    if not isinstance(event.source, GroupSource):
        return

    # ── 重送事件處理 ──────────────────────────────────────────────────────
    # 隧道斷線期間 LINE 送不進來的訊息,恢復後會帶 is_redelivery=True 重新投遞。
    # 不能無腦跳過 — 查 raw_messages,如果我們從沒收過,就要補處理。
    # reply_token 一定已過期,但 _reply 會 fallback 到 push_message。
    dctx = getattr(event, "delivery_context", None)
    is_redelivery = bool(dctx and getattr(dctx, "is_redelivery", False))
    if is_redelivery:
        msg = getattr(event, "message", None)
        msg_id = getattr(msg, "id", None) if msg else None
        gid = event.source.group_id
        if msg_id and gid:
            existing = memory.get_raw_message(gid, msg_id)
            if existing is not None:
                logger.info("skip truly-duplicate redelivery msg_id=%s", msg_id)
                return
            logger.info("processing missed redelivery msg_id=%s group=%s", msg_id, gid)
        else:
            logger.info("skip redelivered event (no msg_id)")
            return

    group_id = event.source.group_id

    # 多群白名單（2026-05-27 multi-group）：空 list = unlock mode 接受所有群（log only）；
    # 非空時只接受清單內。`allowed_group_ids` 是 settings 的 @computed_field：
    # ALLOWED_GROUP_IDS env (csv) > legacy ALLOWED_GROUP_ID（單數）> []。
    if settings.allowed_group_ids and group_id not in settings.allowed_group_ids:
        logger.info("ignoring message from non-allowed group_id=%s", group_id)
        return

    msg = event.message
    sender_user_id = getattr(event.source, "user_id", None)

    # ── quota 爆時：非文字內容不排 pending reply ────────────────────────
    # Andrew 2026-06-06 指示取消 pending 回覆機制；File/Audio 在額度爆時
    # 不再寫入 pending_explicit_reply.json，也不做 reply-token 補回。
    if _quota_exhausted() and isinstance(
        msg, (FileMessageContent, AudioMessageContent)
    ):
        if _pending_reply_enabled():
            _save_pending_any(event, group_id, sender_user_id, msg)
            _try_piggyback_drain_with_reply_token(event.reply_token, group_id)
        else:
            logger.info(
                "pending reply disabled; dropped quota-exhausted %s group=%s",
                type(msg).__name__, group_id,
            )
        return

    # 文字：先記進 raw_messages（供 quote 回查 / burst look-back / Layer 2 抓 trigger）
    if isinstance(msg, TextMessageContent):
        text_body = msg.text or ""
        memory.log_raw_message(group_id, msg.id, sender_user_id, text_body)
        # Piggyback reminder：有人留言時仍可補 reminder；reply pending 已由開關停用。
        _spawn_piggyback_drain(group_id)
        _handle_text_message(event, group_id)
        return

    # Piggyback reminder：非文字內容仍可補 reminder；文字在上方先做衝突 gate。
    _spawn_piggyback_drain(group_id)

    # 2026-06-06 改：圖片/影片只走即時 local analyze；不再存 pending reply。
    if isinstance(msg, ImageMessageContent):
        memory.log_raw_message(group_id, msg.id, sender_user_id, "[圖片]")
        memory.log_raw_message_meta(
            group_id, msg.id, media_type="image", mime_type="image/jpeg"
        )
        if _pending_reply_enabled():
            _save_pending_any(event, group_id, sender_user_id, msg)
        _MEDIA_EXECUTOR.submit(_handle_image_message, event, group_id)
        return
    if isinstance(msg, VideoMessageContent):
        memory.log_raw_message(group_id, msg.id, sender_user_id, "[影片]")
        memory.log_raw_message_meta(
            group_id, msg.id, media_type="video", mime_type="video/mp4"
        )
        if _pending_reply_enabled():
            _save_pending_any(event, group_id, sender_user_id, msg)
        _MEDIA_EXECUTOR.submit(_handle_video_message, event, group_id)
        return
    if isinstance(msg, AudioMessageContent):
        memory.log_raw_message(group_id, msg.id, sender_user_id, "[音訊]")
        memory.log_raw_message_meta(
            group_id, msg.id, media_type="audio", mime_type="audio/m4a"
        )
        _handle_audio_message(event, group_id)
        return

    # 檔案：只處理文字檔，其他婉拒（file 很罕見，不是 burst 的一部分）
    if isinstance(msg, FileMessageContent):
        file_name = getattr(msg, "file_name", "") or "unknown"
        memory.log_raw_message(group_id, msg.id, sender_user_id, f"[檔案: {file_name}]")
        memory.log_raw_message_meta(
            group_id,
            msg.id,
            media_type="file",
            mime_type=_guess_mime_type(file_name),
            file_name=file_name,
        )
        _handle_file_message(event, group_id)
        return

    # Catch-all for unknown message types (Sticker / Location / Template / etc.).
    # 只保留 raw audit；不再存 pending reply，也不送低價值 system fallback。
    logger.info(
        "unknown message type %s group=%s — ignored without pending reply",
        type(msg).__name__, group_id,
    )
    # I7 fix (2026-05-30): 補 log_raw_message（其他分支都有，唯獨 catch-all 漏）。
    # redelivery 去重靠 raw_messages，沒記 → 貼圖/位置等類型 redelivery 會被當「沒收過」
    # 重複回覆；引用該訊息問後續問題也查不到原文。
    try:
        msg_id = getattr(msg, "id", None)
        if msg_id:
            memory.log_raw_message(group_id, msg_id, sender_user_id, f"[{type(msg).__name__}]")
    except Exception as e:
        logger.warning("log_raw_message for unknown msg type failed: %s", e)
    if _pending_reply_enabled():
        try:
            _save_pending_any(event, group_id, sender_user_id, msg)
        except Exception as e:
            logger.warning("save pending for unknown msg type failed: %s", e)


def _handle_audio_message(event: MessageEvent, group_id: str) -> None:
    """語音留言自動分析 — 不需要 @mention，下載後直接丟 Gemini 轉寫 + 回應。"""
    if _quota_exhausted():
        return
    try:
        data = _download_content(event.message.id)
    except Exception as e:
        logger.warning("download audio failed: %s", e)
        return
    if len(data) > _MEDIA_BYTE_LIMIT:
        return
    parts = [
        types.Part.from_bytes(data=bytes(data), mime_type="audio/m4a"),
        "(群組成員傳了一段語音留言，請先完整轉寫內容，再根據系統指令判斷是否有查核或回應價值。若只是閒聊請用一兩句自然回應即可。)",
    ]
    context = memory.get_context(group_id)
    facts = memory.top_facts(group_id)
    pnotes = _get_persona_notes(group_id)
    try:
        with _thinking_indicator(group_id):
            reply_text = _llm_chat(parts, context, facts, pnotes)
    except Exception as e:
        if _is_quota_error(e):
            _mark_quota_exhausted()
        else:
            logger.exception("gemini chat (audio) failed: %s", e)
        return
    if not reply_text or not reply_text.strip():
        return
    memory.log_raw_message_meta(
        group_id,
        event.message.id,
        media_type="audio",
        mime_type="audio/m4a",
        description=reply_text,
    )
    memory.append_turn(group_id, "user", "[語音留言]")
    _append_bot_turn(group_id, reply_text)
    _maybe_extract_facts(group_id)
    _reply(event.reply_token, reply_text, group_id=group_id)


def _handle_image_message(event, group_id):
    """Pure-local image handler. Runs in _MEDIA_EXECUTOR (cap 2 concurrent).

    Flow: download → media_pipeline.analyze_image (50s hard timeout per GP1 §2,
    because reply_token TTL is ~1 min) → reply → remove from pending. On any
    failure, leave the pending entry for the cron retry worker to drain later.
    """
    import threading as _threading
    msg_id = event.message.id
    try:
        data = _download_content(msg_id)
    except Exception as e:
        logger.warning("image download failed for handler: %s", e)
        return
    if len(data) > _MEDIA_BYTE_LIMIT:
        logger.info("image too large for handler: %d bytes", len(data))
        return

    holder = {"reply": None}

    def _run():
        try:
            import media_pipeline
            holder["reply"] = media_pipeline.analyze_image(bytes(data), group_id=group_id)
        except Exception as e:
            logger.warning("image analyze failed: %s", e)

    t = _threading.Thread(target=_run, daemon=True)
    t.start()
    t.join(timeout=50)
    if t.is_alive():
        logger.info("image analyze timeout (>50s), leaving pending for retry worker")
        return
    reply_text = holder["reply"]
    if not reply_text or not reply_text.strip():
        logger.info("image analyze returned empty, leaving pending")
        return
    memory.log_raw_message_meta(
        group_id, msg_id, media_type="image", mime_type="image/jpeg", description=reply_text
    )
    memory.append_turn(group_id, "user", "[圖片]")
    _append_bot_turn(group_id, reply_text)
    try:
        _reply(event.reply_token, reply_text, group_id=group_id)
        _remove_pending_by_msg_id(group_id, msg_id)
    except Exception as e:
        logger.warning("image reply failed: %s", e)


def _handle_video_message(event, group_id):
    """Pure-local video handler. Same shape as _handle_image_message."""
    import threading as _threading
    msg_id = event.message.id
    try:
        data = _download_content(msg_id)
    except Exception as e:
        logger.warning("video download failed for handler: %s", e)
        return
    if len(data) > _MEDIA_BYTE_LIMIT:
        logger.info("video too large for handler: %d bytes", len(data))
        return

    holder = {"reply": None}

    def _run():
        try:
            import media_pipeline
            holder["reply"] = media_pipeline.analyze_video(bytes(data), group_id=group_id)
        except Exception as e:
            logger.warning("video analyze failed: %s", e)

    t = _threading.Thread(target=_run, daemon=True)
    t.start()
    t.join(timeout=50)
    if t.is_alive():
        logger.info("video analyze timeout (>50s), leaving pending for retry worker")
        return
    reply_text = holder["reply"]
    if not reply_text or not reply_text.strip():
        logger.info("video analyze returned empty, leaving pending")
        return
    memory.log_raw_message_meta(
        group_id, msg_id, media_type="video", mime_type="video/mp4", description=reply_text
    )
    memory.append_turn(group_id, "user", "[影片]")
    _append_bot_turn(group_id, reply_text)
    try:
        _reply(event.reply_token, reply_text, group_id=group_id)
        _remove_pending_by_msg_id(group_id, msg_id)
    except Exception as e:
        logger.warning("video reply failed: %s", e)


def _try_piggyback_reminders_fast_path(
    reply_token: str | None, group_id: str
) -> bool:
    """Use an ordinary group message's reply_token to deliver due reminders.

    LINE push_message can hit monthly 429 quota; reply_message does not use that
    quota. Stages are marked only after reply_message succeeds.
    """
    if (
        not _reminder_reply_piggyback_enabled()
        or not reply_token
        or not group_id
        or settings.bot_muted
    ):
        return False
    try:
        import calendar_db
        import event_reminder as _er
        messages: list = []
        pending: list[tuple[str, int]] = []
        for offset in calendar_db.REMINDER_OFFSETS:
            if len(messages) >= 5:
                break
            for e in calendar_db.list_due_for_reminder(group_id, days_ahead=offset):
                if len(messages) >= 5:
                    break
                spec = _er.build_reminder_message_spec(e, offset, allow_mention=True)
                if spec is None:
                    continue
                messages.append(_er.sdk_message_from_spec(spec))
                pending.append((e["event_id"], offset))
        import reminder_push as _rp
        pending_pushes: list[tuple[int, str]] = []
        if len(messages) < 5:
            remaining = 5 - len(messages)
            for item in _rp.due_reminders_for_reply(group_id, limit=remaining):
                messages.append(
                    item.get("message") or TextMessage(text=item["text"][:5000])
                )
                pending_pushes.append((item["reminder_id"], item["stage"]))
        if not messages:
            return False
        try:
            with ApiClient(_get_line_config()) as api_client:
                MessagingApi(api_client).reply_message(
                    ReplyMessageRequest(
                        reply_token=reply_token, messages=messages,
                    )
                )
        except Exception as e:
            logger.warning(
                "reminder fast-path reply failed (reminders preserved): %s",
                str(e)[:200],
            )
            return False
        for ev_id, off in pending:
            calendar_db.mark_reminded(ev_id, off, group_id)
        if pending_pushes:
            _rp.mark_reminders_pushed(pending_pushes)
        logger.info(
            "reminder fast-path: pushed %d calendar + %d reminder_push via reply_token group=%s",
            len(pending), len(pending_pushes), group_id,
        )
        return True
    except Exception as e:
        logger.warning("reminder fast-path failed: %s", e)
        return False


# 自動 capture cheap pre-filter — 含日期 hint + 行程動詞 才考慮跑 Gemini extractor
_AUTO_CAPTURE_DATE_HINT_RE = re.compile(
    r"明天|後天|大後天|今天|這週末|下週末|下週|本週|"
    r"(?:星期|週|周|禮拜)[一二三四五六日天]|"
    r"\d+月\d+日|\d{4}-\d{2}-\d{2}|\d{1,2}/\d{1,2}|"
    r"(?:[01]?\d|2[0-3]):[0-5]\d"
)
_AUTO_CAPTURE_VERB_RE = re.compile(
    r"聚餐|生日|出遊|看(?:.{0,12})?(?:醫生|醫師|牙醫|牙醫師)|牙醫|"
    r"拿(?:蛋糕|藥|包裹|貨|禮物|花)|"
    r"接(?:爸|媽|妹|弟|姊|爺爺|奶奶|小孩|小朋友)|"
    r"陪(?:.{0,5})(?:就醫|看醫生|看病)|"
    r"回(?:台北|新北|台中|台南|高雄|花蓮|宜蘭|新竹|苗栗|嘉義|屏東|台東|老家|家)|"
    r"做(?:胃鏡|大腸鏡|健康檢查|體檢|手術|健檢|LDCT)|"
    r"領(?:藥|處方簽|包裹)|"
    r"婚禮|喜宴|滿月|彌月|"
    r"打(?:疫苗|球|羽球)|羽球|行程|活動|抽血|健檢|出差|北上|南下"
)

_MEDICAL_ACTOR_ACTION_RE = re.compile(
    r"看(?:.{0,12})?(?:醫生|醫師|牙醫|牙醫師)|牙醫|牙醫師|"
    r"就醫|看病|回診|掛(?:號|醫生|醫師)|"
    r"做(?:胃鏡|大腸鏡|健康檢查|體檢|手術|健檢|LDCT|MRI|"
    r"正子斷層掃描|正子斷層|電腦斷層|(?<![A-Za-z])CT(?![A-Za-z])|核磁共振)|"
    r"MRI|核磁共振|正子斷層掃描|正子斷層|PET-?CT|PET|電腦斷層|"
    r"(?<![A-Za-z])CT(?![A-Za-z])|"
    r"打疫苗|抽血|領(?:藥|處方簽)"
)
_FAMILY_ACTOR_TERMS = (
    "媽媽", "爸爸", "姊姊", "姐姐", "妹妹", "弟弟", "哥哥",
    "爺爺", "奶奶", "聖雅", "聖穎", "黃聖雅", "黃聖穎", "黃將修", "全家",
)
_FAMILY_ACTOR_NORMALIZE = {
    "姐姐": "姊姊",
    "妹妹": "黃聖雅",
    "聖雅": "黃聖雅",
    "聖穎": "黃聖穎",
}


def _normalize_family_actor(name: str) -> str:
    return _FAMILY_ACTOR_NORMALIZE.get(name, name)


def _family_name_variants(name: str) -> set[str]:
    variants = {name}
    for raw, normalized in _FAMILY_ACTOR_NORMALIZE.items():
        if normalized == name:
            variants.add(raw)
    return variants


def _family_name_in_text(text: str, name: str) -> bool:
    return any(variant in text for variant in _family_name_variants(name))


def _alias_from_user_id(user_id: str | None) -> str:
    """Best-effort LINE user_id -> family alias using the existing local map."""
    if not user_id:
        return ""
    try:
        import food_extractor
        return food_extractor.alias_from_user_id(user_id) or ""
    except Exception as e:
        logger.debug("alias lookup skipped: %s", e)
        return ""


def _event_actor_role(event_type: str | None) -> str:
    if event_type == "personal_trip":
        return "旅者"
    if event_type == "medical":
        return "就醫"
    return "主角"


def _participants_contain_actor(participants: list[str], actor: str) -> bool:
    return any(actor in p for p in participants)


def _has_specific_family_actor(text: str, ignored: set[str] | None = None) -> bool:
    ignored = ignored or set()
    for term in _FAMILY_ACTOR_TERMS:
        normalized = _normalize_family_actor(term)
        if normalized in ignored:
            continue
        if _family_name_in_text(text, normalized):
            return True
    return False


def _apply_sender_first_person_event(data: dict, sender_user_id: str | None) -> None:
    """Resolve first-person event wording to the LINE sender alias.

    Calendar extraction sees only text, so family shorthand like "我到高雄"
    can otherwise be persisted as "我". The sender id is available in webhook
    context; use the local alias map as a deterministic post-process.
    """
    actor = _alias_from_user_id(sender_user_id)
    if not actor:
        return

    title = str(data.get("title") or "")
    title_changed = False
    if title:
        new_title = re.sub(r"^我的", f"{actor}的", title)
        new_title = re.sub(r"^我(?=[^\s，。！？!?、,.;])", actor, new_title)
        title_changed = new_title != title
        data["title"] = new_title

    participants = [str(p) for p in (data.get("participants") or []) if p]
    new_parts: list[str] = []
    participant_changed = False
    for part in participants:
        new_part = re.sub(r"^我(?=(?:\(|（|$))", actor, part)
        participant_changed = participant_changed or new_part != part
        new_parts.append(new_part)

    if (title_changed or participant_changed) and not _participants_contain_actor(
        new_parts, actor
    ):
        new_parts.insert(0, f"{actor}({_event_actor_role(data.get('event_type'))})")
    data["participants"] = new_parts


def _apply_family_context_defaults(data: dict, combined_text: str) -> None:
    """Apply stable family-specific defaults after model extraction."""
    haystack = " ".join(
        str(x or "")
        for x in (
            combined_text,
            data.get("title"),
            data.get("location"),
        )
    )
    if "東海" not in haystack:
        return
    participants = [str(p) for p in (data.get("participants") or []) if p]
    if _has_specific_family_actor(str(data.get("title") or ""), {"爸爸", "全家"}):
        return
    if any(_has_specific_family_actor(p, {"爸爸", "全家"}) for p in participants):
        return
    if not _participants_contain_actor(participants, "爸爸"):
        data["participants"] = ["爸爸(東海相關)", *participants]


def _infer_medical_actor(text: str, user_id: str | None = None) -> str | None:
    """Infer who the medical reminder/event is for.

    Explicit family names win. For subjectless medical messages, default to the
    sender alias because family chat shorthand like "星期四看牙醫" usually means
    the speaker's own appointment.
    """
    if not text or not _MEDICAL_ACTOR_ACTION_RE.search(text):
        return None
    alias = _alias_from_user_id(user_id)
    companions = set(_infer_medical_companions(text))
    if alias and re.search(r"陪\s*我|陪我|陪同我", text):
        return alias
    for term in _FAMILY_ACTOR_TERMS:
        normalized = _normalize_family_actor(term)
        if normalized in companions:
            continue
        if term in text:
            return normalized
    if alias:
        return alias
    return None


def _infer_medical_companions(text: str) -> list[str]:
    """Return people explicitly described as accompanying, not receiving care."""
    if not text:
        return []
    companions: list[str] = []
    for term in _FAMILY_ACTOR_TERMS:
        normalized = _normalize_family_actor(term)
        if normalized == "全家":
            continue
        escaped = re.escape(term)
        patterns = (
            rf"{escaped}[^。！？\n]{{0,8}}陪(?:我|你|他|她|同|診|著|去|看|就醫)",
            rf"(?:由|請|給|找){escaped}[^。！？\n]{{0,4}}陪",
        )
        if any(re.search(pattern, text) for pattern in patterns):
            if normalized not in companions:
                companions.append(normalized)
    return companions


def _has_family_actor(text: str, ignore_names: list[str] | None = None) -> bool:
    ignored = set(ignore_names or [])
    for term in _FAMILY_ACTOR_TERMS:
        normalized = _normalize_family_actor(term)
        if normalized in ignored:
            continue
        if term in text:
            return True
    return False


def _apply_medical_actor(
    action: str, actor: str | None, companions: list[str] | None = None
) -> str:
    if not action:
        return action
    companions = [name for name in (companions or []) if name and name != actor]
    if actor:
        action = re.sub(r"^我的", f"{actor}的", action)
        action = re.sub(r"^我(?=看|做|掛|回診|就醫|領|打|抽)", actor, action)
        if not _has_family_actor(action, ignore_names=companions):
            action = f"{actor}{action}"
    missing_companions = [
        name for name in companions if not _family_name_in_text(action, name)
    ]
    if missing_companions:
        action = f"{action}（{'、'.join(missing_companions)}陪同）"
    return action


def _medical_mention_aliases(
    actor: str | None, companions: list[str] | None = None
) -> list[str]:
    aliases: list[str] = []
    for name in [actor, *(companions or [])]:
        if name and name not in aliases:
            aliases.append(name)
    return aliases


def _with_medical_actor_participant(participants: list | None, actor: str | None) -> list:
    parts = [str(p) for p in (participants or []) if p]
    if not actor:
        return parts
    if any(actor in p for p in parts):
        return parts
    if actor == "全家":
        return [actor, *parts]
    return [f"{actor}(就醫)", *parts]


def _capture_calendar_events_regex_only(
    group_id: str,
    combined_text: str,
    sender_user_id: str = "",
    message_id: str = "",
) -> int:
    """Persist explicit date/time family events without Gemini.

    This path runs before reply generation, so lite/local mode cannot say an
    event was recorded while the DB still lacks the corresponding reminder.
    """
    try:
        import calendar_db
        import calendar_regex

        today_tw = datetime.now(ZoneInfo("Asia/Taipei")).date()
        events = calendar_regex.extract_many_regex_only(
            combined_text, today_tw, require_time=True
        )
        inserted = 0
        for data in events:
            if not (data.get("has_event") and data.get("title") and data.get("date")):
                continue
            if data.get("event_type") == "medical":
                actor = _infer_medical_actor(combined_text, sender_user_id)
                data["title"] = _apply_medical_actor(data["title"], actor)
                data["participants"] = _with_medical_actor_participant(
                    data.get("participants"), actor
                )
            else:
                _apply_sender_first_person_event(data, sender_user_id)
            _apply_family_context_defaults(data, combined_text)
            event_id = calendar_db.insert_event(
                group_id=group_id,
                title=data["title"],
                event_date=data["date"],
                event_time=data.get("time"),
                location=data.get("location"),
                participants=data.get("participants") or [],
                event_type=data.get("event_type", "family_gathering"),
                source_msg_id=message_id,
            )
            if event_id:
                inserted += 1
                logger.info(
                    "calendar event captured by local regex: %s '%s' on %s type=%s (group=%s)",
                    event_id,
                    data["title"],
                    data["date"],
                    data.get("event_type", "family_gathering"),
                    group_id,
                )
        return inserted
    except Exception as e:
        logger.warning("local calendar regex capture failed: %s", e)
        return 0


def _auto_capture_text_if_important(
    group_id: str,
    text: str,
    sender_user_id: str = "",
    message_id: str = "",
) -> None:
    """每條 text message 來時 cheap pre-filter → 通過才 async 跑 Gemini extractor。

    避免每訊息都打 Gemini 燒 20 req/day quota。
    UNIQUE INDEX 自動 dedup，重跑安全。
    """
    if not text or len(text.strip()) < 4 or len(text) > 500:
        return
    if not (_AUTO_CAPTURE_DATE_HINT_RE.search(text) and _AUTO_CAPTURE_VERB_RE.search(text)):
        return
    if _capture_calendar_events_regex_only(group_id, text, sender_user_id, message_id):
        return
    import threading

    def _bg() -> None:
        try:
            _maybe_capture_calendar_event(group_id, text, sender_user_id, message_id)
        except Exception as e:
            logger.warning("auto capture (every-msg) failed: %s", e)

    threading.Thread(target=_bg, daemon=True).start()


_CALENDAR_CORRECTION_MARKER_RE = re.compile(
    r"(?:更正|修正|校正|改成|改為|改到|改在|改至)"
)
_CALENDAR_CORRECTION_NEGATION_RE = re.compile(
    r"(?:不用|不要|先別|先不要|暫時不要).{0,6}(?:更正|修正|校正|改)"
)
_ZH_NUMERAL = {
    "零": 0,
    "〇": 0,
    "一": 1,
    "二": 2,
    "兩": 2,
    "三": 3,
    "四": 4,
    "五": 5,
    "六": 6,
    "七": 7,
    "八": 8,
    "九": 9,
}
_CALENDAR_CORRECTION_KEYWORDS: tuple[str, ...] = (
    "羽球", "打球", "聚餐", "吃飯", "生日", "出遊", "旅行", "回台北",
    "回老家", "看醫生", "牙醫", "洗牙", "拿蛋糕", "蛋糕", "包裹",
    "接媽媽", "接爸爸", "媽媽", "爸爸", "姊姊", "妹妹", "弟弟", "全家",
)
_CALENDAR_CORRECTION_SYNONYMS: dict[str, tuple[str, ...]] = {
    "羽球": ("羽球", "打球"),
    "打球": ("打球", "羽球"),
    "洗牙": ("洗牙", "牙醫"),
    "牙醫": ("牙醫", "洗牙"),
}
_CALENDAR_CORRECTION_ACTORS: tuple[str, ...] = (
    "全家", "媽媽", "爸爸", "姊姊", "姐姐", "妹妹", "弟弟", "哥哥",
    "爺爺", "奶奶", "黃聖雅", "黃聖穎", "聖雅", "聖穎",
)


def _parse_zh_int(raw: str | None) -> int | None:
    s = (raw or "").strip()
    if not s:
        return None
    if s.isdigit():
        return int(s)
    if s in _ZH_NUMERAL:
        return _ZH_NUMERAL[s]
    if s == "十":
        return 10
    if "十" in s:
        left, _, right = s.partition("十")
        tens = 1 if not left else _ZH_NUMERAL.get(left)
        ones = 0 if not right else _ZH_NUMERAL.get(right)
        if tens is None or ones is None:
            return None
        return tens * 10 + ones
    return None


def _parse_calendar_correction_time(text: str) -> str | None:
    s = (text or "").strip()
    if not s:
        return None

    m = re.search(r"(?<!\d)([01]?\d|2[0-3])\s*[:：]\s*([0-5]\d)(?!\d)", s)
    if m:
        return f"{int(m.group(1)):02d}:{int(m.group(2)):02d}"

    cn = "零〇一二兩三四五六七八九十"
    m = re.search(
        rf"(凌晨|早上|上午|中午|下午|晚上|傍晚|晚間)?\s*"
        rf"(\d{{1,2}}|[{cn}]{{1,3}})\s*(?:點|時)\s*"
        rf"(半|(?:\d{{1,2}}|[{cn}]{{1,3}})\s*分?)?",
        s,
    )
    if m:
        period = m.group(1) or ""
        hour = _parse_zh_int(m.group(2))
        minute_raw = (m.group(3) or "").replace("分", "").strip()
        minute = 30 if minute_raw == "半" else (_parse_zh_int(minute_raw) or 0)
        if hour is None or hour > 23 or minute > 59:
            return None
        if period in ("下午", "晚上", "傍晚", "晚間") and 1 <= hour <= 11:
            hour += 12
        elif period == "中午" and hour == 0:
            hour = 12
        return f"{hour:02d}:{minute:02d}"

    m = re.search(
        r"(?<![\d年月日/\-])([01]?\d|2[0-3])([0-5]\d)(?![\d年月日/\-])",
        s,
    )
    if m:
        return f"{int(m.group(1)):02d}:{int(m.group(2)):02d}"
    return None


def _parse_calendar_absolute_date(text: str):
    s = text or ""
    today = datetime.now(ZoneInfo("Asia/Taipei")).date()
    patterns = (
        r"(?P<y>\d{4})[-/](?P<m>\d{1,2})[-/](?P<d>\d{1,2})",
        r"(?P<m>\d{1,2})\s*月\s*(?P<d>\d{1,2})\s*[日號]?",
        r"(?P<m>\d{1,2})/(?P<d>\d{1,2})",
    )
    for pat in patterns:
        m = re.search(pat, s)
        if not m:
            continue
        try:
            year = int(m.groupdict().get("y") or today.year)
            return datetime(year, int(m.group("m")), int(m.group("d"))).date()
        except ValueError:
            return None
    return None


def _resolve_calendar_correction_date(text: str):
    return _parse_calendar_absolute_date(text) or _resolve_relative_date(text)


def _calendar_correction_keywords(text: str) -> list[str]:
    found: list[str] = []
    candidates = [*_CALENDAR_CORRECTION_KEYWORDS]
    try:
        candidates.extend(_QUERY_NOUN_KEYWORDS)
    except NameError:
        pass
    for kw in candidates:
        if kw and kw in text and kw not in found:
            found.append(kw)
            for syn in _CALENDAR_CORRECTION_SYNONYMS.get(kw, ()):
                if syn not in found:
                    found.append(syn)

    if found:
        return found

    marker = _CALENDAR_CORRECTION_MARKER_RE.search(text)
    prefix = text[: marker.start()] if marker else text
    cleaned = re.sub(
        r"(今天|明天|後天|大後天|這週|本週|下週|週[一二三四五六日天]|"
        r"\d{4}[-/]\d{1,2}[-/]\d{1,2}|\d{1,2}/\d{1,2}|"
        r"\d{1,2}\s*月\s*\d{1,2}\s*[日號]?|"
        r"[早上上午中午下午晚上傍晚晚間凌晨半點時分\d:：\s，,。.!！?？、])",
        "",
        prefix,
    )
    cleaned = cleaned.strip()
    return [cleaned] if len(cleaned) >= 2 else []


def _calendar_correction_content_candidate(text: str) -> str:
    s = (text or "").strip()
    if not s:
        return ""
    s = _CALENDAR_CORRECTION_MARKER_RE.sub(" ", s)
    s = re.sub(r"^[\s為成到在至是:：]+", " ", s)
    s = re.sub(
        r"(今天|明天|後天|大後天|這週|本週|下週|週[一二三四五六日天]|"
        r"\d{4}[-/]\d{1,2}[-/]\d{1,2}|\d{1,2}/\d{1,2}|"
        r"\d{1,2}\s*月\s*\d{1,2}\s*[日號]?)",
        " ",
        s,
    )
    s = re.sub(r"(?<!\d)([01]?\d|2[0-3])\s*[:：]\s*([0-5]\d)(?!\d)", " ", s)
    s = re.sub(r"(?<![\d年月日/\-])([01]?\d|2[0-3])([0-5]\d)(?![\d年月日/\-])", " ", s)
    s = re.sub(
        r"(凌晨|早上|上午|中午|下午|晚上|傍晚|晚間)?\s*"
        r"[\d零〇一二兩三四五六七八九十]{1,3}\s*(?:點|時)"
        r"(?:半|[\d零〇一二兩三四五六七八九十]{1,3}\s*分?)?",
        " ",
        s,
    )
    s = re.sub(r"[\s，,。.!！?？、]+", "", s)
    return s if len(s) >= 2 else ""


def _calendar_correction_title_candidate(text: str) -> str:
    s = (text or "").strip()
    marker = _CALENDAR_CORRECTION_MARKER_RE.search(s)
    if not marker:
        return ""
    suffix = _calendar_correction_content_candidate(s[marker.end():])
    if suffix:
        return suffix
    return _calendar_correction_content_candidate(s[: marker.start()])


def _calendar_correction_new_title(text: str, target_title: str | None) -> str | None:
    candidate = _calendar_correction_title_candidate(text)
    if not candidate:
        return None
    target = (target_title or "").strip()
    if target and candidate in target:
        return None
    if target and candidate == "羽球" and "打球" in target and "羽球" not in target:
        return target.replace("打球", "羽球")
    if target and candidate == "洗牙" and "牙醫" in target and "洗牙" not in target:
        return target.replace("看牙醫", "洗牙").replace("牙醫", "洗牙")
    if target:
        for actor in _CALENDAR_CORRECTION_ACTORS:
            normalized_actor = _normalize_family_actor(actor)
            if normalized_actor in target and normalized_actor not in candidate:
                if len(candidate) <= 8:
                    return f"{normalized_actor}{candidate}"
                break
    return candidate


def _parse_calendar_correction(text: str) -> dict | None:
    s = (text or "").strip()
    if not s or _CALENDAR_CORRECTION_NEGATION_RE.search(s):
        return None
    marker = _CALENDAR_CORRECTION_MARKER_RE.search(s)
    if not marker:
        return None

    prefix = s[: marker.start()]
    suffix = s[marker.end():]
    suffix_date = _resolve_calendar_correction_date(suffix)
    new_date = suffix_date
    target_date = _resolve_calendar_correction_date(prefix)
    if target_date is None and new_date is None:
        target_date = _resolve_calendar_correction_date(s)
    if new_date is None:
        new_date = target_date
    new_time = _parse_calendar_correction_time(
        suffix
    ) or _parse_calendar_correction_time(s)
    title_candidate = _calendar_correction_title_candidate(s)
    if new_date is None and new_time is None and not title_candidate:
        return None
    keywords = _calendar_correction_keywords(prefix or s)
    if not keywords and target_date is None and not title_candidate:
        return None
    return {
        "target_date": target_date.isoformat() if target_date else None,
        "new_date": new_date.isoformat() if new_date else None,
        "new_date_explicit": suffix_date is not None,
        "new_time": new_time,
        "new_title_raw": title_candidate,
        "keywords": keywords,
    }


def _find_calendar_correction_event(
    group_id: str, keywords: list[str], target_date: str | None
) -> dict | None:
    import calendar_db

    seen: set[str] = set()
    candidates: list[dict] = []
    for kw in keywords:
        try:
            rows = calendar_db.search_by_keyword(group_id, [kw], limit=10)
        except Exception as e:
            logger.debug("calendar correction keyword search failed kw=%s: %s", kw, e)
            rows = []
        for row in rows:
            ev_id = str(row.get("event_id") or "")
            if ev_id and ev_id not in seen:
                seen.add(ev_id)
                candidates.append(row)

    if target_date:
        dated = [row for row in candidates if row.get("event_date") == target_date]
        if dated:
            return dated[0]
    if candidates:
        return candidates[0]

    if target_date:
        try:
            events = calendar_db.list_upcoming(group_id, days=90) + calendar_db.list_past(
                group_id, days=30
            )
            same_day = [row for row in events if row.get("event_date") == target_date]
            if len(same_day) == 1:
                return same_day[0]
        except Exception as e:
            logger.debug("calendar correction same-day fallback failed: %s", e)
    return None


def _reminder_matches_calendar_correction(
    reminder: dict,
    keywords: list[str],
    target_date: str | None,
) -> bool:
    if target_date:
        try:
            rdate = datetime.fromtimestamp(
                int(reminder["remind_at"]), tz=ZoneInfo("Asia/Taipei")
            ).date().isoformat()
        except Exception:
            return False
        if rdate != target_date:
            return False
    if not keywords:
        return True
    haystack = f"{reminder.get('action') or ''} {reminder.get('source_text') or ''}"
    return any(kw and kw in haystack for kw in keywords)


def _update_calendar_correction_reminders(
    group_id: str,
    keywords: list[str],
    target_date: str | None,
    new_date: str | None,
    new_time: str | None,
    new_action: str | None,
    source_text: str,
) -> int:
    try:
        reminders = memory.list_pending_reminders(group_id)
    except Exception as e:
        logger.warning("calendar correction reminder lookup failed: %s", e)
        return 0

    updated = 0
    for reminder in reminders:
        if not _reminder_matches_calendar_correction(reminder, keywords, target_date):
            continue
        try:
            old_dt = datetime.fromtimestamp(
                int(reminder["remind_at"]), tz=ZoneInfo("Asia/Taipei")
            )
            date_part = (
                datetime.fromisoformat(new_date).date() if new_date else old_dt.date()
            )
            if new_time:
                hour_s, minute_s = new_time.split(":", 1)
                hour, minute = int(hour_s), int(minute_s)
            else:
                hour, minute = old_dt.hour, old_dt.minute
            new_dt = datetime(
                date_part.year,
                date_part.month,
                date_part.day,
                hour,
                minute,
                tzinfo=ZoneInfo("Asia/Taipei"),
            )
        except Exception as e:
            logger.debug("calendar correction reminder compose failed: %s", e)
            continue
        if memory.update_reminder_schedule(
            int(reminder["reminder_id"]),
            int(new_dt.timestamp()),
            source_text=source_text[:200],
            action=new_action,
        ):
            updated += 1
    return updated


def _try_handle_calendar_correction(
    event: MessageEvent, group_id: str, text: str
) -> bool:
    correction = _parse_calendar_correction(text)
    if correction is None:
        return False

    import calendar_db

    keywords = list(correction["keywords"])
    target = _find_calendar_correction_event(
        group_id, keywords, correction.get("target_date")
    )
    if target:
        for kw in (target.get("title") or "", target.get("location") or ""):
            if kw and kw not in keywords:
                keywords.append(kw)
    new_title = _calendar_correction_new_title(text, (target or {}).get("title"))
    if new_title and new_title not in keywords:
        keywords.append(new_title)

    if target and not correction.get("new_date_explicit"):
        target_date = target.get("event_date")
        new_date = target.get("event_date")
    else:
        target_date = correction.get("target_date") or (target or {}).get("event_date")
        new_date = correction.get("new_date") or target_date
    new_time = correction.get("new_time")

    event_updated = False
    if target and new_date:
        try:
            event_updated = calendar_db.update_event_schedule(
                target["event_id"], new_date, new_time, title=new_title
            )
        except Exception as e:
            logger.warning("calendar correction event update failed: %s", e)

    reminders_updated = _update_calendar_correction_reminders(
        group_id=group_id,
        keywords=keywords,
        target_date=target_date,
        new_date=new_date,
        new_time=new_time,
        new_action=new_title,
        source_text=text,
    )

    if not event_updated and reminders_updated == 0:
        label = " / ".join(correction.get("keywords") or []) or "這筆"
        when = target_date or correction.get("new_date") or "指定日期"
        reply = f"有收到更正，但找不到 {when} 的「{label}」行程或提醒。"
    else:
        title = new_title or (target or {}).get("title") or "相關提醒"
        when_parts = [new_date or ""]
        if new_time:
            when_parts.append(new_time)
        when = " ".join(p for p in when_parts if p).strip()
        reply = f"已更正：{when} {title}".strip()
        details: list[str] = []
        if event_updated:
            details.append("行事曆已更新")
        if reminders_updated:
            details.append(f"提醒已同步更新 {reminders_updated} 筆")
        if details:
            reply += "\n" + "、".join(details)

    burst_filter.cancel_burst(group_id)
    memory.append_turn(group_id, "user", text)
    _append_bot_turn(group_id, reply)
    _reply(event.reply_token, reply, group_id=group_id)
    logger.info(
        "calendar correction handled group=%s event_updated=%s reminders_updated=%d text=%r",
        group_id,
        event_updated,
        reminders_updated,
        text[:80],
    )
    return True


def _handle_text_message(event: MessageEvent, group_id: str) -> None:
    text = event.message.text or ""
    text_with_quote_context = _text_with_quote_context(event.message, group_id, text)
    source = getattr(event, "source", None)
    sender_user_id = getattr(source, "user_id", None) or ""
    message_id = getattr(event.message, "id", "") or ""
    clean_text = _extract_gemini_trigger(text, event.message)

    # 回饋收集：20:00 ~ 02:00 TW 窗口內，將文字訊息存入 pending_feedback.json
    if feedback_collector.in_feedback_window():
        sender = sender_user_id or "unknown"
        try:
            feedback_collector.collect_message(sender, text)
        except Exception as e:
            logger.warning("[Feedback] collect_message failed: %s", e)

    if _try_one_shot_reply(event, group_id):
        return

    # 使用者更正既有行程/提醒時，必須即時回覆並同步改 events + reminders。
    # 放在 auto-capture / reminder extraction 前，避免更正句被誤當成新提醒。
    if _try_handle_calendar_correction(event, group_id, text):
        return

    # Organic 糾正偵測（2026-05-08 加）：user 講「不對 / 你誤會」之類的
    # 自然糾正訊號 → 抓上一輪 user/bot 訊息拼成 correction 寫進 persona_notes
    # 純信號擷取，**不**接管後續路由（用戶可能糾正完還想繼續對話）
    _detect_user_correction(text, group_id, sender_user_id)

    # 自動萃 knowledge graph 三元組（純本機，fire-and-forget）
    try:
        import knowledge_graph
        knowledge_graph.auto_extract_kg_async(group_id, text)
    except (ImportError, Exception) as e:
        logger.debug("knowledge_graph extract skipped: %s", e)

    # 自動偵測重要訊息（含日期+行程動詞）→ 抽 calendar event 寫進 DB
    # (2026-05-21 user directive: 每條留言自動判定重要性)
    # Cheap pre-filter (regex) → 通過才 spin off thread 跑 Gemini extractor
    # 重跑由 UNIQUE INDEX (group_id, title, event_date) 自動 dedup
    _auto_capture_text_if_important(
        group_id, text, sender_user_id, message_id
    )

    # 自動抽飲食 / 採購訊號（純規則 fire-and-forget，存 food_db；2026-05-31）
    # 逐則抽、不需 user_id（v1 家庭層級，GP2 A）、不需 pre-filter（無 Gemini quota 顧慮）
    try:
        import food_signals
        food_signals.extract_and_store_async(group_id, message_id, text)
    except (ImportError, Exception) as e:
        logger.debug("food_signals extract skipped: %s", e)

    # 自動分類：預設只走本機規則，Gemini fallback 必須明確開 env。
    # 這條會在每則文字訊息觸發；若 rule miss 就打 Gemini，會快速吃掉每日 RPD。
    try:
        import message_classifier

        rule_cat = message_classifier.classify_rule(text)
        if rule_cat is not None:
            message_classifier.update_category(group_id, message_id, rule_cat)
        else:
            classifier_fallback = os.environ.get(
                "GEMINI_CLASSIFIER_FALLBACK_ENABLED", ""
            ).lower() in {"1", "true", "yes", "on"}
            if classifier_fallback and _gemini_side_task_allowed(
                "message_classifier"
            ):
                message_classifier.classify_async(group_id, message_id, text)
            else:
                message_classifier.update_category(
                    group_id,
                    message_id,
                    message_classifier.DEFAULT_CATEGORY,
                )
    except (ImportError, Exception) as e:
        logger.debug("message_classifier skipped: %s", e)

    # 自動偵測 reminder（2026-05-08：含日期 + 時間/行動 hint 才抽，缺時間先當 00:00）
    # 失敗 silent，不阻塞主流程
    _maybe_extract_reminder(
        text, group_id, sender_user_id, message_id
    )

    # 1. 指令處理（指令不需要 @mention 也能用，方便管理）
    cmd_reply = _handle_command(
        group_id,
        text,
        sender_user_id,
        message_id,
    )
    if cmd_reply is not None:
        # 指令是 explicit 操作 → 取消任何待處理的 burst
        burst_filter.cancel_burst(group_id)
        _reply(event.reply_token, cmd_reply, group_id=group_id)
        return

    # 2. Explicit poll request — only when the user calls the bot.
    poll_reply = (
        _handle_explicit_poll_text(event, group_id, clean_text)
        if clean_text is not None
        else None
    )
    if poll_reply is not None:
        burst_filter.cancel_burst(group_id)
        _reply(event.reply_token, poll_reply, group_id=group_id)
        return

    # 3. 待辦 / 提醒查詢 — deterministic path，不需要 @mention 也能即時查 DB
    if _is_todo_query(text):
        burst_filter.cancel_burst(group_id)
        _handle_todo_query(event, group_id, text)
        return

    # 4. 晚餐推薦觸發
    if _is_dinner_question(text):
        burst_filter.cancel_burst(group_id)
        _handle_dinner_recommendation(event, group_id)
        return

    # 5. 未點名但明顯是可查資料問句 → 即時查網路資料後回覆
    quoted_has_url = text_with_quote_context != text and bool(
        _extract_prefetch_urls(text_with_quote_context)
    )
    quoted_web_followup = quoted_has_url and bool(
        re.search(r"這個|這篇|這則|真假|真的假的|可以嗎|能信嗎|怎麼看|如何|值得|推薦", text)
    )
    if clean_text is None and (
        _is_web_research_question(text) or quoted_web_followup
    ):
        research_text = text_with_quote_context if quoted_web_followup else text
        if _handle_web_research_question(event, group_id, research_text):
            burst_filter.cancel_burst(group_id)
            return

    quoted_id = getattr(event.message, "quoted_message_id", None)
    quoted_media_followup = bool(
        quoted_id
        and clean_text is None
        and re.search(
            r"這個|這張|這則|這影片|這段|是什麼|怎麼看|幫我看|分析|判斷|真假|可以嗎|哪裡",
            text,
        )
    )
    if quoted_media_followup:
        raw = memory.get_raw_message(group_id, quoted_id)
        if raw is not None and raw[1] in _MEDIA_PLACEHOLDERS:
            burst_filter.cancel_burst(group_id)
            _handle_media_via_quote(event, group_id, text, quoted_id, raw[1])
            return

    # 6. Explicit 觸發（@mention / /ai / /問 ...）→ 立刻處理，並取消 pending burst
    if clean_text is not None:
        burst_filter.cancel_burst(group_id)
        _handle_explicit_text(event, group_id, clean_text)
        return

    # 7. 其他文字訊息 → burst_filter debounce（等對方說完再回）
    # 7a. fast-path：如果有 due reminder 沒推過，搶在 burst_filter 累積前用 reply_token
    # 推 reminder（LINE push quota 爆時的補救路徑 — reply API 不耗月配額）
    if _try_piggyback_reminders_fast_path(event.reply_token, group_id):
        return

    burst_filter.add_to_burst(
        group_id, message_id, text_with_quote_context, sender_user_id, event.reply_token
    )


_IMAGE_GEN_PATTERNS = [
    re.compile(r"^[\s]*[畫繪]一?張?[\s]*[:：]?[\s]*(.+)", re.IGNORECASE),
    re.compile(r"^[\s]*生成圖片?[\s]*[:：]?[\s]*(.+)", re.IGNORECASE),
    re.compile(r"^[\s]*做一?張[圖照]?[\s]*[:：]?[\s]*(.+)", re.IGNORECASE),
    re.compile(r"^[\s]*幫我[畫繪][\s]*[:：]?[\s]*(.+)", re.IGNORECASE),
    re.compile(r"^[\s]*draw\s+(?:me\s+)?(.+)", re.IGNORECASE),
    re.compile(r"^[\s]*imagine\s+(.+)", re.IGNORECASE),
]


def _detect_image_gen_request(text: str) -> str | None:
    """偵測「畫一張 X / 生成圖 Y / draw Z」→ 回主題；無命中回 None。"""
    s = (text or "").strip()
    if not s:
        return None
    for pat in _IMAGE_GEN_PATTERNS:
        m = pat.match(s)
        if m:
            subject = m.group(1).strip()
            if subject and len(subject) >= 2:
                return subject
    return None


def _handle_image_gen(event: MessageEvent, group_id: str, subject: str) -> None:
    """圖片生成 — 本機 mlx SD/FLUX。LINE 需要 public URL，目前先存本機。"""
    try:
        import image_gen_local
    except ImportError:
        logger.info("image_gen_local 未安裝，silent skip")
        return
    try:
        png = image_gen_local.generate(subject, style="photo")
    except Exception as e:
        logger.warning("image_gen_local.generate failed: %s", e)
        _reply(event.reply_token, "咪寶生成圖失敗，等下再試試喔", group_id=group_id)
        return
    if not png:
        _reply(event.reply_token, "咪寶今天畫不出來（model 沒載成功）", group_id=group_id)
        return
    import uuid
    out_dir = _GENIMG_DIR
    os.makedirs(out_dir, exist_ok=True)
    fname = f"{uuid.uuid4().hex}.png"
    out_path = os.path.join(out_dir, fname)
    try:
        with open(out_path, "wb") as f:
            f.write(png)
    except Exception as e:
        logger.warning("save genimg failed: %s", e)
        return
    # 取 public URL：env IMAGE_GEN_PUBLIC_URL 優先，否則讀 /tmp/cloudflared_line_bot_url.txt
    public_url = os.environ.get("IMAGE_GEN_PUBLIC_URL", "").rstrip("/")
    if not public_url:
        try:
            with open("/tmp/cloudflared_line_bot_url.txt") as f:
                public_url = f.read().strip().rstrip("/")
        except Exception:
            public_url = ""

    if public_url:
        # 真的傳 ImageMessage
        img_url = f"{public_url}/static/img/{fname}"
        try:
            with ApiClient(_get_line_config()) as api_client:
                MessagingApi(api_client).reply_message(
                    ReplyMessageRequest(
                        reply_token=event.reply_token,
                        messages=[ImageMessage(
                            original_content_url=img_url,
                            preview_image_url=img_url,
                        )],
                    )
                )
            memory.append_turn(group_id, "user", f"[圖片生成] {subject}")
            _append_bot_turn(group_id, f"[已傳圖] {img_url}")
        except Exception as e:
            logger.warning("reply ImageMessage failed: %s", e)
            _reply(event.reply_token, f"圖生成 OK 但傳 LINE 失敗：{e}", group_id=group_id)
        return

    # 沒設 public URL → 文字 fallback（已存本機）
    msg = (
        f"圖片已生成（{len(png)//1024} KB）\n"
        f"暫存：{out_path}\n"
        f"啟動 cloudflared tunnel 後可直接傳 LINE"
    )
    _reply(event.reply_token, msg, group_id=group_id)
    memory.append_turn(group_id, "user", f"[圖片生成] {subject}")
    _append_bot_turn(group_id, f"[已生成] {out_path}")


_CALENDAR_QUERY_RE = re.compile(
    r"("
    # 1. 日期 + 行程動詞（原有）
    r"(?:今天|明天|後天|昨天|前天|這週末|下週末|下週|本週|週[一二三四五六日天])"
    r".{0,8}"
    r"(?:有事|有什麼|幾點|安排|要幹嘛|要做什麼|計畫|計劃|行程)"
    r"|"
    # 2. 行程動詞 + 日期（原有反向）
    r"(?:有事|有什麼|幾點|安排|要幹嘛|要做什麼|計畫|計劃|行程)"
    r".{0,8}"
    r"(?:今天|明天|後天|昨天|前天|這週末|下週末|下週|本週|週[一二三四五六日天])"
    r"|"
    # 3. 無日期：「什麼時候 / 上次 / 之前 / 哪一天 + 名詞動作」(GP1+GP2 反饋 noun anchor)
    r"(?:什麼時候|上次|之前|哪一天|哪天)"
    r".{0,12}"
    r"(?:回(?:台北|新北|台中|台南|高雄|花蓮|宜蘭|新竹|苗栗|嘉義|屏東|台東|老家)|"
    r"做(?:胃鏡|大腸鏡|健康檢查|體檢|手術|健檢)|"
    r"看(?:醫生|牙醫|皮膚科|眼科|耳鼻喉科)|"
    r"拿(?:蛋糕|藥|包裹|貨|禮物|花)|"
    r"接(?:爸|媽|妹|弟|姊|爺爺|奶奶|小孩|小朋友)|"
    r"陪(?:.{0,5})(?:就醫|看醫生|看病))"
    r"|"
    # 4. 反向：「名詞動作 + 什麼時候/哪一天」
    r"(?:回(?:台北|新北|台中|台南|高雄|花蓮|宜蘭|新竹|苗栗|嘉義|屏東|台東|老家)|"
    r"做(?:胃鏡|大腸鏡|健康檢查|體檢|手術|健檢)|"
    r"看(?:醫生|牙醫))"
    r".{0,12}"
    r"(?:什麼時候|哪一天|哪天|上次|之前)"
    r")"
)


# 名詞 keyword whitelist — branch 2 fallback：text 含這些 keyword → search_by_keyword
_QUERY_NOUN_KEYWORDS: tuple[str, ...] = (
    # 地名
    "台北", "新北", "台中", "台南", "高雄", "花蓮", "宜蘭", "新竹", "苗栗",
    "嘉義", "屏東", "台東", "老家",
    # 醫療
    "胃鏡", "大腸鏡", "健檢", "體檢", "醫生", "牙醫", "看病",
    # 家人
    "媽媽", "爸爸", "姊姊", "妹妹", "弟弟", "爺爺", "奶奶", "全家",
    # 物件
    "蛋糕", "禮物",
    # 場所
    "喜來登",
)

# Verb+noun phrase 抽取 — 對應 title LIKE '%verb%noun%' 順序匹配（中間可有字）
# 「媽媽什麼時候回台北」→ ('回', '台北') → title LIKE '%回%台北%'
# 「哪一天拿蛋糕」→ ('拿', '蛋糕') → 也能命中「拿爸爸生日蛋糕」
_QUERY_PHRASE_RE = re.compile(
    r"(回)(台北|新北|台中|台南|高雄|花蓮|宜蘭|新竹|苗栗|嘉義|屏東|台東|老家)"
    r"|(做)(胃鏡|大腸鏡|健康檢查|體檢|手術|健檢|LDCT)"
    r"|(看)(醫生|牙醫|皮膚科|眼科|耳鼻喉科)"
    r"|(拿)(蛋糕|藥|包裹|貨|禮物|花)"
    r"|(接)(爸|媽|妹|弟|姊|爺爺|奶奶|小孩|小朋友)"
    r"|(陪)(就醫|看醫生|看病)"
    r"|(領)(藥|處方簽|包裹)"
)


def _extract_verb_noun_pairs(text: str) -> list[tuple[str, str]]:
    """從 query 抽 (verb, noun) tuples。Regex 多 group，過濾出非 None pair。"""
    pairs: list[tuple[str, str]] = []
    for m in _QUERY_PHRASE_RE.finditer(text):
        groups = m.groups()
        # groups 是 14 個 (7 個 verb + 7 個 noun)，只有命中的 group 非 None
        for i in range(0, len(groups), 2):
            if groups[i] and groups[i + 1]:
                pairs.append((groups[i], groups[i + 1]))
                break
    return pairs

# 未來指向關鍵字 — 「什麼時候/哪一天/何時」預設只看未來
_FUTURE_LEANING_RE = re.compile(r"(?:什麼時候|哪一天|哪天|何時|什麼日子)")
_PAST_LEANING_RE = re.compile(r"(?:上次|之前|上回|前一次|何時.{0,3}過)")


def _is_calendar_query(text: str) -> bool:
    """偵測「明天有什麼 / 爸爸明天幾點要拿蛋糕 / 後天有事嗎」等行事曆查詢。

    放在 explicit handler 開頭，命中即走 deterministic calendar_db 查詢，
    完全跳過 Gemini quota（GP2 反饋：query 不該綁 lite_reply Stage 1 handler）。

    兩階段偵測（避免 _CALENDAR_QUERY_RE verb list 沒涵蓋全部行程動詞時 miss）：
    1. _CALENDAR_QUERY_RE 直接命中
    2. 含問句詞（什麼時候/哪一天/哪天/何時/上次/之前）+ verb_noun phrase 命中
    """
    if not text:
        return False
    if _CALENDAR_QUERY_RE.search(text):
        return True
    # 二段 fallback：問句詞 + 行程動作 phrase
    if re.search(r"什麼時候|哪一天|哪天|何時|上次|之前", text) and _QUERY_PHRASE_RE.search(text):
        return True
    return False


def _resolve_relative_date(text: str):
    """偵測 text 中的相對日期關鍵字 → TW timezone target date。回 None = 沒命中。

    支援未來：今天/明天/後天/週X、過去：昨天/前天/上週X/N天前。
    """
    from datetime import datetime as _dt, timedelta as _td
    from zoneinfo import ZoneInfo as _ZI
    today_tw = _dt.now(_ZI("Asia/Taipei")).date()
    # 過去先（更具體的字眼優先）
    if "前天" in text:
        return today_tw - _td(days=2)
    if "昨天" in text:
        return today_tw - _td(days=1)
    if "今天" in text:
        return today_tw
    if "明天" in text:
        return today_tw + _td(days=1)
    if "後天" in text:
        return today_tw + _td(days=2)
    # 「N 天前」(數字)
    m = re.search(r"(\d+)\s*天前", text)
    if m:
        n = min(int(m.group(1)), 365)
        return today_tw - _td(days=n)
    wmap = {"一": 0, "二": 1, "三": 2, "四": 3, "五": 4, "六": 5, "日": 6, "天": 6}
    # 「上週X」: 找上一個該星期幾
    m = re.search(r"上週([一二三四五六日天])", text)
    if m:
        target_wd = wmap[m.group(1)]
        delta = (today_tw.weekday() - target_wd) % 7
        if delta == 0:
            delta = 7  # 講「上週X」通常指上一個（即使今天是 X）
        return today_tw - _td(days=delta)
    # 「週X」: 找下一個該星期幾
    m = re.search(r"週([一二三四五六日天])", text)
    if m:
        target_wd = wmap[m.group(1)]
        delta = (target_wd - today_tw.weekday()) % 7
        if delta == 0:
            delta = 7  # 講「週X」通常指下一個
        return today_tw + _td(days=delta)
    return None


# event_type → emoji prefix (plan C: 視覺區分三類)
_TYPE_EMOJI: dict[str, str] = {
    "family_gathering": "🍽️",
    "personal_trip": "🚆",
    "medical": "🏥",
}


def _format_calendar_event(ev: dict) -> str:
    """events row → reply text。Plan C：用 event_type emoji 區分三類。"""
    import json as _json
    title = ev.get("title", "")
    date_s = ev.get("event_date", "")
    time_s = ev.get("event_time") or ""
    location = ev.get("location") or ""
    et = ev.get("event_type") or "family_gathering"
    emoji = _TYPE_EMOJI.get(et, "🗓️")
    parts_raw = ev.get("participants") or "[]"
    try:
        parts = _json.loads(parts_raw) if isinstance(parts_raw, str) else parts_raw
    except Exception:
        parts = []
    lines = [f"{emoji} {date_s}" + (f" {time_s}" if time_s else "") + f" {title}"]
    if location:
        lines.append(f"📍 {location}")
    if parts:
        lines.append(f"👥 {' / '.join(parts)}")
    return "\n".join(lines)


_TODO_QUERY_RE = re.compile(
    r"(?:待辦事項|待辦|提醒事項|提醒清單|提醒列表|會提醒的時間|提醒.*時間|"
    r"有哪些.{0,8}(?:待辦|提醒|要做|事項)|"
    r"目前.{0,8}(?:待辦|提醒|要做)|"
    r"還有.{0,8}(?:待辦|提醒|要做|事項))"
)
_TODO_CREATE_RE = re.compile(
    r"(?:提醒我|幫我提醒|新增|加一個|加入|設定提醒|記得提醒|麻煩提醒)"
)
_TODO_QUERY_INTENT_RE = re.compile(
    r"(?:哪些|目前|還有|清單|列表|查|看看|告訴我|會提醒的時間)"
)
_TODO_DETAIL_QUERY_RE = re.compile(r"(?:細節|詳細|完整|網址|連結|驗證碼|票券|預約編號)")
_REMINDER_TIMING_MARKER_RE = re.compile(r"(?:什麼時候|何時|哪一天|哪天|幾點)")
_REMINDER_TIMING_KEYWORD_ALIASES: dict[str, tuple[str, ...]] = {
    "壁球": ("壁球", "squash"),
}
_REMINDER_KNOWN_MENTION_ALIASES = {"爸爸", "媽媽", "黃聖雅", "黃聖穎"}
_REMINDER_DETAIL_PREFIXES = ("地點", "預約編號", "接送網址", "票券驗證碼", "驗證碼")


def _reminder_timing_query_keywords(text: str) -> list[str]:
    s = (text or "").strip()
    if not s or not _REMINDER_TIMING_MARKER_RE.search(s):
        return []
    lower = s.lower()
    keywords: list[str] = []
    for label, aliases in _REMINDER_TIMING_KEYWORD_ALIASES.items():
        if any(alias.lower() in lower for alias in aliases):
            keywords.append(label)
    return keywords


def _is_reminder_timing_query(text: str) -> bool:
    return bool(_reminder_timing_query_keywords(text))


def _is_todo_query(text: str) -> bool:
    """Detect questions asking the bot to read stored todos/reminders."""
    s = (text or "").strip()
    if not s:
        return False
    if _TODO_CREATE_RE.search(s) and not _TODO_QUERY_INTENT_RE.search(s):
        return False
    if _is_reminder_timing_query(s):
        return True
    return bool(_TODO_QUERY_RE.search(s))


def _fmt_todo_date(value: str | None) -> str:
    return value or "未設定日期"


def _fmt_remind_at(ts: int | float | None) -> str:
    if not ts:
        return "未設定時間"
    try:
        return datetime.fromtimestamp(int(ts), tz=ZoneInfo("Asia/Taipei")).strftime(
            "%Y-%m-%d %H:%M"
        )
    except Exception:
        return "未設定時間"


def _fmt_reminder_report_time(ts: int | float | None) -> str:
    if not ts:
        return "時間待補"
    try:
        dt = datetime.fromtimestamp(int(ts), tz=ZoneInfo("Asia/Taipei"))
    except Exception:
        return "時間待補"
    weekday = "一二三四五六日"[dt.weekday()]
    date_part = f"{dt.month}/{dt.day}（{weekday}）"
    if dt.hour == 0 and dt.minute == 0:
        return f"{date_part}時間待補"
    return f"{date_part}{dt.strftime('%H:%M')}"


def _event_line(ev: dict) -> str:
    title = ev.get("title") or ""
    date_s = ev.get("event_date") or ""
    time_s = ev.get("event_time") or ""
    location = ev.get("location") or ""
    parts_raw = ev.get("participants") or "[]"
    try:
        parts = _json.loads(parts_raw) if isinstance(parts_raw, str) else parts_raw
    except Exception:
        parts = []
    when = f"{date_s}{(' ' + time_s) if time_s else ''}".strip()
    tail = f" @ {location}" if location else ""
    people = f"（{'、'.join(parts)}）" if parts else ""
    return f"- {when} {title}{tail}{people}".strip()


def _wants_todo_details(text: str) -> bool:
    return bool(_TODO_DETAIL_QUERY_RE.search(text or ""))


def _text_message_with_mentions(
    text: str, *, validation_source: str = "text_message_with_mentions"
) -> tuple[str, object]:
    reply_text = _prepare_outbound_text(text, source=validation_source)[:4900]
    try:
        import line_mentions

        aliases = line_mentions.aliases_mentioned_in_text(reply_text)
        targets = []
        seen_user_ids: set[str] = set()
        if "@all" in reply_text or "全家" in reply_text:
            targets.append(
                line_mentions.MentionTarget(key="all", kind="all", label="@all")
            )
        seen: set[str] = set()
        deduped_aliases: list[str] = []
        for alias in aliases:
            clean_alias = str(alias or "").strip().lstrip("@")
            if not clean_alias or clean_alias == "全家" or clean_alias in seen:
                continue
            deduped_aliases.append(clean_alias)
            seen.add(clean_alias)
        for alias in deduped_aliases:
            user_id = line_mentions.user_id_for_alias(alias)
            if not user_id or user_id in seen_user_ids:
                continue
            targets.append(
                line_mentions.MentionTarget(
                    key=f"p{len(targets) + 1}",
                    kind="user",
                    user_id=user_id,
                    label=f"@{alias}",
                )
            )
            seen_user_ids.add(user_id)
        if targets:
            message_dict = line_mentions.text_v2_dict(reply_text, targets)
            return reply_text, line_mentions.sdk_message_from_text_v2_dict(message_dict)
    except Exception as e:
        logger.warning("build mention reply failed; fallback text message: %s", str(e)[:200])
    return reply_text, TextMessage(text=reply_text)


def _split_action_detail(action: str) -> tuple[str, list[str]]:
    action = str(action or "").strip()
    m = re.match(r"^(.*?)[（(]([^()（）]+)[）)]$", action)
    if not m:
        return action, []
    title = m.group(1).strip()
    detail = m.group(2).strip()
    return title or action, [detail] if detail else []


def _reminder_source_detail_lines(item: dict, action_title: str) -> list[str]:
    source_text = str(item.get("source_text") or "").strip()
    if not source_text:
        return []
    parts = [p.strip().strip("。") for p in re.split(r"[；;\n]+", source_text) if p.strip()]
    detail_parts: list[str] = []
    keyed_lines: list[str] = []
    for idx, part in enumerate(parts):
        clean = part.strip().strip("()（）")
        if not clean:
            continue
        if clean == "時間待補":
            continue
        if idx == 0 and (clean == action_title or clean in action_title or action_title in clean):
            continue
        if re.match(r"^(時間|參加人)[:：]", clean):
            continue
        if clean.startswith(_REMINDER_DETAIL_PREFIXES):
            keyed_lines.append(_normalize_reminder_detail_line(clean))
            continue
        detail_parts.append(clean)
    lines: list[str] = []
    if detail_parts:
        lines.append("細節：" + "；".join(detail_parts))
    lines.extend(keyed_lines)
    return lines


def _normalize_reminder_detail_line(clean: str) -> str:
    for prefix in _REMINDER_DETAIL_PREFIXES:
        if not clean.startswith(prefix):
            continue
        value = clean[len(prefix):].strip()
        if value.startswith(("：", ":")):
            value = value[1:].strip()
        return f"{prefix}：{value}" if value else prefix
    return clean


def _format_reminder_participants(mentions: list[str] | tuple[str, ...] | None) -> str:
    labels: list[str] = []
    for raw in mentions or []:
        name = str(raw or "").strip().lstrip("@")
        if not name:
            continue
        if name == "全家":
            label = "@all"
        elif name in _REMINDER_KNOWN_MENTION_ALIASES:
            label = f"@{name}"
        else:
            label = name
        if label not in labels:
            labels.append(label)
    return "、".join(labels)


def _format_reminder_report_item(item: dict, index: int) -> list[str]:
    action_title, action_details = _split_action_detail(str(item.get("action") or ""))
    lines = [
        f"{index}. {_fmt_reminder_report_time(item.get('remind_at'))}",
        f"事項：{action_title}",
    ]
    source_detail_lines = _reminder_source_detail_lines(item, action_title)
    action_details = [
        detail for detail in action_details
        if not any(detail in source_line for source_line in source_detail_lines)
    ]
    detail_lines = ["細節：" + "；".join(action_details)] if action_details else []
    detail_lines.extend(source_detail_lines)
    for detail in detail_lines:
        if detail.startswith("細節：") and any(
            existing.startswith("細節：") and detail.removeprefix("細節：") in existing
            for existing in lines
        ):
            continue
        if detail not in lines:
            lines.append(detail)
    participants = _format_reminder_participants(item.get("mention_aliases") or [])
    if participants:
        lines.append(f"參加人：{participants}")
    return lines


def _reminder_match_terms(keywords: list[str]) -> list[str]:
    terms: list[str] = []
    for keyword in keywords:
        for term in _REMINDER_TIMING_KEYWORD_ALIASES.get(keyword, (keyword,)):
            clean = str(term or "").strip().lower()
            if clean and clean not in terms:
                terms.append(clean)
    return terms


def _reminder_matches_timing_keywords(item: dict, keywords: list[str]) -> bool:
    terms = _reminder_match_terms(keywords)
    if not terms:
        return True
    aliases = item.get("mention_aliases") or []
    if isinstance(aliases, str):
        try:
            aliases = _json.loads(aliases)
        except Exception:
            aliases = [aliases]
    if not isinstance(aliases, (list, tuple, set)):
        aliases = [aliases]
    haystack = " ".join(
        [
            str(item.get("action") or ""),
            str(item.get("source_text") or ""),
            " ".join(str(alias or "") for alias in aliases),
        ]
    ).lower()
    return any(term in haystack for term in terms)


def _build_todo_status_reply(group_id: str, clean_text: str = "") -> str:
    """Read todos + reminders for immediate replies."""
    try:
        import calendar_db

        calendar_db.sync_active_events_to_reminders(group_id)
    except Exception as e:
        logger.warning("failed to sync active events for todo view: %s", e)

    import todo

    target = _resolve_relative_date(clean_text)
    target_iso = target.isoformat() if target else None
    timing_keywords = _reminder_timing_query_keywords(clean_text)

    try:
        todos = todo.list_pending(group_id, limit=10, due_date=target_iso)
    except Exception as e:
        logger.warning("todo query list_pending failed: %s", e)
        todos = []
    if timing_keywords:
        todos = []

    try:
        reminders = memory.list_pending_reminders(
            group_id, within_seconds=90 * 86400
        )
    except Exception as e:
        logger.warning("todo query list_pending_reminders failed: %s", e)
        reminders = []
    if target_iso:
        reminders = [
            r for r in reminders
            if _fmt_remind_at(r.get("remind_at")).startswith(target_iso)
        ]
    else:
        now_ts = int(datetime.now(tz=ZoneInfo("Asia/Taipei")).timestamp())
        reminders = [
            r for r in reminders
            if int(r.get("remind_at") or 0) >= now_ts
        ]
    if timing_keywords:
        reminders = [
            r for r in reminders
            if _reminder_matches_timing_keywords(r, timing_keywords)
        ]

    header = f"{target_iso} 的待辦/提醒：" if target_iso else "目前待辦/提醒："
    if timing_keywords:
        keyword_label = "、".join(timing_keywords)
        header = (
            f"{target_iso} 的{keyword_label}相關提醒事項＆細節："
            if target_iso
            else f"{keyword_label}相關提醒事項＆細節"
        )
    elif reminders and not todos:
        header = f"{target_iso} 的提醒事項＆細節：" if target_iso else "未來提醒事項＆細節"
    lines: list[str] = [header]
    has_any = False

    if todos:
        has_any = True
        lines.append("\n待辦事項：")
        for item in todos[:10]:
            owner = _alias_from_user_id(item.get("sender_user_id") or "")
            owner_part = f"（{owner}）" if owner else ""
            lines.append(
                f"- {_fmt_todo_date(item.get('due_date'))} {item.get('task', '')}{owner_part}"
            )

    if reminders:
        has_any = True
        lines.append("\n提醒事項：")
        for idx, item in enumerate(reminders[:10], start=1):
            lines.append("\n".join(_format_reminder_report_item(item, idx)))

    if not has_any:
        if timing_keywords:
            keyword_label = "、".join(timing_keywords)
            prefix = f"{target_iso} " if target_iso else "目前"
            return f"{prefix}沒有查到{keyword_label}相關 pending 待辦或提醒事項。"
        return (
            f"{target_iso} 沒有查到 pending 待辦或提醒事項。"
            if target_iso
            else "目前沒有查到 pending 待辦或提醒事項。"
        )
    return "\n".join(lines)


def _handle_todo_query(
    event: MessageEvent, group_id: str, clean_text: str
) -> None:
    """deterministic todo/reminder query path — read DB, skip LLM."""
    reply = _build_todo_status_reply(group_id, clean_text)
    reply_text, message = _text_message_with_mentions(reply)
    try:
        if not settings.bot_muted:
            with ApiClient(_get_line_config()) as api_client:
                MessagingApi(api_client).reply_message(
                    ReplyMessageRequest(
                        reply_token=event.reply_token,
                        messages=[message],
                    )
                )
        logger.info("todo query reply sent group=%s", group_id)
    except Exception as e:
        logger.warning("todo query reply failed: %s", e)

    memory.append_turn(group_id, "user", clean_text)
    _append_bot_turn(group_id, reply)


def _handle_calendar_query(
    event: MessageEvent, group_id: str, clean_text: str
) -> None:
    """deterministic 行事曆查詢 — 3 branch (plan C v3)：
    1. 解析出具體日期 → list_past + list_upcoming 取交集
    2. 無日期但含名詞 keyword (台北/胃鏡/媽媽...) → search_by_keyword
    3. 無日期無名詞 → 列未來 30 天

    Format: 用 event_reminder._format_event（含 🔔 header）統一格式，
    user 問「明天要幹嘛」時收到的訊息跟 launchd 主動推的 reminder 一致。
    """
    import calendar_db
    import event_reminder as _er
    from datetime import datetime as _dt
    from zoneinfo import ZoneInfo as _ZI
    today_tw = _dt.now(_ZI("Asia/Taipei")).date()
    target = _resolve_relative_date(clean_text)

    def _fmt(ev: dict) -> str:
        from datetime import date as _date
        try:
            ed = _date.fromisoformat(ev.get("event_date", ""))
            offset = (ed - today_tw).days
        except Exception:
            offset = 7  # fallback default
        return _er._format_event(ev, offset)

    if target:
        # branch 1: 具體日期 → past+future 都掃，命中該日的列出
        try:
            past = calendar_db.list_past(group_id, days=90)
            future = calendar_db.list_upcoming(group_id, days=90)
        except Exception as e:
            logger.warning("calendar query list_past/upcoming failed: %s", e)
            past, future = [], []
        target_iso = target.isoformat()
        hits = [e for e in (past + future) if e.get("event_date") == target_iso]
        if hits:
            reply = "\n\n".join(_fmt(e) for e in hits)
        else:
            reply = f"{target_iso} 沒有家族行程喔～"
    else:
        # branch 2: 無日期 → **預設只看未來**（user 2026-05-21 directive:
        # 「任何詢問通常都在問未來，不用回應過去」）
        # 只有明確「上次/之前/上回/前一次」才看 past
        today_iso = today_tw.isoformat()
        strict_past = bool(re.search(r"上次|之前|上回|前一次", clean_text))

        def _future(events: list) -> list:
            return [e for e in events if e.get("event_date", "") >= today_iso]

        def _past(events: list) -> list:
            return [e for e in events if e.get("event_date", "") < today_iso]

        # Step 1: phrase 精準 match (title LIKE '%verb%noun%')
        vn_pairs = _extract_verb_noun_pairs(clean_text)
        phrases = [f"{v}{n}" for v, n in vn_pairs]
        hits_phrase: list = []
        if vn_pairs:
            try:
                hits_phrase = calendar_db.search_by_title_phrase(
                    group_id, vn_pairs, limit=10
                )
            except Exception as e:
                logger.warning("calendar query phrase search failed: %s", e)
        phrase_label = "「" + "/".join(phrases) + "」" if phrases else None

        if strict_past:
            # 「上次媽媽什麼時候回台北」→ 只看 past (user 明確要過去)
            phrase_past = _past(hits_phrase)
            if phrase_past:
                reply = "\n\n".join(_fmt(e) for e in phrase_past[:3])
            else:
                reply = f"沒有過去{phrase_label or '相關'}的紀錄～"
        else:
            # default: future-only — phrase 命中未來就列；沒未來改 noun fallback 找未來
            phrase_future = _future(hits_phrase)
            if phrase_future:
                reply = "\n\n".join(_fmt(e) for e in phrase_future[:3])
            else:
                # Soft fallback: 優先 family member 的未來事件
                family_kw = ("媽媽", "爸爸", "姊姊", "妹妹", "弟弟", "爺爺", "奶奶")
                family_nouns = [kw for kw in family_kw if kw in clean_text]
                other_nouns = [
                    kw for kw in _QUERY_NOUN_KEYWORDS
                    if kw in clean_text and kw not in family_kw
                ]
                search_nouns = family_nouns if family_nouns else other_nouns
                noun_future: list = []
                if search_nouns:
                    try:
                        hits_noun = calendar_db.search_by_keyword(
                            group_id, search_nouns, limit=10
                        )
                        noun_future = _future(hits_noun)
                    except Exception as e:
                        logger.warning("calendar noun fallback failed: %s", e)
                if noun_future:
                    related_who = (
                        "/".join(family_nouns) if family_nouns else "近期"
                    )
                    prefix = (
                        f"未來沒有{phrase_label}的安排。{related_who}相關行程：\n\n"
                        if phrase_label else "找到相關行程：\n\n"
                    )
                    reply = prefix + "\n\n".join(_fmt(e) for e in noun_future[:3])
                elif phrase_label or search_nouns:
                    label = phrase_label or ("「" + "/".join(search_nouns) + "」")
                    reply = f"未來沒有{label}的安排～"
                else:
                    # 無 phrase 無 noun → 列未來
                    try:
                        events = calendar_db.list_upcoming(group_id, days=30)
                    except Exception as e:
                        logger.warning("calendar list_upcoming failed: %s", e)
                        events = []
                    if events:
                        reply = "最近的家族行程：\n\n" + "\n\n".join(
                            _fmt(e) for e in events[:5]
                        )
                    else:
                        reply = "目前沒有登記的家族行程～"
    logger.info(
        "calendar query reply built: len=%d preview=%r",
        len(reply), reply[:120],
    )
    # 直接呼 LINE reply API，跳過 _reply 內的 pending piggyback drain
    # （pending drain 會跑 local LLM / vision_llm，CPU heavy 害 reply 慢 1 分鐘）
    # 行事曆查詢應該 < 5 秒回，piggyback 留給其他 reply 路徑做
    reply_text = _prepare_outbound_text(reply, source="calendar_query")
    try:
        if not settings.bot_muted:
            with ApiClient(_get_line_config()) as api_client:
                MessagingApi(api_client).reply_message(
                    ReplyMessageRequest(
                        reply_token=event.reply_token,
                        messages=[TextMessage(text=reply_text)],
                    )
                )
        logger.info("calendar query reply sent group=%s", group_id)
    except Exception as e:
        logger.warning("calendar query reply failed: %s", e)

    memory.append_turn(group_id, "user", clean_text)
    _append_bot_turn(group_id, reply)


def _handle_explicit_text(event: MessageEvent, group_id: str, clean_text: str) -> None:
    """使用者明確叫 bot（@mention / /ai 等），立刻丟 Gemini 回覆。"""
    sender_user_id = getattr(event.source, "user_id", None) or ""

    # 若引用了媒體訊息（圖片 / 影片 / 音訊），走 multimodal 路徑
    quoted_id = getattr(event.message, "quoted_message_id", None)
    if quoted_id:
        raw = memory.get_raw_message(group_id, quoted_id)
        if raw is not None and raw[1] in _MEDIA_PLACEHOLDERS:
            _handle_media_via_quote(event, group_id, clean_text, quoted_id, raw[1])
            return

    # 圖片生成請求（畫一張 / 生成圖片 / draw me ...）→ 純本機 mlx SD
    gen_subject = _detect_image_gen_request(clean_text)
    if gen_subject:
        _handle_image_gen(event, group_id, gen_subject)
        return

    # clean_text 空且沒引用 → 用戶只打「咪寶」等觸發詞 → 問候回應
    if not clean_text and not quoted_id:
        _reply(event.reply_token, "嗯？\n怎麼了嗎\n要找我什麼啦", group_id=group_id)
        return

    # 群組民調 — 僅限 explicit bot trigger，不再由普通聊天自動開 poll / 記 vote。
    if clean_text:
        poll_reply = _handle_explicit_poll_text(event, group_id, clean_text)
        if poll_reply is not None:
            _reply(event.reply_token, poll_reply, group_id=group_id)
            return

    # 待辦 / 提醒查詢 — deterministic path，不依賴 Gemini quota
    if clean_text and _is_todo_query(clean_text):
        logger.info(
            "todo query routed: text=%r group=%s",
            clean_text[:50], group_id,
        )
        _handle_todo_query(event, group_id, clean_text)
        return

    # 行事曆查詢 — deterministic path，不依賴 Gemini quota
    # （GP2 反饋：query 不該綁 lite_reply Stage 1，layer 對齊）
    if clean_text and _is_calendar_query(clean_text):
        logger.info(
            "calendar query routed: text=%r group=%s",
            clean_text[:50], group_id,
        )
        _handle_calendar_query(event, group_id, clean_text)
        return

    context = memory.get_context(group_id)
    market_quote_reply = _get_explicit_market_quote_reply(
        clean_text,
        context=context,
    )
    if market_quote_reply:
        logger.info(
            "explicit market quote routed deterministically: text=%r group=%s",
            clean_text[:50], group_id,
        )
        memory.append_turn(group_id, "user", clean_text)
        _append_bot_turn(group_id, market_quote_reply)
        _reply(
            event.reply_token,
            market_quote_reply,
            group_id=group_id,
            allow_push_fallback=False,
        )
        return

    # 純文字 + 可能的文字引用
    quoted_block = _build_quoted_block(event.message, group_id)
    # 空 clean_text + 有引用 → 讓 Gemini 針對原文回應
    if not clean_text and quoted_block:
        user_input = (
            f"{quoted_block}\n\n(使用者只輸入觸發詞呼叫你，請針對上面引用的原文做回應)"
        )
    else:
        user_input = (
            clean_text if not quoted_block else f"{quoted_block}\n\n{clean_text}"
        )

    quote_policy_input = user_input
    # URL 預抓取：先用 Python 抓網頁內容塞進 prompt，繞過 Gemini url_context 的限制
    user_input = _prefetch_urls(user_input)

    # Gemini quota 爆時仍要進 _llm_chat；內部會先跑 deterministic lite_reply，
    # miss 後直接走 local_llm，避免家人 @咪寶時 bot 沉默。
    quote_reply_only = _is_market_quote_request(
        quote_policy_input,
        context=context,
    )
    facts = memory.top_facts(group_id, user_id=sender_user_id)
    pnotes = _get_persona_notes(group_id)
    try:
        with _thinking_indicator(group_id):
            reply_text = _llm_chat(user_input, context, facts, pnotes)
    except Exception as e:
        if _is_quota_error(e):
            _mark_quota_exhausted()
            logger.warning(
                "gemini chat (explicit) quota exhausted, retry via lite_reply"
            )
            # Retry once — 這次 _quota_exhausted()=True，_llm_chat 內部走
            # deterministic lite_reply → direct local_llm fallback。
            try:
                reply_text = _llm_chat(user_input, context, facts, pnotes)
            except Exception as e2:
                logger.warning("lite_reply retry failed: %s", e2)
                reply_text = ""
            if not reply_text:
                _maybe_capture_calendar_event(
                    group_id,
                    clean_text,
                    sender_user_id,
                    getattr(event.message, "id", "") if getattr(event, "message", None) else "",
                )
                if _pending_reply_enabled():
                    _save_pending_burst_text(group_id, clean_text or user_input)
                else:
                    logger.info(
                        "pending reply disabled; sending explicit visible degraded fallback group=%s",
                        group_id,
                    )
                    _reply(
                        event.reply_token,
                        _visible_llm_degraded_reply(),
                        group_id=group_id,
                        allow_push_fallback=not quote_reply_only,
                    )
                return
        elif _is_gemini_unavailable_error(e):
            logger.warning(
                "gemini chat (explicit) unavailable, retry via local text fallback: %s",
                e,
            )
            reply_text = _local_text_llm_fallback(user_input, context=context)
            if not reply_text:
                _reply(
                    event.reply_token,
                    "我剛剛接不上雲端，本機也沒生出內容；等一下再問我一次。",
                    group_id=group_id,
                    allow_push_fallback=not quote_reply_only,
                )
                return
        else:
            logger.exception("gemini chat (explicit) failed: %s", e)
            _reply(
                event.reply_token,
                _friendly_gemini_error(e),
                group_id=group_id,
                allow_push_fallback=not quote_reply_only,
            )
            return

    # local fallback 全敗時 _llm_chat 才會回空；pending reply 停用時改送可見降級回覆。
    if not reply_text or not reply_text.strip():
        if _quota_exhausted():
            _maybe_capture_calendar_event(
                group_id,
                clean_text,
                sender_user_id,
                getattr(event.message, "id", "") if getattr(event, "message", None) else "",
            )
            if _pending_reply_enabled():
                _save_pending_burst_text(group_id, clean_text or user_input)
            else:
                logger.info(
                    "pending reply disabled; sending explicit empty visible degraded fallback group=%s",
                    group_id,
                )
                _reply(
                    event.reply_token,
                    _visible_llm_degraded_reply(),
                    group_id=group_id,
                    allow_push_fallback=not quote_reply_only,
                )
            return
        _reply(
            event.reply_token,
            "我剛剛沒生出內容，等一下再問我一次好嗎？",
            group_id=group_id,
            allow_push_fallback=not quote_reply_only,
        )
        return

    # 即時糾正偵測：使用者如果在糾正 bot，自動記住
    _try_save_correction(group_id, clean_text)

    memory.append_turn(group_id, "user", user_input)
    _append_bot_turn(group_id, reply_text)
    _maybe_extract_facts(group_id, user_id=sender_user_id)
    _reply(
        event.reply_token,
        reply_text,
        group_id=group_id,
        allow_push_fallback=not quote_reply_only,
    )
    # 14:51 case: 家人在 explicit 路徑手打「YYYY-MM-DD HH:MM 拿蛋糕」，
    # 之前 _maybe_capture_calendar_event 只在 burst 路徑跑，explicit 完全沒抽。
    _maybe_capture_calendar_event(
        group_id,
        clean_text,
        sender_user_id,
        getattr(event.message, "id", "") if getattr(event, "message", None) else "",
    )


def _handle_burst_flush(group_id: str, combined_text: str, reply_token: str) -> None:
    """burst_filter 判定「值得主動回應」時觸發。跑在 Timer 的 thread 裡。

    quota/lite/local miss 時不排 pending reply，但仍送一則可見的降級回覆，
    避免群組誤判 bot 掛掉。
    """
    logger.info(
        "burst flush triggered group=%s text_len=%d",
        group_id,
        len(combined_text),
    )

    # 2026-05-18 加：抽家族財經觀點寫入 finance_views（fire-and-forget，不擋主流程）
    try:
        if _gemini_side_task_allowed("finance_view_extract"):
            import finance_view_extractor
            finance_view_extractor.maybe_extract_and_save_async(
                group_id, combined_text
            )
    except Exception as e:
        logger.debug("finance_view extract skipped: %s", e)

    # quota 爆時也繼續走 _llm_chat；內部會用 deterministic lite_reply，
    # miss 後接 direct local_llm fallback。

    context = memory.get_context(group_id)
    quote_reply_only = _is_market_quote_request(combined_text, context=context)

    # cache 命中：謠言快取直接回，省 LLM 呼叫
    cached = memory.check_fact_cache(group_id, combined_text)
    if cached:
        logger.info("burst flush cache hit group=%s", group_id)
        _reply(
            reply_token,
            cached,
            group_id=group_id,
            allow_push_fallback=not quote_reply_only,
        )
        return

    facts = memory.top_facts(group_id)
    pnotes = _get_persona_notes(group_id)

    # URL 預抓取：先用 Python 抓網頁內容塞進 prompt，繞過 Gemini url_context 的限制
    prefetched = _prefetch_urls(combined_text)

    user_input = (
        "(下面是群組裡最近累積的訊息，已經被過濾器判定值得主動回應。"
        "請根據系統指令中的規則，針對其中有查證價值或爭議點的部份做一次"
        "精簡的回應；若只是閒聊請用一句話帶過。)\n\n"
        f"{prefetched}"
    )
    try:
        with _thinking_indicator(group_id):
            reply_text = _llm_chat(user_input, context, facts, pnotes)
    except Exception as e:
        if _is_quota_error(e):
            _mark_quota_exhausted()
            logger.warning(
                "gemini chat (burst) quota exhausted, retry via lite_reply"
            )
            try:
                reply_text = _llm_chat(user_input, context, facts, pnotes)
            except Exception as e2:
                logger.warning("lite_reply retry (burst) failed: %s", e2)
                reply_text = ""
            if not reply_text:
                logger.warning("burst quota retry miss group=%s", group_id)
                _maybe_capture_calendar_event(group_id, combined_text, message_id="")
                if _pending_reply_enabled():
                    _save_pending_burst_text(group_id, combined_text)
                else:
                    logger.info(
                        "pending reply disabled; sending burst visible degraded fallback group=%s",
                        group_id,
                    )
                    _reply(
                        reply_token,
                        _visible_llm_degraded_reply(),
                        group_id=group_id,
                        allow_push_fallback=not quote_reply_only,
                    )
                return
        elif _is_gemini_unavailable_error(e):
            logger.warning(
                "gemini chat (burst) unavailable, retry via local text fallback: %s",
                e,
            )
            reply_text = _local_text_llm_fallback(user_input, context=context)
            if not reply_text:
                logger.warning(
                    "burst unavailable local fallback miss group=%s",
                    group_id,
                )
                return
        else:
            logger.exception("gemini chat (burst) failed: %s", e)
            _reply(
                reply_token,
                "Gemini 那邊好像塞車了，等一下再回你～",
                group_id=group_id,
                allow_push_fallback=not quote_reply_only,
            )
            return

    logger.info(
        "burst llm reply len=%d text=%s",
        len(reply_text) if reply_text else 0,
        repr(reply_text[:200]) if reply_text else "(empty)",
    )
    if not reply_text or not reply_text.strip():
        # Quota exhausted 且 pending reply 停用時，直接短回，不再入隊補答。
        # 非 quota empty reply 才保留 2026-05-25 的短訊 fallback。
        if _quota_exhausted():
            logger.info(
                "burst empty while quota exhausted group=%s",
                group_id,
            )
            _maybe_capture_calendar_event(group_id, combined_text, message_id="")
            if _pending_reply_enabled():
                _save_pending_burst_text(group_id, combined_text)
            else:
                logger.info(
                    "pending reply disabled; sending burst empty visible degraded fallback group=%s",
                    group_id,
                )
                _reply(
                    reply_token,
                    _visible_llm_degraded_reply(),
                    group_id=group_id,
                    allow_push_fallback=not quote_reply_only,
                )
            return
        logger.warning(
            "burst gemini returned empty reply — sending fallback msg group=%s",
            group_id,
        )
        try:
            _reply(
                reply_token,
                "咪寶聽到了但這個話題不太接得上~",
                group_id=group_id,
                allow_push_fallback=not quote_reply_only,
            )
        except Exception as e:
            logger.warning("burst empty-reply fallback failed: %s", e)
        return

    memory.store_fact_cache(group_id, combined_text, reply_text)
    memory.append_turn(group_id, "user", f"[burst]\n{combined_text}")
    _append_bot_turn(group_id, reply_text)
    _maybe_extract_facts(group_id)
    _maybe_capture_calendar_event(group_id, combined_text, message_id="")
    _reply(
        reply_token,
        reply_text,
        group_id=group_id,
        allow_push_fallback=not quote_reply_only,
    )


def _maybe_capture_calendar_event(
    group_id: str,
    combined_text: str,
    sender_user_id: str = "",
    message_id: str = "",
) -> None:
    """從 burst 抽出家族活動 → 寫 events / 取消 events。失敗不擋主流程。"""
    try:
        import calendar_db
        import calendar_extractor

        if _gemini_side_task_allowed("calendar_capture"):
            data = calendar_extractor.extract(combined_text)
        else:
            data = {
                "has_event": False,
                "is_cancellation": False,
                "title": None,
                "date": None,
                "time": None,
                "location": None,
                "participants": [],
                "cancel_target_keyword": None,
                "event_type": "family_gathering",
            }
        extracted = calendar_extractor.extract_many(combined_text, primary=data)
        if extracted["is_cancellation"]:
            kw = extracted.get("cancel_target_keyword")
            target = calendar_db.find_active_event(
                group_id, keyword=kw, near_date=extracted.get("date")
            )
            if target:
                if extracted.get("date") and extracted["date"] != target["event_date"]:
                    calendar_db.update_event_date(
                        target["event_id"], extracted["date"], extracted.get("time")
                    )
                    logger.info(
                        "calendar event rescheduled: %s → %s (group=%s)",
                        target["event_id"],
                        extracted["date"],
                        group_id,
                    )
                else:
                    calendar_db.cancel_event(target["event_id"])
                    logger.info(
                        "calendar event cancelled: %s (kw=%s, group=%s)",
                        target["event_id"],
                        kw,
                        group_id,
                    )
            return

        for data in extracted.get("events") or []:
            actor = None
            if data.get("event_type") == "medical":
                actor = _infer_medical_actor(combined_text, sender_user_id)
                data["title"] = _apply_medical_actor(data["title"], actor)
                data["participants"] = _with_medical_actor_participant(
                    data.get("participants"), actor
                )
            else:
                _apply_sender_first_person_event(data, sender_user_id)
            _apply_family_context_defaults(data, combined_text)
            event_id = calendar_db.insert_event(
                group_id=group_id,
                title=data["title"],
                event_date=data["date"],
                event_time=data.get("time"),
                location=data.get("location"),
                participants=data.get("participants") or [],
                event_type=data.get("event_type", "family_gathering"),
                source_msg_id=message_id,
            )
            if event_id:
                logger.info(
                    "calendar event captured: %s '%s' on %s type=%s (group=%s)",
                    event_id,
                    data["title"],
                    data["date"],
                    data.get("event_type", "family_gathering"),
                    group_id,
                )
            else:
                logger.info(
                    "calendar event dedup hit (skipped): '%s' on %s (group=%s)",
                    data["title"],
                    data["date"],
                    group_id,
                )
    except Exception as e:
        logger.warning("calendar capture failed: %s", e)


# 在 module load 時把 callback 注入 burst_filter
burst_filter.register_on_flush(_handle_burst_flush)


# ── 媒體 quote 觸發（唯一會分析圖片/影片/音訊的路徑）──────────────────────

# placeholder → (mime_type, 中文名)；dispatch 時寫入 raw_messages，explicit
# 路徑遇到對應 quote 時用這張表查 mime 再重新下載。
_MEDIA_PLACEHOLDERS: dict[str, tuple[str, str]] = {
    "[圖片]": ("image/jpeg", "圖片"),
    "[影片]": ("video/mp4", "影片"),
    "[音訊]": ("audio/m4a", "音訊"),
}

# 單次可下載的媒體上限（LINE 上傳原本就有限制，這裡做二層保護）
_MEDIA_BYTE_LIMIT = 20 * 1024 * 1024  # 20 MB

# Bounded worker pool for image/video webhook handlers (GP1 §6: thread leak guard).
# webhook arrival rate * vision LLM latency could exhaust threads → cap at 2 concurrent.
from concurrent.futures import ThreadPoolExecutor as _ThreadPoolExecutor
_MEDIA_EXECUTOR = _ThreadPoolExecutor(max_workers=2, thread_name_prefix="media-handler")


def _media_pipeline_fallback(
    event: MessageEvent,
    group_id: str,
    clean_text: str,
    quoted_message_id: str,
    mime_type: str,
    media_name: str,
) -> bool:
    """quota 爆時的 local fallback：呼叫 media_pipeline.analyze_image / analyze_video。

    任何失敗（download / import / analyze）都沉默退出，避免回給使用者一堆錯誤訊息。
    """
    # 1. download bytes — 失敗就沉默退出（連 Gemini 路徑也是這個 fail mode）
    try:
        data = _download_content(quoted_message_id)
    except Exception as e:
        logger.warning("media_pipeline fallback: download failed: %s", e)
        return _handle_quoted_media_description_fallback(
            event, group_id, clean_text, quoted_message_id, media_name
        )

    if len(data) > _MEDIA_BYTE_LIMIT:
        logger.info(
            "media_pipeline fallback: %s too large (%d bytes), skipping",
            media_name,
            len(data),
        )
        return False

    # 2. 路由到 image / video pipeline
    try:
        if mime_type.startswith("image/"):
            from media_pipeline import analyze_image

            reply = analyze_image(data, user_prompt=clean_text or "", group_id=group_id)
        elif mime_type.startswith("video/"):
            from media_pipeline import analyze_video

            reply = analyze_video(data, user_prompt=clean_text or "", group_id=group_id)
        else:
            return False
    except ImportError as e:
        logger.warning("media_pipeline not available: %s", e)
        return False
    except Exception as e:
        logger.warning("media_pipeline %s fallback failed: %s", media_name, e)
        return False

    # 3. 回應給 user — 沒結果就沉默
    if not reply or not str(reply).strip():
        logger.info("media_pipeline fallback: no reply for %s", media_name)
        return False

    try:
        memory.append_turn(group_id, "user", f"[{media_name} + 問題]\n{clean_text}")
        memory.log_raw_message_meta(
            group_id,
            quoted_message_id,
            media_type="image" if mime_type.startswith("image/") else "video",
            mime_type=mime_type,
            description=str(reply),
        )
        _append_bot_turn(group_id, reply)
    except Exception as e:
        logger.warning("media_pipeline fallback: memory append failed: %s", e)
    _reply(event.reply_token, reply, group_id=group_id)
    return True


def _handle_quoted_media_description_fallback(
    event: MessageEvent,
    group_id: str,
    clean_text: str,
    quoted_message_id: str,
    media_name: str,
) -> bool:
    """Answer a media quote from a cached prior media description when bytes are gone."""
    meta = memory.get_raw_message_meta(group_id, quoted_message_id) or {}
    description = str(meta.get("description") or "").strip()
    if not description:
        return False
    prompt_text = (
        f"(使用者引用了一則{media_name}向你提問，但原始媒體目前下載不到。"
        "以下是先前分析這則媒體時留下的內容摘要；請合併目前問題回答，"
        "並在不確定處明確說明是根據既有摘要判斷。)\n\n"
        f"--- {media_name}既有摘要 開始 ---\n{description[:2500]}\n"
        f"--- {media_name}既有摘要 結束 ---\n\n"
        f"使用者目前問題：{clean_text or '請根據這則媒體內容回應。'}"
    )
    try:
        if media_name == "圖片":
            from local_llm import chat as local_chat
            from media_pipeline import (
                _build_image_argument_prompt,
                _ensure_image_argument_structure,
            )

            image_task = _build_image_argument_prompt(
                clean_text or "請根據這則圖片既有摘要回應。"
            )
            image_prompt = (
                f"{image_task}\n\n"
                f"--- 圖片既有摘要 開始 ---\n{description[:2500]}\n"
                f"--- 圖片既有摘要 結束 ---"
            )

            reply = local_chat(
                image_prompt,
                context=memory.get_context(group_id),
                system_prompt=(
                    "你是 LINE 群組助理咪寶。圖片內容與圖片摘要必須維持本機處理，"
                    "不可要求雲端查圖。請根據既有圖片摘要與使用者問題回答。"
                    "固定使用四段：圖片內容、正方、反方、統一論點。"
                    "資訊不足就明確說不足，不要編造。"
                ),
                max_tokens=700,
            )
            reply = _ensure_image_argument_structure(
                reply,
                desc=description,
                ocr_text=clean_text or "",
            )
        else:
            reply = _llm_chat(
                prompt_text,
                memory.get_context(group_id),
                memory.top_facts(group_id),
                _get_persona_notes(group_id),
            )
    except Exception as e:
        logger.warning("quoted media description fallback chat failed: %s", e)
        if media_name == "圖片":
            try:
                from media_pipeline import _ensure_image_argument_structure

                reply = _ensure_image_argument_structure(
                    description,
                    desc=description,
                    ocr_text=clean_text or "",
                )
                if reply:
                    memory.append_turn(group_id, "user", prompt_text)
                    _append_bot_turn(group_id, reply)
                    _reply(event.reply_token, reply, group_id=group_id)
                    return True
            except Exception as fallback_error:
                logger.warning(
                    "quoted image deterministic fallback failed: %s", fallback_error
                )
        return False
    if not reply or not reply.strip():
        return False
    memory.append_turn(group_id, "user", prompt_text)
    _append_bot_turn(group_id, reply)
    _reply(event.reply_token, reply, group_id=group_id)
    return True


def _audio_asr_fallback(
    event: MessageEvent,
    group_id: str,
    clean_text: str,
    quoted_message_id: str,
    media_name: str,
) -> None:
    """quota 爆時的 audio fallback — 走本機 mlx-whisper。

    User policy：純本機（不打 Groq / OpenAI Whisper）。下載 audio bytes →
    audio_local.transcribe → 把轉出的文字當 user 訊息餵 fallback_chat → reply。
    任一步失敗沉默退出。
    """
    try:
        data = _download_content(quoted_message_id)
    except Exception as e:
        logger.warning("audio fallback download failed: %s", e)
        return
    if len(data) > _MEDIA_BYTE_LIMIT:
        logger.info("audio fallback skipped: file too large (%d bytes)", len(data))
        return
    try:
        import audio_local
        text = audio_local.transcribe(bytes(data), language="zh")
    except ImportError:
        logger.info("audio_local not available, silent skip")
        return
    except Exception as e:
        logger.warning("audio_local.transcribe failed: %s", e)
        return
    if not text:
        return
    user_input = f"{clean_text}\n\n（語音轉文字）{text}" if clean_text else text
    try:
        # I2 fix (2026-05-30): llm_router 已於 2026-05-18 移除（import 會 ModuleNotFoundError
        # → 被 except 吞掉 → 語音永遠不回覆）。改走 _llm_chat，與文字主鏈一致：
        # quota 爆時內部自動退 lite_reply。
        reply = _llm_chat(
            user_input,
            memory.get_context(group_id),
            memory.top_facts(group_id),
            _get_persona_notes(group_id),
        )
    except Exception as e:
        logger.warning("audio fallback chat failed: %s", e)
        return
    if not reply:
        return
    memory.append_turn(group_id, "user", f"[語音轉文字] {text}")
    _append_bot_turn(group_id, reply)
    _reply(event.reply_token, reply, group_id=group_id)


def _handle_media_via_quote(
    event: MessageEvent,
    group_id: str,
    clean_text: str,
    quoted_message_id: str,
    placeholder: str,
) -> None:
    """使用者 @AI 並引用了一則圖片/影片/音訊。

    圖片 → **永遠走 local media_pipeline**（OCR + vision LLM），不消耗 Gemini quota（user policy 2026-05-08）。
    影片 → Gemini multimodal；quota 爆才 fallback local keyframes + vision LLM。
    音訊 → Gemini multimodal；quota 爆時 `_audio_asr_fallback` 沉默退出
            （純本機策略下不打雲端 ASR）。
    """
    mime_type, media_name = _MEDIA_PLACEHOLDERS[placeholder]

    # 圖片：永遠 local，不問 quota
    if mime_type.startswith("image/"):
        logger.info("media quote: routing image to local media_pipeline")
        _media_pipeline_fallback(
            event, group_id, clean_text, quoted_message_id, mime_type, media_name
        )
        return

    # 影片 / 音訊：quota 爆 → 影片走 local，音訊先 Groq Whisper ASR 再餵 chat fallback
    if _quota_exhausted():
        if mime_type.startswith("video/"):
            logger.info("media quote (video): quota exhausted → local fallback")
            _media_pipeline_fallback(
                event, group_id, clean_text, quoted_message_id, mime_type, media_name
            )
        else:
            _audio_asr_fallback(
                event, group_id, clean_text, quoted_message_id, media_name
            )
        return

    try:
        data = _download_content(quoted_message_id)
    except Exception as e:
        logger.warning("download quoted media failed: %s", e)
        if _handle_quoted_media_description_fallback(
            event, group_id, clean_text, quoted_message_id, media_name
        ):
            return
        _reply(
            event.reply_token,
            f"這則{media_name}下載不到（LINE 最多保留 7 天，可能已經過期）。"
            "要分析的話請重新貼一次。",
            group_id=group_id,
        )
        return

    if len(data) > _MEDIA_BYTE_LIMIT:
        if _handle_quoted_media_description_fallback(
            event, group_id, clean_text, quoted_message_id, media_name
        ):
            return
        _reply(
            event.reply_token,
            f"這則{media_name}太大（{len(data) / 1024 / 1024:.1f} MB），"
            f"超過 {_MEDIA_BYTE_LIMIT // 1024 // 1024} MB 上限，沒辦法分析。",
            group_id=group_id,
        )
        return

    prompt_text = clean_text or f"請分析這則{media_name}的內容並回應。"
    parts = [
        types.Part.from_bytes(data=bytes(data), mime_type=mime_type),
        f"(使用者引用了一則{media_name}向你提問)\n\n{prompt_text}",
    ]

    context = memory.get_context(group_id)
    facts = memory.top_facts(group_id)
    pnotes = _get_persona_notes(group_id)
    try:
        with _thinking_indicator(group_id):
            reply_text = _llm_chat(parts, context, facts, pnotes)
    except Exception as e:
        if _is_quota_error(e):
            _mark_quota_exhausted()
            logger.warning("gemini chat (quoted-%s) quota exhausted", media_name)
        else:
            logger.exception("gemini chat (quoted-%s) failed: %s", media_name, e)
            _reply(event.reply_token, _friendly_gemini_error(e), group_id=group_id)
        return

    memory.append_turn(group_id, "user", f"[{media_name} + 問題]\n{prompt_text}")
    _append_bot_turn(group_id, reply_text)
    _maybe_extract_facts(group_id)
    _reply(event.reply_token, reply_text, group_id=group_id)


def _build_quoted_block(message: TextMessageContent, group_id: str) -> str | None:
    """如果訊息有引用原始訊息，回傳「原始訊息」block；否則回 None。"""
    quoted_id = getattr(message, "quoted_message_id", None)
    if not quoted_id:
        return None
    raw = memory.get_raw_message(group_id, quoted_id)
    if raw is not None:
        sender_user_id, original_text = raw
        sender_name = _get_member_display_name(group_id, sender_user_id)
        meta = memory.get_raw_message_meta(group_id, quoted_id) or {}
        meta_lines: list[str] = []
        media_type = str(meta.get("media_type") or "").strip()
        mime_type = str(meta.get("mime_type") or "").strip()
        file_name = str(meta.get("file_name") or "").strip()
        description = str(meta.get("description") or "").strip()
        if media_type or mime_type or file_name:
            parts = [p for p in (media_type, mime_type, file_name) if p]
            meta_lines.append("媒體資訊：" + " / ".join(parts))
        if description:
            meta_lines.append(f"已知內容摘要：{description[:1500]}")
        meta_block = ("\n" + "\n".join(meta_lines)) if meta_lines else ""
        return (
            "(使用者引用了下面這則原始訊息向你提問)\n"
            f"--- 原始訊息 開始 ---\n"
            f"[{sender_name}]: {original_text}{meta_block}\n"
            f"--- 原始訊息 結束 ---"
        )

    # 找不到精確的那則 → 撈最近對話當上下文，讓 Gemini 自己判斷被引用的是哪一則
    recent = memory.get_recent_raw_messages(group_id, limit=20)
    if not recent:
        return (
            "(使用者引用了群組裡的一則訊息,但原文不在記憶中，"
            "也沒有近期對話紀錄。請根據使用者自己寫的文字盡力回應。)"
        )
    lines = []
    for _mid, uid, text, _ts in recent:
        name = _get_member_display_name(group_id, uid)
        lines.append(f"[{name}]: {text}")
    ctx_block = "\n".join(lines)
    return (
        "(使用者引用了群組裡的一則訊息,但該則原文不在記憶中。\n"
        "以下是群組最近的對話紀錄,請從中推斷使用者引用的是哪一則,\n"
        "並據此回應。不要跟使用者說你找不到原文。)\n"
        f"--- 最近對話 開始 ---\n{ctx_block}\n--- 最近對話 結束 ---"
    )


def _text_with_quote_context(
    message: TextMessageContent, group_id: str, text: str
) -> str:
    """Combine the current text with the original quoted message, if any."""
    quoted_block = _build_quoted_block(message, group_id)
    if not quoted_block:
        return text
    return (
        f"{quoted_block}\n\n"
        "(下面是使用者目前這則回覆；回答時必須合併理解原始訊息與目前回覆，"
        "不要只看目前這一句。)\n"
        f"--- 目前回覆 開始 ---\n{text}\n--- 目前回覆 結束 ---"
    )


def _get_member_display_name(group_id: str, user_id: str | None) -> str:
    """查群組成員的顯示名稱;失敗就用 fallback。"""
    if user_id is None:
        return "某人"
    if user_id == "__bot__":
        return "我 (bot)"
    try:
        with ApiClient(_get_line_config()) as api_client:
            profile = MessagingApi(api_client).get_group_member_profile(
                group_id, user_id
            )
            return getattr(profile, "display_name", None) or "群組成員"
    except Exception as e:
        logger.debug("get_group_member_profile failed: %s", e)
        return "群組成員"


# 文字類 mime 白名單 — 這些 decode 成字串丟給 Gemini
_TEXT_LIKE_MIMES = {
    "application/json",
    "application/xml",
    "application/javascript",
    "application/x-yaml",
    "application/x-sh",
    "application/x-python",
    "application/x-python-code",
}
# Gemini 原生支援直接送 bytes 的 MIME
_GEMINI_NATIVE_MIMES = {
    "application/pdf",
    "image/jpeg",
    "image/png",
    "image/gif",
    "image/webp",
    "image/heic",
    "image/heif",
}
# 80k 字中文 ≈ 100〜120k tokens，留 buffer 給 system prompt + context
_TEXT_CHAR_LIMIT = 80_000


def _extract_office_text(data: bytes, file_name: str) -> str | None:
    """從 Word/Excel/PPT bytes 抽出純文字，失敗回 None。"""
    ext = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else ""
    try:
        if ext in ("docx",):
            from docx import Document
            import io

            doc = Document(io.BytesIO(data))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if ext in ("xlsx", "xls"):
            import openpyxl  # type: ignore[import-untyped]
            import io

            wb = openpyxl.load_workbook(
                io.BytesIO(data), read_only=True, data_only=True
            )
            lines = []
            for sheet in wb.worksheets:
                lines.append(f"[工作表：{sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    row_str = "\t".join("" if v is None else str(v) for v in row)
                    if row_str.strip():
                        lines.append(row_str)
            return "\n".join(lines)
        if ext in ("pptx",):
            from pptx import Presentation
            import io

            prs = Presentation(io.BytesIO(data))
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"[第 {i} 頁]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
            return "\n".join(lines)
    except Exception as e:
        logger.warning("office extract failed (%s): %s", file_name, e)
    return None


def _handle_file_message(event: MessageEvent, group_id: str) -> None:
    """檔案訊息 — 支援 PDF/圖片/Word/Excel/PPT/文字檔，其餘婉拒。"""
    msg = event.message
    file_name = getattr(msg, "file_name", "") or "unknown"
    mime_type = _guess_mime_type(file_name)
    is_text_like = mime_type.startswith("text/") or mime_type in _TEXT_LIKE_MIMES
    is_native = mime_type in _GEMINI_NATIVE_MIMES
    ext = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else ""
    is_office = ext in ("docx", "xlsx", "xls", "pptx")

    if not (is_text_like or is_native or is_office):
        _reply(
            event.reply_token,
            f"這個檔案格式 ({ext or mime_type}) 目前還不支援。\n"
            f"可以處理的格式：PDF、Word、Excel、PPT、圖片、txt/csv/json。",
            group_id=group_id,
        )
        return

    # 2026-05-16 改：刪 quota 短路。所有 file 路徑都應該嘗試 local fallback —
    # local LLM / vision / OCR 不受 Gemini quota 限制，能跑就跑。
    # 對齊 feedback_quota_fallback_never_skip.md (binary input 也不該無聲)
    try:
        data = _download_content(msg.id)
    except Exception as e:
        logger.exception("download file failed: %s", e)
        _reply(event.reply_token, f"下載檔案失敗:{type(e).__name__}", group_id=group_id)
        return

    context = memory.get_context(group_id)
    facts = memory.top_facts(group_id)
    pnotes = _get_persona_notes(group_id)

    # ── PDF / 圖片 ───────────────────────────────────────────────────────────
    if is_native:
        # 圖片永遠走純本機 media_pipeline，不把 image bytes 或圖片摘要送 Gemini。
        if mime_type.startswith("image/"):
            try:
                from media_pipeline import analyze_image

                reply_text = analyze_image(
                    data,
                    user_prompt=(
                        "請分析這張圖片檔案，說明圖片內容、正方、反方與統一論點。"
                        f"\n\n[檔名：{file_name}]"
                    ),
                    group_id=group_id,
                )
                logger.info(
                    "file image local analyze: %s -> %d chars",
                    file_name,
                    len(reply_text or ""),
                )
            except Exception as e:
                logger.warning("file image local analyze failed: %s", e)
                reply_text = None
            if reply_text and reply_text.strip():
                memory.log_raw_message_meta(
                    group_id,
                    msg.id,
                    media_type="file",
                    mime_type=mime_type,
                    file_name=file_name,
                    description=reply_text,
                )
                memory.append_turn(group_id, "user", f"[file image: {file_name}]")
                _append_bot_turn(group_id, reply_text)
                _maybe_extract_facts(group_id)
                _reply(event.reply_token, reply_text, group_id=group_id)
            else:
                _reply(
                    event.reply_token,
                    "這張圖片我這邊暫時分析不到，可能是本機圖片模型還沒載好。",
                    group_id=group_id,
                )
            return

        # PDF quota 爆時走純本機 fallback：pypdf 抽文字 → _llm_chat (自帶 fallback chain)
        if _quota_exhausted():
            reply_text = None
            if mime_type == "application/pdf":
                try:
                    from pypdf import PdfReader
                    import io
                    reader = PdfReader(io.BytesIO(data))
                    total_pages = len(reader.pages)
                    content = "\n".join(
                        (p.extract_text() or "") for p in reader.pages[:30]
                    )
                    if content.strip():
                        # 抽到文字 → 走 local LLM；prompt 把檔名收尾 metadata，
                        # 開頭明確指示避免 local 14B echo opener (規則 0 違規)
                        content = content[:_TEXT_CHAR_LIMIT]
                        page_note = (
                            f"（PDF 共 {total_pages} 頁，只看前 30 頁）"
                            if total_pages > 30 else ""
                        )
                        prompt_text = (
                            "請根據以下 PDF 內容做分析回應，用繁體中文，"
                            "第一句必須是具體判斷或結論，"
                            "不要以「使用者」「我看到」「咪寶」「這份檔案」等空話開頭。\n\n"
                            f"--- PDF 內容開始 ---\n{content}\n--- PDF 內容結束 ---\n\n"
                            f"[檔名：{file_name}]{page_note}"
                        )
                        reply_text = _llm_chat(prompt_text, context, facts, pnotes)
                        logger.info(
                            "file pdf local fallback: %s -> %d pages, %d chars text, reply %d chars",
                            file_name, total_pages, len(content), len(reply_text or ""),
                        )
                    else:
                        # 抽不到文字 (scanned PDF) → rasterize 第一頁走 vision (per user 2026-05-16)
                        logger.info(
                            "pdf local fallback: %s 抽不到文字 → rasterize page 0 走 analyze_image",
                            file_name,
                        )
                        try:
                            import fitz
                            doc = fitz.open(stream=data, filetype="pdf")
                            fitz_pages = len(doc)
                            if fitz_pages > 0:
                                pix = doc[0].get_pixmap(dpi=150)
                                img_bytes = pix.tobytes("png")
                                doc.close()
                                from media_pipeline import analyze_image
                                reply_text = analyze_image(
                                    img_bytes,
                                    user_prompt=(
                                        f"這是 scanned PDF「{file_name}」的第一頁。"
                                        "請分析內容，第一句必須是具體判斷，"
                                        "不要以「使用者」「我看到」「咪寶」開頭。"
                                    ),
                                    group_id=group_id,
                                )
                                if reply_text and fitz_pages > 1:
                                    reply_text = (
                                        f"{reply_text}\n\n"
                                        f"（scanned PDF 共 {fitz_pages} 頁，只分析第一頁）"
                                    )
                            else:
                                doc.close()
                        except Exception as e:
                            logger.warning("pdf rasterize fallback failed: %s", e)
                except Exception as e:
                    logger.warning("file pdf local fallback failed: %s", e)
            if reply_text and reply_text.strip():
                memory.log_raw_message_meta(
                    group_id,
                    msg.id,
                    media_type="file",
                    mime_type=mime_type,
                    file_name=file_name,
                    description=reply_text,
                )
                memory.append_turn(group_id, "user", f"[file: {file_name}]")
                _append_bot_turn(group_id, reply_text)
                _maybe_extract_facts(group_id)
                _reply(event.reply_token, reply_text, group_id=group_id)
            else:
                # local 也失敗（scanned PDF 抽不到 / vision LLM 掛了）才友善訊息
                _reply(event.reply_token, _quota_exhausted_message(), group_id=group_id)
            return

        # quota OK → Gemini Part bytes
        from google.genai import types as _gtypes

        parts = [
            _gtypes.Part.from_bytes(data=data, mime_type=mime_type),
            f"使用者傳了一個檔案：{file_name}。請分析其內容並回應。",
        ]
        try:
            with _thinking_indicator(group_id):
                reply_text = _llm_chat(parts, context, facts, pnotes)
        except Exception as e:
            if _is_quota_error(e):
                _mark_quota_exhausted()
            else:
                logger.exception("gemini chat (file-native) failed: %s", e)
            _reply(
                event.reply_token,
                _friendly_gemini_error(e, file_name),
                group_id=group_id,
            )
            return
        memory.log_raw_message_meta(
            group_id,
            msg.id,
            media_type="file",
            mime_type=mime_type,
            file_name=file_name,
            description=reply_text,
        )
        memory.append_turn(group_id, "user", f"[file: {file_name}]")
        _append_bot_turn(group_id, reply_text)
        _maybe_extract_facts(group_id)
        _reply(event.reply_token, reply_text, group_id=group_id)
        return

    # ── Office 文件 → 抽文字再送 ──────────────────────────────────────────────
    if is_office:
        content = _extract_office_text(data, file_name)
        if content is None:
            _reply(
                event.reply_token,
                f"讀取 {file_name} 失敗，檔案可能損毀或格式不符。",
                group_id=group_id,
            )
            return
    else:
        # 文字檔
        content = data.decode("utf-8", errors="replace")

    original_len = len(content)
    note = ""
    if original_len > _TEXT_CHAR_LIMIT:
        content = content[:_TEXT_CHAR_LIMIT]
        note = f"\n\n(原始檔案共 {original_len:,} 字，只看前 {_TEXT_CHAR_LIMIT:,} 字)"

    prompt_text = (
        f"(使用者丟了一個檔案：{file_name}){note}\n\n"
        f"--- 內容開始 ---\n{content}\n--- 內容結束 ---\n\n請分析這個檔案的內容並回應。"
    )

    try:
        with _thinking_indicator(group_id):
            reply_text = _llm_chat(prompt_text, context, facts, pnotes)
    except Exception as e:
        if _is_quota_error(e):
            _mark_quota_exhausted()
            logger.warning("gemini chat (file) quota exhausted")
        else:
            logger.exception("gemini chat (file) failed: %s", e)
        _reply(
            event.reply_token, _friendly_gemini_error(e, file_name), group_id=group_id
        )
        return

    memory.log_raw_message_meta(
        group_id,
        msg.id,
        media_type="file",
        mime_type=mime_type,
        file_name=file_name,
        description=reply_text,
    )
    memory.append_turn(group_id, "user", f"[file: {file_name}]")
    _append_bot_turn(group_id, reply_text)
    _maybe_extract_facts(group_id)
    _reply(event.reply_token, reply_text, group_id=group_id)


_PT_TZ = ZoneInfo("America/Los_Angeles")
_TW_TZ = ZoneInfo("Asia/Taipei")


def _next_gemini_reset_tw() -> tuple[str, str]:
    """算下一次 Gemini free-tier quota 重置的台灣時間。

    Gemini 免費層每天 00:00 PT 重置。DST 期間台灣 = 15:00，非 DST = 16:00。
    回傳 (絕對時間字串, 相對倒數字串)，例如 ("今天 15:00", "還有 6 小時 24 分鐘")。
    """
    now_pt = datetime.now(tz=_PT_TZ)
    next_midnight_pt = (now_pt + timedelta(days=1)).replace(
        hour=0, minute=0, second=0, microsecond=0
    )
    reset_tw = next_midnight_pt.astimezone(_TW_TZ)
    now_tw = datetime.now(tz=_TW_TZ)

    today_tw = now_tw.date()
    if reset_tw.date() == today_tw:
        prefix = "今天"
    elif reset_tw.date() == today_tw + timedelta(days=1):
        prefix = "明天"
    else:
        prefix = reset_tw.strftime("%m/%d")
    abs_str = f"{prefix} {reset_tw.strftime('%H:%M')}"

    delta = reset_tw - now_tw
    total_min = max(0, int(delta.total_seconds() // 60))
    hours, mins = divmod(total_min, 60)
    if hours > 0:
        rel_str = f"還有 {hours} 小時 {mins} 分鐘"
    else:
        rel_str = f"還有 {mins} 分鐘"
    return abs_str, rel_str


# ── Gemini quota cache ────────────────────────────────────────────────────────
# 第一次遇到 429 PerDay 就記住「今天都是爆的」，之後所有 handler 在打 Gemini 之前
# 先看這個 cache，直接短路回 quota 訊息，不浪費網路 round-trip。
# 下一個 00:00 PT（= 台灣 15:00 夏令 / 16:00 非夏令）自動失效。
_QUOTA_STATE_FILE = os.path.join(os.path.dirname(__file__), "quota_state.json")
_quota_exhausted_until_ts: float = 0.0
_quota_notified_for_ts: float = 0.0
_quota_last_probe_ts: float = 0.0


def _load_quota_state() -> None:
    """從磁碟還原 quota exhausted 狀態，避免重啟後重複嘗試已耗盡的 quota。"""
    global _quota_exhausted_until_ts, _quota_notified_for_ts, _quota_last_probe_ts
    try:
        with open(_QUOTA_STATE_FILE) as f:
            d = _json.load(f)
        _quota_exhausted_until_ts = float(d.get("exhausted_until_ts", 0))
        _quota_notified_for_ts = float(d.get("notified_for_ts", 0))
        _quota_last_probe_ts = float(d.get("last_probe_ts", 0))
    except FileNotFoundError:
        pass
    except Exception as e:
        logger.warning("load quota state failed: %s", e)


def _save_quota_state() -> None:
    """Atomic write 防 mid-write process kill 造成 quota state 損毀。

    I1 fix (2026-05-30): 用 tempfile.mkstemp 產唯一 tmp 名，不再用固定共享
    '<file>.tmp' — 否則 uvicorn handler 與獨立 cron process 幾乎同時寫時，
    兩者 os.replace 互搶會把寫到一半的 tmp 搬成正式檔 → JSON 損毀 → 解析失敗
    → exhausted_until 歸 0 → 誤判 quota 已恢復 → 狂打已爆 Gemini。
    """
    import tempfile
    tmp = None
    try:
        fd, tmp = tempfile.mkstemp(
            prefix=".quota_state.", suffix=".tmp",
            dir=os.path.dirname(_QUOTA_STATE_FILE) or ".",
        )
        with os.fdopen(fd, "w") as f:
            _json.dump(
                {
                    "exhausted_until_ts": _quota_exhausted_until_ts,
                    "notified_for_ts": _quota_notified_for_ts,
                    "last_probe_ts": _quota_last_probe_ts,
                },
                f,
            )
        os.replace(tmp, _QUOTA_STATE_FILE)
    except Exception as e:
        if tmp:
            try:
                os.unlink(tmp)
            except OSError:
                pass
        logger.warning("save quota state failed: %s", e)


def _mark_quota_exhausted() -> None:
    """記錄 Gemini quota 已爆,到下一個 00:00 PT 前都不要再打了；同時推播一次性提醒到群組。"""
    global _quota_exhausted_until_ts, _quota_notified_for_ts
    now_pt = datetime.now(tz=_PT_TZ)
    next_midnight_pt = (now_pt + timedelta(days=1)).replace(
        hour=0, minute=0, second=0, microsecond=0
    )
    _quota_exhausted_until_ts = next_midnight_pt.timestamp()
    gemini_client.mark_quota_exhausted_in_usage()
    _save_quota_state()
    logger.warning(
        "gemini quota marked exhausted until %s",
        next_midnight_pt.astimezone(_TW_TZ).strftime("%Y-%m-%d %H:%M TW"),
    )

    # quota 爆掉不再 push 通知群組（2026-05-06 用戶要求：不要顯示「⚠️ 今日免費額度已用完」之類訊息）
    # 只更新內部 notified_for_ts state，避免邏輯混亂；不發任何 LINE message
    if _quota_notified_for_ts != _quota_exhausted_until_ts:
        _quota_notified_for_ts = _quota_exhausted_until_ts
        _save_quota_state()


def _quota_exhausted() -> bool:
    """True = 本機記錄到 429 PerDay，且重置時間還沒到。
    只靠 Google 實際回 429 判斷，不靠計數器預測。"""
    return time.time() < _quota_exhausted_until_ts


def _quota_recheck_allowed() -> bool:
    """Allow a sparse Gemini retry in case cached quota exhaustion is stale."""
    if not _quota_exhausted():
        return False
    if _GEMINI_QUOTA_RECHECK_INTERVAL_SEC <= 0:
        return False
    return time.time() - _quota_last_probe_ts >= _GEMINI_QUOTA_RECHECK_INTERVAL_SEC


def _record_quota_recheck_attempt() -> None:
    global _quota_last_probe_ts
    _quota_last_probe_ts = time.time()
    _save_quota_state()


def _clear_quota_exhausted_after_recheck() -> None:
    global _quota_exhausted_until_ts, _quota_notified_for_ts
    _quota_exhausted_until_ts = 0.0
    _quota_notified_for_ts = 0.0
    _save_quota_state()
    logger.warning("gemini quota cache cleared after successful recheck")


def _quota_exhausted_message() -> str:
    """quota 爆時統一的內部狀態訊息(含動態台灣重置時間)。"""
    abs_str, rel_str = _next_gemini_reset_tw()
    return (
        f"今日請求額度已用完。\n"
        f"可以再使用的時間:{abs_str}(台灣時間,{rel_str})\n"
        f"想馬上恢復，需開啟 pay-as-you-go。"
    )


def _visible_llm_degraded_reply() -> str:
    """User-visible fallback when cloud and local LLM paths both miss."""
    return "我有收到，但現在只能先回簡短模式；複雜問題等一下再問我一次，我再補完整。"


# ── Legacy pending reply helpers（2026-06-06 起產品路徑預設停用）───────────

_PENDING_EXPLICIT_PATH = os.path.join(
    os.path.dirname(__file__), "pending_explicit_reply.json"
)
_PENDING_MEDIA_DIR = os.path.join(os.path.dirname(__file__), "pending_media")
_PENDING_DLQ_PATH = os.path.join(os.path.dirname(__file__), "pending_dlq.jsonl")
_ONE_SHOT_REPLY_PATH = os.path.join(os.path.dirname(__file__), "one_shot_replies.json")
_PENDING_MAX_AGE_SEC = 7 * 86400  # 7 天沒被 drain → 進 DLQ，避免 PDF/stuck entry 永久卡住
_PENDING_REPLY_ENABLED = _DEFAULT_PENDING_REPLY_ENABLED
_REMINDER_REPLY_PIGGYBACK_ENABLED = True


def _load_one_shot_replies() -> dict[str, str]:
    try:
        with open(_ONE_SHOT_REPLY_PATH, encoding="utf-8") as f:
            data = _json.load(f)
    except FileNotFoundError:
        return {}
    except Exception as e:
        logger.warning("load one-shot replies failed: %s", str(e)[:200])
        return {}
    if not isinstance(data, dict):
        return {}
    return {
        str(group_id): str(text)
        for group_id, text in data.items()
        if str(group_id).strip() and str(text).strip()
    }


def _save_one_shot_replies(data: dict[str, str]) -> None:
    tmp_path = f"{_ONE_SHOT_REPLY_PATH}.tmp"
    with open(tmp_path, "w", encoding="utf-8") as f:
        _json.dump(data, f, ensure_ascii=False, indent=2)
        f.write("\n")
    os.replace(tmp_path, _ONE_SHOT_REPLY_PATH)


def queue_one_shot_reply(group_id: str, text: str) -> None:
    group_id = (group_id or "").strip()
    text = (text or "").strip()
    if not group_id or not text:
        raise ValueError("group_id and text are required")
    data = _load_one_shot_replies()
    data[group_id] = text
    _save_one_shot_replies(data)


def _try_one_shot_reply(event: MessageEvent, group_id: str) -> bool:
    data = _load_one_shot_replies()
    text = data.get(group_id)
    if not text:
        return False
    reply_text, message = _text_message_with_mentions(text)
    if settings.bot_muted:
        logger.info("[MUTED] would one-shot reply group=%s len=%d", group_id, len(reply_text))
        return False
    try:
        with ApiClient(_get_line_config()) as api_client:
            MessagingApi(api_client).reply_message(
                ReplyMessageRequest(
                    reply_token=event.reply_token,
                    messages=[message],
                )
            )
    except Exception as e:
        logger.warning("one-shot reply failed; preserved group=%s: %s", group_id, str(e)[:300])
        return False
    data.pop(group_id, None)
    _save_one_shot_replies(data)
    try:
        _append_bot_turn(group_id, reply_text)
    except Exception:
        pass
    logger.info("one-shot reply sent and cleared group=%s", group_id)
    return True


def _pending_reply_enabled() -> bool:
    return _PENDING_REPLY_ENABLED


def _reminder_reply_piggyback_enabled() -> bool:
    """Allow due reminders to ride on reply_token when someone talks in-group."""
    return _REMINDER_REPLY_PIGGYBACK_ENABLED


def _load_pending_explicit() -> dict:
    import pending_store as _ps
    return _ps.load()


def _save_pending_explicit_raw(data: dict) -> None:
    import pending_store as _ps
    try:
        _ps.save_full(data)
    except Exception as e:
        logger.warning("save pending raw failed: %s", str(e)[:200])


def _save_pending_any(event, group_id: str, user_id: str | None, msg) -> None:
    """quota 爆時把任何訊息（文字/檔案/音訊）存進佇列，每則獨立不合併。恢復時由 Gemini 判語意分組。"""
    if not _pending_reply_enabled():
        logger.info(
            "pending reply disabled; skip saving message type=%s group=%s",
            type(msg).__name__, group_id,
        )
        return
    try:
        data = _load_pending_explicit()
        if group_id not in data or not isinstance(data[group_id], list):
            data[group_id] = []

        entry = {
            "user_id": user_id,
            "message_id": msg.id,
            "quote_token": getattr(msg, "quote_token", None),
            "timestamp": time.time(),
        }

        if isinstance(msg, TextMessageContent):
            entry["type"] = "text"
            entry["text"] = msg.text or ""
            # 若引用他人訊息，帶上被引用的原文，讓 Gemini 恢復時有脈絡
            qid = getattr(msg, "quoted_message_id", None)
            if qid:
                raw = memory.get_raw_message(group_id, qid)
                if raw:
                    entry["quoted_original"] = raw[1]

        elif isinstance(msg, ImageMessageContent):
            entry["type"] = "image"
            try:
                content = _download_content(msg.id)
                if len(content) > _MEDIA_BYTE_LIMIT:
                    logger.info("image too large for pending: %d bytes", len(content))
                    entry["download_failed"] = True
                else:
                    os.makedirs(_PENDING_MEDIA_DIR, exist_ok=True)
                    path = os.path.join(_PENDING_MEDIA_DIR, f"{_uuid.uuid4().hex}.jpg")
                    with open(path, "wb") as f:
                        f.write(bytes(content))
                    entry["media_path"] = path
                    memory.log_raw_message_meta(
                        group_id,
                        msg.id,
                        media_type="image",
                        mime_type="image/jpeg",
                        media_path=path,
                    )
            except Exception as e:
                logger.warning("download image for pending failed: %s", e)
                entry["download_failed"] = True

        elif isinstance(msg, VideoMessageContent):
            entry["type"] = "video"
            try:
                content = _download_content(msg.id)
                if len(content) > _MEDIA_BYTE_LIMIT:
                    logger.info("video too large for pending: %d bytes", len(content))
                    entry["download_failed"] = True
                else:
                    os.makedirs(_PENDING_MEDIA_DIR, exist_ok=True)
                    path = os.path.join(_PENDING_MEDIA_DIR, f"{_uuid.uuid4().hex}.mp4")
                    with open(path, "wb") as f:
                        f.write(bytes(content))
                    entry["media_path"] = path
                    memory.log_raw_message_meta(
                        group_id,
                        msg.id,
                        media_type="video",
                        mime_type="video/mp4",
                        media_path=path,
                    )
            except Exception as e:
                logger.warning("download video for pending failed: %s", e)
                entry["download_failed"] = True

        elif isinstance(msg, FileMessageContent):
            entry["type"] = "file"
            entry["file_name"] = getattr(msg, "file_name", "unknown")
            try:
                content = _download_content(msg.id)
                os.makedirs(_PENDING_MEDIA_DIR, exist_ok=True)
                path = os.path.join(_PENDING_MEDIA_DIR, f"{_uuid.uuid4().hex}.bin")
                with open(path, "wb") as f:
                    f.write(bytes(content))
                entry["media_path"] = path
                memory.log_raw_message_meta(
                    group_id,
                    msg.id,
                    media_type="file",
                    file_name=entry["file_name"],
                    mime_type=_guess_mime_type(entry["file_name"]),
                    media_path=path,
                )
            except Exception as e:
                logger.warning("download file for pending failed: %s", e)
                entry["download_failed"] = True

        elif isinstance(msg, AudioMessageContent):
            entry["type"] = "audio"
            try:
                content = _download_content(msg.id)
                os.makedirs(_PENDING_MEDIA_DIR, exist_ok=True)
                path = os.path.join(_PENDING_MEDIA_DIR, f"{_uuid.uuid4().hex}.m4a")
                with open(path, "wb") as f:
                    f.write(bytes(content))
                entry["media_path"] = path
                entry["mime_type"] = "audio/m4a"
                memory.log_raw_message_meta(
                    group_id,
                    msg.id,
                    media_type="audio",
                    mime_type="audio/m4a",
                    media_path=path,
                )
            except Exception as e:
                logger.warning("download audio for pending failed: %s", e)
                entry["download_failed"] = True
        else:
            return

        # C1 fix (2026-05-30): 走 pending_store.add 單鎖 load+append+save，
        # 不再 load 全 dict → append → save_full 整段覆寫（跨 process/thread lost-update）。
        import pending_store as _ps
        _ps.add(group_id, entry)
    except Exception as e:
        logger.warning("save pending any failed: %s", str(e)[:200])


def _save_pending_burst_text(group_id: str, text: str) -> None:
    """Save a combined burst as pending when quota exhaustion prevents a useful reply."""
    if not _pending_reply_enabled():
        logger.info("pending reply disabled; skip saving burst text group=%s", group_id)
        return
    text = (text or "").strip()
    if not text:
        return
    try:
        import pending_store as _ps
        _ps.add(
            group_id,
            {
                "user_id": None,
                "message_id": f"burst-{_uuid.uuid4().hex}",
                "quote_token": None,
                "timestamp": time.time(),
                "type": "text",
                "text": text,
            },
        )
    except Exception as e:
        logger.warning("save pending burst text failed: %s", str(e)[:200])


def _clear_pending_explicit(group_id: str) -> None:
    # C1 fix (2026-05-30): 走 pending_store.clear_group 單鎖（含刪 media 檔），
    # 不再 load 全 dict → pop → save_full 整段覆寫（會蓋掉別 group 並發寫入）。
    import pending_store as _ps
    _ps.clear_group(group_id)


def _commit_pending_removal(group_id: str, msg_ids: list[str]) -> int:
    """Remove specified message_ids from group pending list, atomic save.
    回 removed 數。給 peek-then-confirm pattern 用 (§3 GP1 C5)。
    """
    if not msg_ids:
        return 0
    import pending_store as _ps
    return _ps.remove_by_message_ids(group_id, msg_ids)


def _commit_pending_entries(group_id: str, entries: list[dict]) -> int:
    """Remove only the supplied snapshot entries by message_id."""
    if not entries:
        return 0
    return _commit_pending_removal(
        group_id, [it.get("message_id") for it in entries if it.get("message_id")]
    )


def _pending_push_retry_key(group_id: str, msg_ids: list[str]) -> str:
    seed = f"line_bot:pending_reply:{group_id}:{','.join(mid for mid in msg_ids if mid)}"
    return str(_uuid.uuid5(_uuid.NAMESPACE_URL, seed))


def _dlq_entry(group_id: str, entry: dict, reason: str) -> None:
    """Append-only DLQ: 一行一筆 JSON。永遠不再 retry，供之後 forensic 查看。"""
    record = {
        "group_id": group_id,
        "reason": reason,
        "dlq_at": time.time(),
        "entry": entry,
    }
    try:
        with open(_PENDING_DLQ_PATH, "a", encoding="utf-8") as f:
            f.write(_json.dumps(record, ensure_ascii=False) + "\n")
    except Exception as e:
        logger.warning("DLQ append failed: %s", str(e)[:200])


def _drop_stale_pending(group_id: str, max_age_sec: int = _PENDING_MAX_AGE_SEC) -> int:
    """掃 pending 把 timestamp 超齡的 entry 移到 DLQ + 刪 media file，回 dropped 數。

    保護機制：avoid PDF / non-text / failed-LLM entry 永久卡住消耗 piggyback slot。
    timestamp 缺失的舊 entry 視為 now（不誤刪）。
    """
    import pending_store as _ps
    drop = _ps.pop_stale(group_id, max_age_sec)
    if not drop:
        return 0
    for it in drop:
        _dlq_entry(group_id, it, reason=f"age>{max_age_sec}s")
    logger.info("DLQ drop: group=%s dropped=%d", group_id, len(drop))
    return len(drop)


def _remove_pending_by_msg_id(group_id: str, message_id: str) -> bool:
    """Idempotent removal of a single pending entry by message_id. Goes through
    pending_store (fcntl + RLock + atomic) so safe across handler thread + cron
    subprocess. Returns whether something was removed."""
    try:
        import pending_store as _ps
        return _ps.remove_by_message_id(group_id, message_id)
    except Exception as e:
        logger.warning("remove pending by msg_id failed: %s", str(e)[:200])
        return False


def _heuristic_group_messages(items: list[dict]) -> list[dict]:
    """Fallback：每則各自一組。"""
    return [{"idxs": [i], "reply_to": i} for i in range(len(items))]


def _gemini_group_messages(items: list[dict]) -> list[dict]:
    """讓 Gemini 依內容判斷哪些訊息屬於同一話題/討論串，決定分組與回覆目標。
    回傳：[{"idxs": [int,...], "reply_to": int}, ...]。失敗退回每則各自一組。"""
    if not items:
        return []
    try:
        from google.genai import types

        client = gemini_client._client
        lines = []
        for i, it in enumerate(items):
            from datetime import datetime as _dt

            ts = it.get("timestamp", 0)
            ts_str = _dt.fromtimestamp(ts).strftime("%H:%M") if ts else "??"
            who = (it.get("user_id") or "?")[:8]
            t = it.get("type", "text")
            content = it.get("text", "")
            if t == "file":
                content = f"[檔案: {it.get('file_name', '')}]"
            elif t == "audio":
                content = "[語音留言]"
            lines.append(f"[{i}] {ts_str} ({who}) {content[:200]}")

        prompt = (
            "以下是 LINE 群組在額度耗盡期間積累的訊息（依時間順序，格式：[索引] 時間 (用戶) 內容）。\n"
            "你的任務：把這些訊息分組，每組對應「值得單獨回覆一次」的內容。\n\n"
            "分組規則：\n"
            "1. 同一人連續發的多則訊息，若講的是同一件事（只是分段打），合為一組\n"
            "2. 不同人在討論同一個話題（你問我答、辯論、補充），合為一組\n"
            "3. 話題明顯轉換（新主題、新問題、無關內容）→ 新的一組\n"
            "4. 純閒聊回應（『哈哈』『好喔』『讚』等）可以單獨一組，也可以和觸發它的那則合併\n"
            "5. 一組不宜超過 8 則，若同人連說了很多且話題明顯轉移，請適時切開\n\n"
            "reply_to：從每組中選一則最具代表性的（最能讓回覆有所依附），必須是該組的索引之一。\n\n"
            "訊息列表：\n"
            + "\n".join(lines)
            + "\n\n只回傳 JSON，不要說明（每個索引恰好出現一次）：\n"
            '{"groups":[{"idxs":[int,...], "reply_to": int}, ...]}'
        )
        # 優先用 flash（分組更準確），quota 爆時降回 flash-lite
        group_model = (
            settings.gemini_model
            if not _quota_exhausted()
            else settings.gemini_light_model
        )
        resp = client.models.generate_content(
            model=group_model,
            contents=prompt,
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                temperature=0.0,
                thinking_config=types.ThinkingConfig(thinking_budget=0),
            ),
        )
        data = _json.loads(resp.text or "")
        groups_raw = data.get("groups", [])
        seen: set[int] = set()
        clean: list[dict] = []
        for g in groups_raw:
            idxs = g.get("idxs") if isinstance(g, dict) else None
            if not isinstance(idxs, list):
                continue
            ok = [
                i
                for i in idxs
                if isinstance(i, int) and 0 <= i < len(items) and i not in seen
            ]
            if not ok:
                continue
            seen.update(ok)
            reply_to = g.get("reply_to")
            if not isinstance(reply_to, int) or reply_to not in ok:
                reply_to = max(ok, key=lambda i: len(items[i].get("text") or ""))
            clean.append({"idxs": ok, "reply_to": reply_to})
        for i in range(len(items)):
            if i not in seen:
                clean.append({"idxs": [i], "reply_to": i})
        logger.info("group gemini: %d items → %d groups", len(items), len(clean))
        return clean
    except Exception as e:
        logger.warning("gemini group failed, fallback to heuristic: %s", str(e)[:200])
        return _heuristic_group_messages(items)


def _build_group_parts(items: list[dict], group_id: str) -> list:
    """把一組 pending 訊息合成 Gemini parts（文字 + 檔案/音訊 bytes）。"""
    from google.genai import types

    parts: list[object] = []
    texts = []
    for it in items:
        t = it.get("type", "text")
        if t == "text":
            txt = it.get("text", "")
            quoted_original = it.get("quoted_original")
            if quoted_original:
                texts.append(f"(引用了：『{quoted_original[:80]}』)\n{txt}")
            else:
                texts.append(txt)
        elif t == "file":
            path = it.get("media_path")
            fname = it.get("file_name", "unknown")
            if path and os.path.exists(path):
                try:
                    with open(path, "rb") as f:
                        data = f.read()
                    mime, _ = mimetypes.guess_type(fname)
                    parts.append(
                        types.Part.from_bytes(
                            data=data, mime_type=mime or "application/octet-stream"
                        )
                    )
                    texts.append(f"(使用者傳了檔案：{fname}，請分析其內容)")
                except Exception as e:
                    logger.warning("read pending file failed: %s", e)
                    texts.append(f"(使用者傳了檔案 {fname}，但讀取失敗)")
            else:
                texts.append(f"(使用者傳了檔案 {fname}，但原始內容已遺失)")
        elif t == "audio":
            path = it.get("media_path")
            mime = it.get("mime_type", "audio/m4a")
            if path and os.path.exists(path):
                try:
                    with open(path, "rb") as f:
                        data = f.read()
                    parts.append(types.Part.from_bytes(data=data, mime_type=mime))
                    texts.append("(使用者傳了語音，請先完整轉寫再回應)")
                except Exception as e:
                    logger.warning("read pending audio failed: %s", e)
                    texts.append("(使用者傳了語音但讀取失敗)")
            else:
                texts.append("(使用者傳了語音但原始內容已遺失)")

    combined_text = "\n".join(texts).strip()
    if combined_text:
        # URL 預抓
        combined_text = _prefetch_urls(combined_text)
        parts.append(combined_text)
    return parts


_GLOBAL_GATE_CACHE_TTL_SEC = 60                                # LINE quota API call 60s 內共用結果
_global_gate_cache: tuple[float, bool] | None = None
_global_gate_cache_lock = threading.Lock()


def _global_pending_drain_ready() -> bool:
    """Drain pending 前的全域 gate：mute / Gemini quota / LINE 月額度。

    True 才可以動。三個 caller（startup / retry worker / piggyback）共用。
    Fail-closed：API check 失敗一律回 False，避免不確定狀態白燒 Gemini。
    60 秒 cache 防 cross-group webhook flood 放大 LINE API 呼叫次數。
    """
    global _global_gate_cache
    now = time.time()
    with _global_gate_cache_lock:
        if _global_gate_cache and now - _global_gate_cache[0] < _GLOBAL_GATE_CACHE_TTL_SEC:
            return _global_gate_cache[1]

    _load_quota_state()
    if settings.bot_muted:
        with _global_gate_cache_lock:
            _global_gate_cache = (now, False)
        return False
    if _quota_exhausted():
        logger.info("drain pending: Gemini exhausted, defer")
        with _global_gate_cache_lock:
            _global_gate_cache = (now, False)
        return False

    try:
        import requests as _req
        from line_token_refresh import get_line_token as _glt
        _tok = _glt()
        if not _tok:
            logger.warning("drain pending: no LINE token, fail-closed skip")
            with _global_gate_cache_lock:
                _global_gate_cache = (now, False)
            return False
        _h = {"Authorization": f"Bearer {_tok}"}
        _ru = _req.get(
            "https://api.line.me/v2/bot/message/quota/consumption",
            headers=_h, timeout=5,
        )
        _rl = _req.get(
            "https://api.line.me/v2/bot/message/quota",
            headers=_h, timeout=5,
        )
        if not (_ru.ok and _rl.ok):
            logger.warning(
                "drain pending: LINE quota API non-200 (used=%s limit=%s), fail-closed skip",
                _ru.status_code, _rl.status_code,
            )
            with _global_gate_cache_lock:
                _global_gate_cache = (now, False)
            return False
        _used = _ru.json().get("totalUsage", 0)
        _limit = _rl.json().get("value", 200)
        if _used >= _limit:
            logger.info(
                "drain pending: LINE quota %d/%d exhausted, defer",
                _used, _limit,
            )
            with _global_gate_cache_lock:
                _global_gate_cache = (now, False)
            return False
    except Exception as _e:
        logger.warning(
            "drain pending: LINE quota precheck failed, fail-closed skip: %s",
            str(_e)[:120],
        )
        with _global_gate_cache_lock:
            _global_gate_cache = (now, False)
        return False

    with _global_gate_cache_lock:
        _global_gate_cache = (now, True)
    return True


# Per-group drain lock：三 caller（startup / retry / piggyback）共用，避免同 group 兩個
# thread/process 同時 drain 重複 push。thread lock 擋同 process；fcntl 擋 uvicorn/launchd
# overlap 的跨 process race。
_drain_locks: dict[str, threading.Lock] = {}
_drain_lock_factory_lock = threading.Lock()
_DRAIN_LOCK_DIR = os.path.join(os.path.dirname(__file__), ".drain_locks")


class _DrainSlot:
    def __init__(self, thread_lock: threading.Lock, file_handle) -> None:
        self._thread_lock = thread_lock
        self._file_handle = file_handle
        self._released = False

    def release(self) -> None:
        if self._released:
            return
        self._released = True
        try:
            fcntl.flock(self._file_handle.fileno(), fcntl.LOCK_UN)
        finally:
            try:
                self._file_handle.close()
            finally:
                self._thread_lock.release()


def _try_acquire_drain_slot(group_id: str) -> _DrainSlot | None:
    """Non-blocking acquire 該 group 的 drain slot。拿不到表示已有 drain 在跑。"""
    with _drain_lock_factory_lock:
        if group_id not in _drain_locks:
            _drain_locks[group_id] = threading.Lock()
        lock = _drain_locks[group_id]
    if not lock.acquire(blocking=False):
        return None
    try:
        os.makedirs(_DRAIN_LOCK_DIR, exist_ok=True)
        digest = hashlib.sha256(group_id.encode("utf-8")).hexdigest()[:32]
        fh = open(os.path.join(_DRAIN_LOCK_DIR, f"{digest}.lock"), "w")
        try:
            fcntl.flock(fh.fileno(), fcntl.LOCK_EX | fcntl.LOCK_NB)
        except BlockingIOError:
            fh.close()
            lock.release()
            return None
        return _DrainSlot(lock, fh)
    except Exception as e:
        logger.warning("drain slot acquire failed group=%s: %s", group_id, str(e)[:120])
        lock.release()
        return None


def _drain_pending_for_group(group_id: str, source: str = "startup") -> bool:
    """處理單一 group 的 pending：分組 → 逐組 LLM 回覆 → 引用推送。

    呼叫前要先 _global_pending_drain_ready() == True。
    source ∈ {"startup", "retry_worker", "piggyback"}，僅用於 log。
    回傳 True 表示有試著 drain（不論成功失敗），False 表示因另一 caller 已在 drain 同 group
    而 skip — 用於 piggyback 判斷要不要寫 throttle ts。
    Gemini 中途重新爆走 _mark_quota_exhausted + 將未處理 items 寫回 pending（不 raise），
    caller 需自己重查 `_quota_exhausted()` 決定要不要繼續下一 group（startup wrapper 有做）。
    """
    if not _pending_reply_enabled():
        logger.info("%s pending: disabled, skip group=%s", source, group_id)
        return False
    slot = _try_acquire_drain_slot(group_id)
    if slot is None:
        logger.info("%s: skip, another drain in progress for group=%s", source, group_id)
        return False
    try:
        import pending_store as _ps
        _ps.ensure_message_ids(group_id)
        items = _ps.list_for_group(group_id)
        if isinstance(items, dict):  # 舊格式相容
            items = [items]

        # D1 TTL: drain 前先清超齡 entry 到 DLQ
        dropped = _drop_stale_pending(group_id)
        if dropped:
            items = _ps.list_for_group(group_id)
        if not items:
            return True

        # __bot__ 條目不應出現在 pending（recovery 污染防護）
        bot_items = [it for it in items if it.get("user_id") == "__bot__"]
        if bot_items:
            _commit_pending_entries(group_id, bot_items)
        items = [it for it in items if it.get("user_id") != "__bot__"]
        if not items:
            return True

        groups = _gemini_group_messages(items)
        logger.info(
            "%s: group=%s items=%d groups=%d", source, group_id, len(items), len(groups)
        )

        processed_count = 0
        for g in groups:
            idxs = [
                i for i in g.get("idxs", [])
                if isinstance(i, int) and 0 <= i < len(items)
            ]
            if not idxs:
                continue
            reply_to_idx = g.get("reply_to")
            if not isinstance(reply_to_idx, int) or reply_to_idx not in idxs:
                reply_to_idx = idxs[0]
            group_items = [items[i] for i in idxs]
            msg_ids = [it.get("message_id") for it in group_items if it.get("message_id")]
            try:
                parts = _build_group_parts(group_items, group_id)
                if not parts:
                    removed = _commit_pending_removal(group_id, msg_ids)
                    processed_count += removed
                    continue

                context = memory.get_context(group_id)
                facts = memory.top_facts(group_id)
                pnotes = _get_persona_notes(group_id)
                reply_text = _llm_chat(parts, context, facts, pnotes)

                text = _prepare_outbound_text(reply_text, source="pending_push")
                footer = _get_quota_footer()
                text = text[: 4900 - len(footer)] + footer
                if _is_system_status_outbound(text):
                    logger.info(
                        "%s pending: suppressed system-status push group=%s preview=%r",
                        source, group_id, text[:120],
                    )
                    return True

                msg_kwargs = {"text": text}
                qt = items[reply_to_idx].get("quote_token")
                if qt:
                    msg_kwargs["quote_token"] = qt

                with ApiClient(_get_line_config()) as api_client:
                    MessagingApi(api_client).push_message(
                        PushMessageRequest(
                            to=group_id,
                            messages=[TextMessage(**msg_kwargs)],
                        ),
                        x_line_retry_key=_pending_push_retry_key(group_id, msg_ids),
                    )
            except Exception as e:
                if _is_quota_error(e):
                    _mark_quota_exhausted()
                    logger.warning(
                        "%s pending: quota re-exhausted, kept unprocessed pending for group=%s",
                        source, group_id,
                    )
                else:
                    logger.warning(
                        "%s pending: push/build failed (%s), kept unprocessed pending for group=%s",
                        source, str(e)[:120], group_id,
                    )
                return True

            try:
                removed = _commit_pending_removal(group_id, msg_ids)
                processed_count += removed
                if removed < len(msg_ids):
                    logger.error(
                        "%s pending: pushed but removed %d/%d pending ids for group=%s",
                        source, removed, len(msg_ids), group_id,
                    )
            except Exception as e:
                logger.exception(
                    "%s pending: pushed but commit removal failed for group=%s: %s",
                    source, group_id, e,
                )

            try:
                memory.append_turn(
                    group_id,
                    "user",
                    "\n".join(
                        it.get("text", "")
                        for it in group_items
                        if it.get("type") == "text"
                    )[:500]
                    or "[非文字訊息]",
                )
                _append_bot_turn(group_id, reply_text)
            except Exception as e:
                logger.warning(
                    "%s pending: memory append failed after push group=%s: %s",
                    source, group_id, str(e)[:120],
                )

        logger.info("%s: group=%s removed %d processed pending", source, group_id, processed_count)
        return True
    finally:
        slot.release()


def _process_pending_on_startup() -> None:
    """uvicorn 啟動時處理所有 pending：thin wrapper，所有實作在 helper 裡。"""
    if not _pending_reply_enabled():
        logger.info("pending reply disabled; startup drain skipped")
        return
    pending = _load_pending_explicit()
    if not pending:
        return
    if not _global_pending_drain_ready():
        return
    for group_id in list(pending.keys()):
        # Gemini 中途又爆，後面的 group 直接放棄這輪
        if _quota_exhausted():
            return
        _drain_pending_for_group(group_id, source="startup")


_PENDING_RETRY_INTERVAL_SEC = 6 * 60 * 60        # 6 小時跑一次（從 30 min 降頻，避免吃光每日 quota）
_PENDING_RETRY_QUOTA_RESERVE = 0.40              # 至少留 40% quota 給新訊息（不讓 retry 把 quota 全吃光）
_PIGGYBACK_DRAIN_THROTTLE_SEC = 30 * 60          # 同 group piggyback drain 冷卻時間（成功一次後）
_last_piggyback_drain_ts: dict[str, float] = {}  # group_id → 上次 piggyback drain epoch
_piggyback_drain_lock = threading.Lock()
# ThreadPoolExecutor 限 piggyback 並行度 = 2，防 webhook 洪水製造無限 thread（GP2 critical #1）。
# 啟動時建立 — 因為 ThreadPoolExecutor 在 import 時間建會在 multiprocessing fork 行為下有問題，
# 但這 module 不被 fork，所以 module-level 建構安全；放在這裡讓 piggyback 函式都看得到。
from concurrent.futures import ThreadPoolExecutor as _PiggybackTPE  # noqa: E402
_PIGGYBACK_EXECUTOR = _PiggybackTPE(max_workers=2, thread_name_prefix="piggyback")


def _has_enough_quota_for_retry() -> bool:
    """retry 前先檢查：今日 quota 用了 ≥ 60% 就停（保留 40% 給新訊息）。"""
    info = gemini_client.get_gemini_quota_info()
    if info is None:
        return True
    used_ratio = info["used_requests"] / max(info["limit_requests"], 1)
    return used_ratio < (1.0 - _PENDING_RETRY_QUOTA_RESERVE)


def _piggyback_drain_pending(group_id: str) -> None:
    """webhook 收到該 group 訊息時呼叫；補該 group 的 pending。

    Gates 順序（任一不過就跳過，且**不寫 throttle ts**，下次 webhook 立刻可再試）：
      1. throttle：同 group 30 分鐘內已成功 drain
      2. 該 group 沒 pending（fast-path bail，省 quota gate）
      3. retry quota gate（用量 ≥ 60% 保 40% 給新訊息）
      4. 全域 gate（mute / Gemini quota / LINE quota，含 60s cache）
      5. per-group drain lock（拿不到 = retry worker / startup 正在 drain 同 group）

    Throttle ts **僅在 _drain_pending_for_group 回傳 True 後寫**，避免 gate fail 燒掉
    30 分鐘冷卻（GP1 critical #1）。
    """
    now = time.time()
    with _piggyback_drain_lock:
        last = _last_piggyback_drain_ts.get(group_id, 0.0)
        if now - last < _PIGGYBACK_DRAIN_THROTTLE_SEC:
            return

    # 2026-05-30: 順帶補 reminder pending（獨立 gate + atomic claim，與 reply drain
    # 解耦——GP2 A1）。放在 reply pending 檢查前，沒 reply pending 的 group 也會補。
    try:
        _drain_pending_reminders(group_id)
    except Exception as e:
        logger.warning("reminder drain (piggyback) failed group=%s: %s", group_id, e)

    try:
        pending = _load_pending_explicit()
        if not _pending_reply_enabled():
            return
        if group_id not in pending or not pending[group_id]:
            return
        if not _has_enough_quota_for_retry():
            logger.info(
                "piggyback drain: quota usage > 60%%, skip group=%s", group_id
            )
            return
        if not _global_pending_drain_ready():
            return
        logger.info("piggyback drain: triggered by message in group=%s", group_id)
        attempted = _drain_pending_for_group(group_id, source="piggyback")
        if attempted:
            with _piggyback_drain_lock:
                _last_piggyback_drain_ts[group_id] = time.time()
    except Exception as e:
        logger.exception("piggyback drain failed for group=%s: %s", group_id, e)


def _spawn_piggyback_drain(group_id: str) -> None:
    """ThreadPoolExecutor submit piggyback drain，不阻塞 webhook。

    Fast-path bail：若 group 仍在 30 分鐘 throttle 內就不 submit，省下 executor 排隊壓力
    （webhook 洪水時 throttle 一定先 hit）。Executor cap=2 保護 cross-group flood。
    """
    if not group_id:
        return
    with _piggyback_drain_lock:
        last = _last_piggyback_drain_ts.get(group_id, 0.0)
        if time.time() - last < _PIGGYBACK_DRAIN_THROTTLE_SEC:
            return
    try:
        _PIGGYBACK_EXECUTOR.submit(_piggyback_drain_pending, group_id)
    except RuntimeError:
        # Executor shutdown（process 在 teardown）— drop，next restart 會清 pending
        pass


def _start_pending_retry_worker() -> None:
    """背景執行緒定期重試 pending；加 quota gate 避免吃光每日額度。"""
    import threading as _threading

    def _worker():
        while True:
            _threading.Event().wait(_PENDING_RETRY_INTERVAL_SEC)
            # reminder pending backstop（2026-05-30；獨立於 reply pending，沒人留言時
            # 也能補。各 group 的 _drain_pending_reminders 自帶 quota gate + per-cycle cap）
            try:
                if not _quota_exhausted():
                    for gid in memory.list_pending_reminder_groups():
                        # 每個 group 前重檢 reserve gate：累計用量達 60% 門檻即 break，
                        # 保 40% 額度給新訊息（Phase6 GP2 A1：per-group cap=5 在多
                        # group 下單檢一次會破壞 reserve 不變式 + 餓死 reply drain）
                        if not _has_enough_quota_for_retry():
                            break
                        _drain_pending_reminders(gid)
            except Exception as e:
                logger.warning("reminder pending backstop failed: %s", e)
            if not _load_pending_explicit():
                continue
            if not _pending_reply_enabled():
                continue
            if _quota_exhausted():
                continue
            if not _has_enough_quota_for_retry():
                logger.info("pending retry worker: quota usage > 60%%, skip to preserve for new messages")
                continue
            logger.info("pending retry worker: quota available + sufficient, processing pending")
            try:
                _process_pending_on_startup()
            except Exception as e:
                logger.exception("pending retry worker failed: %s", e)

    t = _threading.Thread(target=_worker, daemon=True, name="pending-retry")
    t.start()


_local_text_prewarm_started = False


def _prewarm_local_text_llm_if_needed() -> None:
    """Preload local LLM in background during quota outage to avoid cold reply."""
    global _local_text_prewarm_started
    if _local_text_prewarm_started:
        return
    if not _quota_exhausted():
        return
    disabled = os.environ.get("LOCAL_LLM_PREWARM_DISABLED", "").lower()
    if disabled in {"1", "true", "yes", "on"}:
        return
    _local_text_prewarm_started = True

    def _run() -> None:
        start = time.time()
        try:
            import local_llm

            ensure_loaded = getattr(local_llm, "_ensure_loaded", None)
            ok = bool(ensure_loaded()) if callable(ensure_loaded) else False
            model_name = getattr(local_llm, "loaded_model_name", lambda: None)()
            logger.info(
                "local text fallback prewarm ok=%s model=%s seconds=%.1f",
                ok,
                model_name,
                time.time() - start,
            )
        except Exception as e:
            logger.warning("local text fallback prewarm failed: %s", e)

    threading.Thread(target=_run, daemon=True, name="local-llm-prewarm").start()


def _init_on_startup() -> None:
    _load_quota_state()
    _prewarm_local_text_llm_if_needed()
    _start_pending_retry_worker()


@contextmanager
def _thinking_indicator(group_id: str | None, delay: float = 3.0):
    yield


def _get_persona_notes(group_id: str) -> list[dict]:
    """取出 persona notes，供 gemini_client.chat() 注入 system prompt。"""
    return memory.list_persona_notes(group_id)


# 糾正偵測：使用者 @mention bot 時如果內容像糾正，自動存下來
_CORRECTION_KEYWORDS = (
    "不要",
    "不准",
    "別再",
    "下次",
    "記住",
    "以後",
    "不可以",
    "禁止",
    "不能",
    "改掉",
    "不用",
)


def _try_save_correction(group_id: str, user_text: str) -> None:
    """如果 user_text 看起來像是在糾正 bot 的行為，存成 persona correction。"""
    t = user_text.strip()
    if len(t) < 4 or len(t) > 100:
        return
    if any(kw in t for kw in _CORRECTION_KEYWORDS):
        memory.add_persona_note(group_id, "correction", "使用者糾正", t)
        logger.info("persona correction saved: %s", t[:60])


# ── Organic 糾正偵測（2026-05-08 加）────────────────────────────────────────
# user 主動講「不對 / 你誤會 / 我意思不是這個」時，把上一輪 user msg + bot reply
# + 這次糾正三段拼起來寫進 persona_notes，source='organic'，下一輪 prompt
# 就會把它當 negative example 放在規則 0 延伸區。
#
# 跟 _CORRECTION_KEYWORDS 不同：
# - _CORRECTION_KEYWORDS 抓「未來規則」（不要 X、別再 X、以後 X）
# - _ORGANIC_CORRECTION_KEYWORDS 抓「即時糾正上一輪回覆錯了」
_ORGANIC_CORRECTION_KEYWORDS = (
    "不對",
    "不是這樣",
    "不是這意思",
    "你誤會",
    "妳誤會",
    "我說的是",
    "我問的是",
    "我是說",
    "我意思是",
    "我意思不是",
    "重來",
    "請重答",
    "你答錯",
    "妳答錯",
    "答錯了",
    "胡說",
    "亂講",
    "不是我要的",
)


def _weak_correction_signals(
    text: str, prev_user_msg: str, prev_bot_msg: str
) -> dict:
    """除了強 keyword 以外的 5 種弱糾正信號 — 自動判斷不靠 user 配合特定詞。

    Signals:
      - opening_neg: 短訊息 (< 30 字) AND 開頭「不/沒/錯」
      - corrective_marker: 含「應該/才對/才是/正確的是」
      - sarcasm: 含「呵呵 / ㄏㄏ / ... / 。。。 / 笑死 / ｗ」
      - repeat_question: 跟前一輪 user 訊息 token jaccard > 0.5（暗示重述）
      - contrastive_negation: 「不是 X 是 Y」型句式
      - negative_sentiment: chinese_nlp.detect_sentiment 判定負面

    回 {score: int, signals: list[str]}。score >= 2 → 視為 correction。
    """
    signals: list[str] = []
    t = text.strip()
    if len(t) < 30 and t[:1] in ("不", "沒", "錯"):
        signals.append("opening_neg")
    if any(p in t for p in ("應該", "才對", "才是", "正確的是", "明明")):
        signals.append("corrective_marker")
    if any(p in t for p in ("呵呵", "ㄏㄏ", "...", "。。。", "ｗｗ", "wwww", "笑死")):
        signals.append("sarcasm")
    if "不是" in t and "是" in t.replace("不是", "", 1):
        signals.append("contrastive_negation")
    # repeat_question (跟前一輪 user 訊息 jaccard)
    if prev_user_msg:
        try:
            from chinese_nlp import tokenize
            now = set(tokenize(t))
            prev = set(tokenize(prev_user_msg))
            if now and prev and len(now & prev) / len(now | prev) > 0.5:
                signals.append("repeat_question")
        except Exception:
            pass
    # 負面情感
    try:
        from chinese_nlp import detect_sentiment
        if detect_sentiment(t) == "negative":
            signals.append("negative_sentiment")
    except Exception:
        pass
    return {"score": len(signals), "signals": signals}


def _detect_user_correction(
    text: str, group_id: str, sender_user_id: str = ""
) -> bool:
    """偵測 user 訊息是否在糾正 bot 上一輪回覆。命中 → 寫進 organic correction。

    Layer 1（強 keyword）：17 條糾正詞之一命中 → 直接視為 correction
    Layer 2（弱信號）：6 種弱信號（情感 / 重述 / 否定句式 / 嘲諷 / corrective marker / opening neg）
                     score >= 2 → 視為 correction
    寫進 persona_notes（kind='correction', source='organic'）。
    任何步驟例外都吞掉、回 False，主流程不受影響。
    """
    try:
        t = (text or "").strip()
        if len(t) < 2 or len(t) > 200:
            return False

        # 抓上一輪 user / bot 訊息
        ctx = memory.get_context(group_id)
        prev_user_msg = ""
        prev_bot_msg = ""
        # 由新→舊找最近一筆 bot，再抓它前面最近一筆 user
        last_bot_idx = None
        for i in range(len(ctx) - 1, -1, -1):
            if ctx[i][0] == "bot":
                last_bot_idx = i
                break
        if last_bot_idx is not None:
            prev_bot_msg = ctx[last_bot_idx][1] or ""
            for j in range(last_bot_idx - 1, -1, -1):
                if ctx[j][0] == "user":
                    prev_user_msg = ctx[j][1] or ""
                    break

        # 沒有上一輪 bot reply 可糾正 → 跳過（user 可能在糾正其他人類發言）
        if not prev_bot_msg:
            return False

        # Layer 1：強 keyword
        hit_keyword = any(kw in t for kw in _ORGANIC_CORRECTION_KEYWORDS)
        # Layer 2：弱信號 score >= 2
        weak = _weak_correction_signals(t, prev_user_msg, prev_bot_msg)
        is_correction = hit_keyword or weak["score"] >= 2

        if not is_correction:
            return False

        logger.info(
            "organic correction detected (keyword=%s, weak_signals=%s, text=%r)",
            hit_keyword, weak.get("signals", []), t[:60],
        )

        # 嘗試用 light Gemini call 抽一句話總結 — 失敗就直接存 raw
        summary = _summarize_correction(prev_user_msg, prev_bot_msg, t)

        note_id = memory.add_organic_correction(
            group_id=group_id,
            prev_user_msg=prev_user_msg,
            prev_bot_msg=prev_bot_msg,
            correction_msg=t,
            summary=summary,
        )
        logger.info(
            "organic correction saved (group=%s, note_id=%s, summary=%r, "
            "correction=%r)",
            group_id, note_id, (summary or "")[:50], t[:60],
        )
        return True
    except Exception as e:
        logger.warning("_detect_user_correction failed: %s", e)
        return False


def _summarize_correction(
    prev_user_msg: str, prev_bot_msg: str, correction_msg: str
) -> str:
    """用 Gemini 一句話總結「bot 具體做錯什麼」。quota 爆 / 失敗回空字串。

    刻意 light call：不過 system prompt、不過完整 chat 流程，
    避免占用主對話 quota 預算。失敗 silent，caller 直接存 raw 即可。
    """
    if _quota_exhausted():
        return ""
    try:
        prompt = (
            "下面是一段 LINE 群對話。user 在最後一句糾正 bot 的回覆。"
            "請用一句話（25 字內、繁體中文）總結 bot 到底做錯什麼，"
            "讓 bot 下次不要再犯。直接給結論，不要前綴。\n\n"
            f"user 原問：{(prev_user_msg or '')[:200]}\n"
            f"bot 答：{(prev_bot_msg or '')[:200]}\n"
            f"user 糾正：{(correction_msg or '')[:200]}"
        )
        out = gemini_client.chat(prompt, [], [])
        if not out:
            return ""
        # 截到單句、去除多餘換行
        line = out.strip().splitlines()[0].strip() if out.strip() else ""
        return line[:120]
    except Exception as e:
        logger.warning("_summarize_correction failed (silent): %s", e)
        return ""


def _is_quota_error(e: Exception) -> bool:
    """判斷是不是 Gemini 日額度爆的 429。"""
    s = str(e)
    return ("429" in s or "RESOURCE_EXHAUSTED" in s) and (
        "PerDay" in s or "free_tier_requests" in s
    )


def _is_gemini_unavailable_error(e: Exception) -> bool:
    """True for transient Gemini availability failures that local LLM can cover."""
    s = str(e)
    lower = s.lower()
    return "503" in s or "unavailable" in lower or "high demand" in lower


def _is_definite_reply_token_error(e: Exception) -> bool:
    """Only fallback push when LINE definitely rejected the reply token.

    Network timeouts/5xx are ambiguous: LINE may have accepted reply_message, so
    pushing again can duplicate the primary response.
    """
    status = getattr(e, "status", None) or getattr(e, "status_code", None)
    body = str(e).lower()
    try:
        status_int = int(status) if status is not None else None
    except (TypeError, ValueError):
        status_int = None
    if status_int is not None and status_int >= 500:
        return False
    tokenish = (
        "reply_token" in body
        or "reply token" in body
        or ("token" in body and ("expired" in body or "invalid" in body))
    )
    definite = any(
        marker in body
        for marker in (
            "expired",
            "invalid",
            "already been used",
            "used reply token",
            "reply token not found",
        )
    )
    return bool(tokenish and definite)


def _friendly_gemini_error(e: Exception, file_name: str | None = None) -> str:
    """把 google-genai SDK 的錯誤翻成使用者友善訊息。"""
    err_str = str(e)
    if "429" in err_str or "RESOURCE_EXHAUSTED" in err_str:
        # 日額度爆 (RequestsPerDay)
        if "PerDay" in err_str or "free_tier_requests" in err_str:
            return _quota_exhausted_message()
        # 分鐘級限制（transient）→ 靜默，不告知使用者
        return ""
    if "401" in err_str or "403" in err_str:
        return "Gemini API key 有問題,請檢查設定。"
    if "400" in err_str:
        return f"Gemini 說這個輸入有問題:{err_str[:200]}"
    if (
        "500" in err_str
        or "503" in err_str
        or "UNAVAILABLE" in err_str
        or "Server disconnected" in err_str
        or "RemoteProtocolError" in type(e).__name__
        or "Connection reset" in err_str
        or "ReadTimeout" in err_str
    ):
        return "Gemini 那邊暫時斷線,等一下再試。"
    return f"分析失敗:{type(e).__name__}"


def _maybe_extract_facts(group_id: str, user_id: str = "") -> None:
    """每 N 輪抽一次長期事實，user_id 有值時存為 per-user 事實。"""
    if not _gemini_side_task_allowed("fact_extract"):
        return
    if not memory.bump_and_should_extract(group_id):
        return
    try:
        new_facts = gemini_client.extract_facts(memory.get_context(group_id))
    except Exception as e:
        if _is_quota_error(e):
            _mark_quota_exhausted()
        else:
            logger.warning("auto fact extract failed: %s", e)
        return
    added = 0
    for f in new_facts:
        if memory.add_fact(group_id, f, user_id=user_id):
            added += 1
    logger.info(
        "auto-extracted facts: %d new (total=%d)",
        added,
        len(memory.list_facts(group_id)),
    )


# ── 自動偵測 reminder（2026-05-08 加，用戶要求自動記住日常事項）────────────

# 候選訊息 hint：含「日期 keyword」+「時間 keyword」才送 Gemini，省 quota
_REMINDER_DATE_HINT = re.compile(
    r"(\d+\s*月\s*\d+|\d+/\d+|\d+號|\d+日|"
    r"今天|今晚|明天|明日|明晚|後天|大後天|"
    r"(?:星期|週|周|禮拜)[一二三四五六日天]|"
    r"下\s*(週|周|星期|月)|這\s*(週|周|星期))"
)
_REMINDER_TIME_OR_ACTION_HINT = re.compile(
    r"(\d{1,2}\s*[:：點時]|\d{1,2}\s*:\s*\d{2}|"
    r"早上|上午|中午|下午|晚上|今晚|明晚|"
    r"提醒|記得|別忘|要|開會|會議|預約|回診|看診|看醫生|"
    r"訂|買|拿|取|接送|出發|到站|繳|付款|聚餐|報到)"
)

_REMINDER_DRAIN_CAP = 5  # GP2 D1b: 每次 drain 最多抽幾筆，攤平 backlog 不燒爆當天額度


def _has_calendar_event_like_content(text: str, today_tw: object | None = None) -> bool:
    """判斷文字是否像家族事件，避免同訊息同時進 events + reminders。"""
    if not text:
        return False
    try:
        import calendar_regex

        if today_tw is None:
            from datetime import date as _date

            today_tw = _date.today()
        return bool(calendar_regex.extract_many_regex_only(text, today_tw, require_time=False))
    except Exception as e:
        logger.debug("calendar event-like parse failed in reminder gate: %s", e)
        return False


def _calendar_regex_to_reminder_result(
    text: str, today_tw, user_id: str | None = None
) -> dict | None:
    """Use deterministic calendar regex as a reminder extractor fallback."""
    try:
        import calendar_regex
        events = calendar_regex.extract_many_regex_only(text, today_tw, require_time=False)
    except Exception as e:
        logger.debug("calendar regex reminder fallback skipped: %s", e)
        return None

    data = events[0] if events else None
    if not (data is not None and data.get("has_event") and data.get("date")):
        return None

    try:
        year_s, month_s, day_s = str(data["date"]).split("-", 2)
        if data.get("time"):
            hour_s, minute_s = str(data["time"]).split(":", 1)
        else:
            hour_s, minute_s = "0", "0"
        actor = _infer_medical_actor(text, user_id)
        companions = _infer_medical_companions(text)
        action = _apply_medical_actor(
            data.get("title") or text[:30], actor, companions
        )
        return {
            "action": action,
            "mention_aliases": _medical_mention_aliases(actor, companions),
            "year": int(year_s),
            "month": int(month_s),
            "day": int(day_s),
            "hour": int(hour_s),
            "minute": int(minute_s),
        }
    except (ValueError, TypeError):
        return None


def _enqueue_reminder_if_candidate(
    text: str, group_id: str, user_id: str, message_id: str | None
) -> None:
    """quota 爆時的 reminder 補救入隊：只對含日期 + 時間/行動 hint 的訊息入隊，等額度恢復
    由 _drain_pending_reminders 重抽（forward-only，絕不重掃 raw_messages）。失敗
    silent、自包 try/except，絕不可炸掉 caller（GP2 S4-sec：site 1 緊鄰
    _save_pending_any，炸了會連 reply pending 一起丟）。"""
    try:
        if not text or len(text) > 500:
            return
        if not _REMINDER_DATE_HINT.search(text):
            return
        if not _REMINDER_TIME_OR_ACTION_HINT.search(text):
            return
        memory.enqueue_pending_reminder(group_id, user_id, text, message_id)
    except Exception as e:
        logger.warning("_enqueue_reminder_if_candidate failed: %s", e)


def _drain_pending_reminders(group_id: str, limit: int = _REMINDER_DRAIN_CAP) -> None:
    """額度恢復後補抽該 group 的 pending reminder。

    GP1 R1: today_iso 用每筆 created_at 還原，相對日期（明天/今晚）才不會對到 drain
    當天。GP2 D1: _has_enough_quota_for_retry gate + per-cycle cap，避免燒爆當天 20 次
    + 跨餓 reply drain。GP2 A2/S1: 用 DB atomic claim（不重用 _try_acquire_drain_slot，
    那 key 無 namespace 會 starve reply drain）。
    """
    if _quota_exhausted() or not _has_enough_quota_for_retry():
        return
    from datetime import datetime as _dt
    try:
        memory.drop_stale_pending_reminders(_PENDING_MAX_AGE_SEC, group_id)
        rows = memory.list_pending_reminder_retries(group_id, limit=limit)
    except Exception as e:
        logger.warning("drain reminders: list failed group=%s: %s", group_id, e)
        return
    for row in rows:
        pid = row["pending_id"]
        if not memory.claim_pending_reminder(pid):  # atomic claim：搶不到表示別的 drain 在處理
            continue
        terminal_written = False
        try:
            # R1: 用訊息「當時」的 created_at 還原 today_iso，解相對日期
            today_iso = _dt.fromtimestamp(row["created_at"]).strftime("%Y-%m-%d %A")
            msg_dt = _dt.fromtimestamp(row["created_at"])
            result = gemini_client.extract_reminder(row["text"], today_iso=today_iso)
            if result is None:
                result = _calendar_regex_to_reminder_result(
                    row["text"], msg_dt.date(), row["user_id"]
                )
                if result is None:
                    memory.mark_pending_reminder(pid, "dropped")  # Gemini 判定非提醒
                    terminal_written = True
                    continue
            else:
                actor = _infer_medical_actor(row["text"], row["user_id"])
                if actor and "action" in result:
                    companions = _infer_medical_companions(row["text"])
                    result["action"] = _apply_medical_actor(
                        str(result["action"]), actor, companions
                    )
                    result["mention_aliases"] = _medical_mention_aliases(
                        actor, companions
                    )
            try:
                remind_dt = _dt(
                    int(result["year"]), int(result["month"]), int(result["day"]),
                    int(result["hour"]), int(result["minute"]),
                )
            except (ValueError, KeyError, TypeError):
                memory.mark_pending_reminder(pid, "dropped")
                terminal_written = True
                continue
            remind_at = int(remind_dt.timestamp())
            if remind_at < _dt.now().timestamp() - 3600:  # 抽出來已過期
                memory.mark_pending_reminder(pid, "dropped")
                terminal_written = True
                continue
            rid = memory.add_reminder(
                group_id, row["user_id"], result["action"], remind_at,
                source_text=row["text"][:200],
                mention_aliases=result.get("mention_aliases") or [],
            )
            memory.mark_pending_reminder(pid, "done")
            terminal_written = True
            if rid:
                logger.info(
                    "reminder drained: rid=%d action=%r at=%s",
                    rid, result["action"], remind_dt.strftime("%Y-%m-%d %H:%M"),
                )
        except Exception as e:
            if _is_quota_error(e):
                _mark_quota_exhausted()
                memory.release_pending_reminder(pid)
                break  # 額度又爆 → 停本輪，剩下的留待下次
            logger.warning(
                "drain reminders: release pending_id=%s after failure: %s",
                pid, str(e)[:160],
            )
            if not terminal_written:
                memory.release_pending_reminder(pid)
            continue


def _maybe_extract_reminder(
    text: str, group_id: str, user_id: str = "", message_id: str | None = None
) -> None:
    """偵測 user 訊息中的時間性提醒事項，存進 reminders table。失敗 silent。

    流程：regex pre-filter → Gemini light 抽取 → memory.add_reminder
    quota 爆時 extract_reminder raise 429 → 入隊 pending 等恢復補抽（site 2）。
    """
    if not text or len(text) > 500:
        return
    # 必須有日期 + 時間/行動 hint 才送 Gemini（避免「今天天氣」也燒 quota）
    if not _REMINDER_DATE_HINT.search(text):
        return
    if not _REMINDER_TIME_OR_ACTION_HINT.search(text):
        return
    gemini_allowed = _gemini_side_task_allowed("reminder_extract")
    try:
        result = gemini_client.extract_reminder(text) if gemini_allowed else None
        if result is None:
            from datetime import datetime as _dt
            result = _calendar_regex_to_reminder_result(
                text, _dt.now().date(), user_id
            )
            if result is None:
                if not gemini_allowed:
                    _enqueue_reminder_if_candidate(
                        text, group_id, user_id, message_id
                    )
                return
        else:
            actor = _infer_medical_actor(text, user_id)
            if actor and "action" in result:
                companions = _infer_medical_companions(text)
                result["action"] = _apply_medical_actor(
                    str(result["action"]), actor, companions
                )
                result["mention_aliases"] = _medical_mention_aliases(
                    actor, companions
                )
        # 算 remind_at（local timezone）
        from datetime import datetime as _dt
        try:
            remind_dt = _dt(
                int(result["year"]),
                int(result["month"]),
                int(result["day"]),
                int(result["hour"]),
                int(result["minute"]),
            )
        except (ValueError, KeyError, TypeError) as e:
            logger.info("reminder datetime parse failed: %s, result=%s", e, result)
            return
        remind_at = int(remind_dt.timestamp())
        # 過去事件不存（已過 1 hr 以上）
        if remind_at < _dt.now().timestamp() - 3600:
            return
        rid = memory.add_reminder(
            group_id, user_id, result["action"], remind_at, source_text=text[:200],
            mention_aliases=result.get("mention_aliases") or [],
        )
        if rid:
            logger.info(
                "reminder saved: rid=%d action=%r at=%s",
                rid, result["action"], remind_dt.strftime("%Y-%m-%d %H:%M"),
            )
    except Exception as e:
        if _is_quota_error(e):
            # site 2 (intra-request flip): 日額度爆 → mark + 入隊等恢復補抽
            _mark_quota_exhausted()
            _enqueue_reminder_if_candidate(text, group_id, user_id, message_id)
        elif "429" in str(e) or "RESOURCE_EXHAUSTED" in str(e):
            # transient per-minute 429 → 入隊但不 mark（不可壓死全天額度）
            _enqueue_reminder_if_candidate(text, group_id, user_id, message_id)
        else:
            logger.warning("_maybe_extract_reminder failed: %s", e)


# ── Command 處理 ──────────────────────────────────────────────────────────────

_DINNER_KEYWORDS = [
    "晚餐吃什麼",
    "晚餐吃哪",
    "吃什麼晚餐",
    "晚餐去哪",
    "晚餐要吃什麼",
    "今晚吃什麼",
    "今天吃什麼",
]


def _is_dinner_question(text: str) -> bool:
    return any(kw in text for kw in _DINNER_KEYWORDS)


_DINNER_PROMPT = """你是台北美食達人，以善導寺捷運站（台北市中正區）為中心，推薦附近步行可達的晚餐餐廳。

以下餐廳請勿推薦：喜來登、阜杭豆漿、雙月食品社。

請推薦 4～5 間，盡量多樣（台菜、日式、韓式、異國料理、麵食等皆可），格式如下（用換行分隔每間）：
🍽 餐廳名稱
📍 地址（簡短）
🍴 料理類型 ＋ 招牌菜或特色一句話
💰 價位（每人約 NT$XXX）

回覆風格：親切自然，像朋友推薦，繁體中文，不要加多餘的前言或結語。"""


_WEB_RESEARCH_QUESTION_HINTS = (
    "?",
    "？",
    "嗎",
    "呢",
    "如何",
    "怎麼",
    "怎樣",
    "哪",
    "什麼",
    "為什麼",
    "可不可以",
    "能不能",
    "可以",
    "支援",
    "相容",
    "推薦",
    "比較",
    "差在哪",
    "值得",
    "好不好",
    "有沒有",
    "有什麼",
    "最近",
)
_WEB_RESEARCH_PUBLIC_HINTS = (
    # time-sensitive / public facts
    "新聞",
    "近況",
    "趨勢",
    "行情",
    "走勢",
    # markets / economics
    "美股",
    "台股",
    "港股",
    "股票",
    "股市",
    "大盤",
    "匯率",
    "利率",
    "油價",
    "金價",
    "房市",
    "經濟",
    # place / weather / travel / local recommendations
    "氣候",
    "天氣",
    "溫度",
    "降雨",
    "季節",
    "好玩",
    "景點",
    "旅遊",
    "旅行",
    "自由行",
    "行程",
    "必去",
    "交通",
    "簽證",
    "餐廳",
    "住宿",
    # products / software / compatibility
    "支援",
    "相容",
    "規格",
    "版本",
    "安裝",
    "價格",
    "評價",
    "評測",
    "m1",
    "m2",
    "m3",
    "m4",
    "apple silicon",
    "mac",
    "晶片",
    # official / policy / availability
    "政策",
    "法規",
    "規定",
    "開放",
    "營業",
    "票價",
    "官方",
    "來源",
    "資料",
)
_WEB_RESEARCH_RECOMMENDATION_RE = re.compile(
    r"(哪裡|哪邊|哪兒|哪里|哪個|有什麼|有哪些|推薦|好玩|景點|怎麼去|怎麼玩)"
)
_WEB_RESEARCH_DEFINITION_RE = re.compile(r"(什麼是|是什麼|介紹一下|解釋一下)")
_WEB_RESEARCH_SUBJECT_STOPWORDS_RE = re.compile(
    r"(請問|你覺得|你認為|可以|可不可以|能不能|嗎|呢|如何|怎麼|怎樣|"
    r"哪裡|哪邊|哪兒|哪里|哪個|有什麼|有哪些|推薦|好玩|景點|怎麼去|怎麼玩|"
    r"爸爸|媽媽|阿公|阿嬤|奶奶|爺爺|哥哥|姐姐|妹妹|弟弟|你|我|他|她|它|這個|那個)"
)


def _has_web_research_subject(text: str) -> bool:
    subject = _WEB_RESEARCH_SUBJECT_STOPWORDS_RE.sub("", text or "")
    subject = re.sub(r"[\s?？。！!，,、：:；;（）()]+", "", subject)
    return len(subject) >= 2


def _is_web_research_question(text: str) -> bool:
    """Detect public-info questions worth answering immediately with web research."""
    s = (text or "").strip()
    if not s or len(s) > 180:
        return False
    lower = s.lower()
    has_question = any(h in lower for h in _WEB_RESEARCH_QUESTION_HINTS)
    if not has_question:
        return False
    if _extract_prefetch_urls(s):
        return True
    if _WEB_RESEARCH_DEFINITION_RE.search(s):
        return _has_web_research_subject(s)
    if _WEB_RESEARCH_RECOMMENDATION_RE.search(s):
        return _has_web_research_subject(s)
    return any(h in lower for h in _WEB_RESEARCH_PUBLIC_HINTS)


def _build_web_research_queries(text: str) -> list[str]:
    base = re.sub(r"\s+", " ", (text or "").strip(" \t\r\n?？。！!"))
    if not base:
        return []
    lower = base.lower()
    queries = [base]
    if any(h in base for h in ("最近", "最新", "今天", "現在", "目前", "新聞", "行情", "走勢")):
        queries.append(f"{base} 最新 2026")
    if any(h in base for h in ("氣候", "天氣", "旅遊", "旅行", "景點", "行程", "簽證", "交通")):
        queries.append(f"{base} 官方 旅遊資訊")
    if any(h in lower for h in ("支援", "相容", "規格", "版本", "安裝", "m1", "m2", "m3", "m4", "apple silicon", "mac")):
        queries.append(f"{base} official support")

    deduped: list[str] = []
    seen: set[str] = set()
    for q in queries:
        if q and q not in seen:
            deduped.append(q)
            seen.add(q)
    return deduped[:3]


def _collect_web_research_sources(text: str) -> list[dict]:
    queries = _build_web_research_queries(text)
    if not queries:
        return []
    try:
        import source_aggregator
        sources = source_aggregator.aggregate_sources(queries, total_max=6)
    except Exception as e:
        logger.info("web research source aggregation failed: %s", e)
        return []
    if not sources:
        return []
    try:
        import fulltext_fetcher
        return fulltext_fetcher.fetch_top_sources(
            sources,
            top_n=2,
            max_chars_per=1200,
            timeout_per_task=3.0,
            max_workers=2,
        )
    except Exception as e:
        logger.info("web research full-text enrichment skipped: %s", e)
        return sources


def _format_web_research_sources(sources: list[dict], limit: int = 5) -> str:
    if not sources:
        return ""
    blocks: list[str] = []
    for idx, src in enumerate(sources[:limit], 1):
        title = str(src.get("title") or "").strip()
        url = str(src.get("url") or "").strip()
        domain = str(src.get("domain") or "").strip()
        body = str(src.get("full_text") or src.get("snippet") or "").strip()
        if len(body) > 700:
            body = body[:700] + "..."
        parts = [f"[{idx}] {title or domain or url}"]
        if domain:
            parts.append(f"domain: {domain}")
        if url:
            parts.append(f"url: {url}")
        if body:
            parts.append(f"內容摘錄: {body}")
        blocks.append("\n".join(parts))
    return "\n\n".join(blocks)


def _build_web_research_prompt(text: str, sources: list[dict] | None = None) -> str:
    today_tw = datetime.now(ZoneInfo("Asia/Taipei")).date().isoformat()
    source_block = _format_web_research_sources(sources or [])
    if not source_block:
        source_block = (
            "本機爬蟲沒有抓到足夠資料。請使用 Google Search grounding 補查，"
            "再用查到的資料回答；不要只靠模型記憶。"
        )
    return (
        "這是一則 LINE 群組裡沒有明確 @ bot、但明顯在問公開資料或推薦的問題。\n"
        "請根據本機爬蟲資料回答；如果資料不足或題目涉及最新狀態，"
        "請再使用 Google Search / 可用網路資料補查。\n"
        "回答規則：第一句直接給具體判斷或建議；用繁體中文；"
        "控制在 5-9 行；不要空泛重述問題；事實或近期資訊要附簡短來源名稱或網址。\n"
        "如果問題範圍太大，先分成幾個最實用方向回答，最後指出要縮小範圍才能排細節。\n\n"
        f"今天台灣日期：{today_tw}\n"
        f"使用者原文：{text.strip()}\n\n"
        f"本機爬蟲資料：\n{source_block}"
    )


def _handle_web_research_question(
    event: MessageEvent,
    group_id: str,
    text: str,
) -> bool:
    """Immediate no-mention path for public-info questions; returns True when consumed."""
    sender_user_id = getattr(event.source, "user_id", None) or ""
    context = memory.get_context(group_id)
    facts = memory.top_facts(group_id, user_id=sender_user_id)
    pnotes = _get_persona_notes(group_id)
    try:
        with _thinking_indicator(group_id):
            sources = _collect_web_research_sources(text)
            research_text = _prefetch_urls(text)
            prompt_text = _build_web_research_prompt(research_text, sources)
            reply_text = _llm_chat(prompt_text, context, facts, pnotes)
    except Exception as e:
        if _is_quota_error(e):
            _mark_quota_exhausted()
            logger.warning("web research question quota exhausted, retry via fallback")
            try:
                sources = _collect_web_research_sources(text)
                research_text = _prefetch_urls(text)
                prompt_text = _build_web_research_prompt(research_text, sources)
                reply_text = _llm_chat(prompt_text, context, facts, pnotes)
            except Exception as e2:
                logger.warning("web research question fallback failed: %s", e2)
                return True
        else:
            logger.warning("web research immediate path failed: %s", e)
            return False

    if not reply_text or not reply_text.strip():
        logger.info("web research question returned empty group=%s", group_id)
        return True

    memory.append_turn(group_id, "user", text)
    _append_bot_turn(group_id, reply_text)
    _reply(event.reply_token, reply_text, group_id=group_id)
    return True


def _handle_dinner_recommendation(event: MessageEvent, group_id: str) -> None:
    # 2026-05-16 改：刪 quota 短路。dinner prompt 是純 text，llm_router.fallback_chat
    # 4-tier waterfall (local LLM → RAG → lite_reply) 能處理；對齊
    # feedback_quota_fallback_never_skip.md。注意 file handler 因含 PDF/image bytes、
    # _extract_text 會丟，仍保留短路 + 友善訊息（known exception）。
    context = memory.get_context(group_id)
    facts = memory.top_facts(group_id)
    pnotes = _get_persona_notes(group_id)
    try:
        with _thinking_indicator(group_id):
            reply_text = _llm_chat(_DINNER_PROMPT, context, facts, pnotes)
    except Exception as e:
        if _is_quota_error(e):
            _mark_quota_exhausted()
            logger.warning("dinner recommendation quota exhausted, retry via local")
            reply_text = _local_text_llm_fallback(_DINNER_PROMPT, context=context)
        elif _is_gemini_unavailable_error(e):
            logger.warning(
                "dinner recommendation unavailable, retry via local: %s",
                e,
            )
            reply_text = _local_text_llm_fallback(_DINNER_PROMPT, context=context)
        else:
            logger.exception("dinner recommendation failed: %s", e)
            _reply(event.reply_token, _friendly_gemini_error(e), group_id=group_id)
            return
    if not reply_text or not reply_text.strip():
        # fallback_chat 全敗回空 → 給使用者明確訊息（不能送空到 LINE SDK）
        _reply(event.reply_token, "晚餐推薦今天罷工了，等一下再試試", group_id=group_id)
        return
    _reply(event.reply_token, reply_text, group_id=group_id)


_CLASSIFY_EMOJIS = {
    "財經": "💰", "健康": "🏥", "警示": "🚨", "影片": "🎬", "行程": "📅", "其他": "📌",
}
_CLASSIFY_CMD_TO_CAT = {
    "/今日財經": "財經",
    "/今日健康": "健康",
    "/今日警示": "警示",
    "/今日影片": "影片",
    "/今日行程": "行程",
    "/今日其他": "其他",
}


def _handle_classify_command(group_id: str, text: str) -> str | None:
    """處理 /分類 與 /今日X 系列指令；無命中回 None。2026-05-10 加。"""
    t = text.strip()
    if t not in _CLASSIFY_CMD_TO_CAT and t != "/分類":
        return None
    try:
        import message_classifier
        import time as _time
        if t == "/分類":
            counts = message_classifier.today_category_counts(group_id)
            lines = ["📊 今日訊息分類"]
            for cat in message_classifier.CATEGORIES:
                n = counts.get(cat, 0)
                lines.append(f"{_CLASSIFY_EMOJIS.get(cat, '')} {cat}: {n} 則")
            lines.append("")
            lines.append("查詢：/今日財經 /今日健康 /今日警示 /今日影片 /今日行程 /今日其他")
            return "\n".join(lines)
        cat = _CLASSIFY_CMD_TO_CAT[t]
        rows = message_classifier.list_today_in_category(group_id, cat, limit=20)
        emoji = _CLASSIFY_EMOJIS.get(cat, "")
        if not rows:
            return f"{emoji} 今日{cat}：沒有訊息"
        out = [f"{emoji} 今日{cat} ({len(rows)} 則)"]
        for i, r in enumerate(rows, 1):
            ts = _time.strftime("%H:%M", _time.localtime(r["created_at"]))
            txt = (r["text"] or "")[:60]
            out.append(f"{i}. {ts} {txt}")
        return "\n".join(out)
    except Exception as e:
        logger.warning("_handle_classify_command failed: %s", e)
        return None


def _handle_finance_view_command(group_id: str, text: str) -> str | None:
    """處理 /觀點 [可選: 家人名 | 標的]。純 SQL 聚合，不過 Gemini。

    第一句必須是判斷句（AGENTS.md 的 LINE bot 規則 0）。
    """
    s = (text or "").strip()
    if not s.startswith("/觀點"):
        return None

    try:
        import finance_view_db
        import stock_quote
    except ImportError as e:
        logger.warning("finance_view import failed: %s", e)
        return None

    args = s.replace("/觀點", "", 1).strip()

    if not args:
        views = finance_view_db.list_recent(group_id, limit=10)
        if not views:
            return (
                "目前家族觀點庫還是空的。\n"
                "（聊到具體標的 + 方向時會自動記下，例「我覺得 0050 會漲到 180」）"
            )
        return _format_finance_views(views, header="📈 家族最近財經觀點")

    tokens = args.split()
    ticker = None
    person = None
    for tok in tokens:
        try:
            resolved = stock_quote.detect_symbols(tok)
        except Exception:
            resolved = None
        if resolved:
            ticker = resolved[0]
        elif not person:
            person = tok

    if ticker:
        views = finance_view_db.list_by_ticker(group_id, ticker, limit=10)
        header = f"📈 {ticker} 相關家族觀點"
    elif person:
        views = finance_view_db.list_by_person(group_id, person, limit=10)
        header = f"📈 {person} 的財經觀點"
    else:
        return None

    if not views:
        return f"{header}\n\n沒找到相關觀點記錄。"
    return _format_finance_views(views, header=header)


def _format_finance_views(views: list[dict], header: str) -> str:
    """純 SQL 聚合輸出。第一句判斷句（規則 0）。"""
    if not views:
        return f"{header}\n\n沒有記錄。"

    from datetime import datetime as _dt
    from zoneinfo import ZoneInfo as _ZI

    tw = _ZI("Asia/Taipei")

    hit = sum(1 for v in views if v.get("validation_result") == "hit")
    miss = sum(1 for v in views if v.get("validation_result") == "miss")
    pending = sum(1 for v in views if v.get("validation_result") == "pending")
    na = sum(1 for v in views if v.get("validation_result") == "na")

    if hit and miss:
        summary = f"共 {len(views)} 條觀點，已驗證 {hit} 命中 / {miss} 落空，{pending} 仍待驗。"
    elif hit:
        summary = f"共 {len(views)} 條觀點，已驗證 {hit} 命中（其餘待驗）。"
    elif miss:
        summary = f"共 {len(views)} 條觀點，已驗證 {miss} 落空。"
    else:
        summary = f"共 {len(views)} 條觀點，{pending + na} 仍待驗。"

    lines = [header, "", summary, ""]
    for v in views[:10]:
        label = v.get("ticker") or v.get("macro_topic") or "?"
        d = v.get("direction") or ""
        dir_str = {"bull": "看多", "bear": "看空", "neutral": "持平"}.get(d, "")
        target = ""
        if v.get("target_price"):
            target = f" 目標 {v['target_price']}"
        elif v.get("target_pct"):
            target = f" 目標 {v['target_pct']:+.0f}%"
        result = v.get("validation_result") or ""
        result_str = {"hit": "✅", "miss": "❌", "pending": "⏳", "na": "—"}.get(result, "")
        try:
            created = _dt.fromtimestamp(v["created_at"] / 1000, tz=tw).strftime("%m-%d")
        except Exception:
            created = "??"
        speaker = v.get("display_name") or "家人"
        lines.append(f"• {created} {speaker} {label} {dir_str}{target} {result_str}")
    return "\n".join(lines)


def _handle_food_command(group_id: str, text: str) -> str | None:
    """家族飲食 / 食材媒合指令（純 DB 查詢，跳過 Gemini）。無命中回 None。

    嚴格 == 比對（GP1 C4：寬鬆比對會短路後面所有既有指令）。
    輸出只給全家層級、不點名個人（GP2 A）。
    """
    t = text.strip()
    if t in ("/今晚煮什麼", "/煮什麼", "/食材媒合"):
        import food_db
        import food_recipes
        inventory = food_db.query_inventory(group_id)
        prefs = food_db.query_prefs(group_id)
        msg = food_recipes.format_suggestions(inventory, prefs.get("dislikes"))
        if not msg:
            return (
                "🍽️ 目前還沒記錄到家裡有什麼食材～\n"
                "等大家在群組聊到「冰箱有…」「買了…」我就會記起來囉"
            )
        return msg
    if t in ("/該買什麼", "/採購清單"):
        import food_db
        shopping = food_db.query_shopping(group_id)
        if not shopping:
            return "🛒 目前沒有待買的食材"
        return "🛒 待買清單：\n" + "\n".join(f"・{f}" for f in shopping)
    if t in ("/家裡有什麼", "/庫存"):
        import food_db
        inventory = food_db.query_inventory(group_id)
        if not inventory:
            return "🧊 目前還沒記錄到家裡的食材"
        return "🧊 家裡現有：\n" + "、".join(inventory)
    return None


def _handle_poll_command(
    group_id: str,
    text: str,
    user_id: str | None = None,
    message_id: str = "",
) -> str | None:
    """Group poll commands. Slash commands remain available without a bot mention."""
    t = (text or "").strip()
    poll_command_prefixes = (
        "/民調",
        "/投票",
        "/催民調",
        "/提醒民調",
        "/關閉民調",
        "/結束民調",
        "/取消民調",
        "幫我做民調",
        "幫我們做民調",
        "幫大家做民調",
        "幫我開民調",
        "幫我們開民調",
        "做民調",
        "開民調",
        "建立民調",
        "發起民調",
        "弄民調",
        "關掉民調",
        "關閉民調",
        "結束民調",
        "取消民調",
        "停掉民調",
        "停止民調",
        "民調關掉",
        "民調關閉",
        "民調結束",
        "民調取消",
    )
    if not t.startswith(poll_command_prefixes):
        return None
    try:
        import family_poll

        return family_poll.handle_explicit_message(
            group_id,
            t,
            user_id=user_id,
            sender_alias=_get_member_display_name(group_id, user_id) if user_id else "",
            source_msg_id=message_id,
        )
    except Exception as e:
        logger.warning("poll command failed: %s", e)
        return None


def _handle_explicit_poll_text(
    event: MessageEvent, group_id: str, text: str
) -> str | None:
    """Poll creation/votes only after an explicit bot trigger."""
    try:
        import family_poll

        user_id = getattr(event.source, "user_id", None) or ""
        msg_id = getattr(event.message, "id", "") or ""
        return family_poll.handle_explicit_message(
            group_id,
            text,
            user_id=user_id,
            sender_alias=_get_member_display_name(group_id, user_id) if user_id else "",
            source_msg_id=msg_id,
        )
    except Exception as e:
        logger.warning("explicit poll handler failed: %s", e)
        return None


def _handle_poll_text(event: MessageEvent, group_id: str, text: str) -> str | None:
    """Legacy natural poll handler; ordinary chat no longer calls this path."""
    try:
        import family_poll

        user_id = getattr(event.source, "user_id", None) or ""
        msg_id = getattr(event.message, "id", "") or ""
        return family_poll.handle_natural_message(
            group_id,
            text,
            user_id=user_id,
            sender_alias=_get_member_display_name(group_id, user_id) if user_id else "",
            source_msg_id=msg_id,
        )
    except Exception as e:
        logger.warning("poll text handler failed: %s", e)
        return None


def _handle_command(
    group_id: str,
    text: str,
    user_id: str | None = None,
    message_id: str = "",
) -> str | None:
    """有對應到指令回 str；沒有回 None。"""
    t = text.strip()

    poll_reply = _handle_poll_command(group_id, t, user_id, message_id)
    if poll_reply is not None:
        return poll_reply

    # 家族飲食 / 食材媒合（2026-05-31 加；嚴格 == 無命中回 None，純 DB 跳過 Gemini）
    food_reply = _handle_food_command(group_id, t)
    if food_reply is not None:
        return food_reply

    # 訊息分類查詢（2026-05-10 加）
    classify_reply = _handle_classify_command(group_id, t)
    if classify_reply is not None:
        return classify_reply

    # 家族財經觀點查詢（2026-05-18 加）
    fv_reply = _handle_finance_view_command(group_id, t)
    if fv_reply is not None:
        return fv_reply

    if t == "/group_id":
        return f"本群 group_id：\n{group_id}"

    if t == "/help" or t == "/指令":
        return _HELP_TEXT

    # ── 長期記憶（facts）──────────────────────────────────────────
    if t == "/看記憶":
        facts = memory.list_facts(group_id)
        if not facts:
            return "目前沒有任何記憶。要讓我記住什麼，用：\n/記住 <內容>"
        return "目前的記憶：\n" + "\n".join(f"• {f}" for f in facts)

    if t == "/記住":
        return "用法：/記住 <要記住的內容>"
    if t.startswith("/記住 "):
        fact = t[len("/記住 ") :].strip()
        if not fact:
            return "用法：/記住 <要記住的內容>"
        if memory.add_fact(group_id, fact):
            return f"好，記住了：{fact}"
        return f"這條已經在記憶裡了：{fact}"

    if t == "/忘記":
        return "用法：/忘記 <關鍵字>"
    if t.startswith("/忘記 "):
        keyword = t[len("/忘記 ") :].strip()
        if not keyword:
            return "用法：/忘記 <關鍵字>"
        n = memory.remove_fact(group_id, keyword)
        return (
            f"刪除了 {n} 條含「{keyword}」的記憶。"
            if n
            else f"沒有找到含「{keyword}」的記憶。"
        )

    if t == "/清除記憶":
        n = memory.clear_facts(group_id)
        return f"已清除 {n} 條記憶。"

    # ── Layer 1：使用者手動管理過濾規則 ──────────────────────────────
    if t == "/不要回":
        return "用法：/不要回 <這類訊息的特徵，例如「早安」「中午吃什麼」>"
    if t.startswith("/不要回 "):
        pattern = t[len("/不要回 ") :].strip()
        if not pattern:
            return "用法：/不要回 <這類訊息的特徵，例如「早安」「中午吃什麼」>"
        rid = memory.add_filter_rule(group_id, "skip", pattern, source="user")
        return f"好，以後訊息裡有「{pattern}」就不主動回。(規則 #{rid})"

    if t == "/以後要查":
        return "用法：/以後要查 <這類訊息的特徵，例如「某醫師說」「疫苗」>"
    if t.startswith("/以後要查 "):
        pattern = t[len("/以後要查 ") :].strip()
        if not pattern:
            return "用法：/以後要查 <這類訊息的特徵，例如「某醫師說」「疫苗」>"
        rid = memory.add_filter_rule(group_id, "must_answer", pattern, source="user")
        return f"好，以後訊息裡有「{pattern}」就會主動查證。(規則 #{rid})"

    if t == "/規則":
        rules = memory.list_filter_rules(group_id)
        if not rules:
            return "目前沒有過濾規則。\n新增：/不要回 <特徵>  或  /以後要查 <特徵>"
        lines = ["目前的過濾規則："]
        for r in rules:
            tag = "不要回" if r["kind"] == "skip" else "要查"
            src = "手動" if r["source"] == "user" else "自動學"
            lines.append(f"#{r['rule_id']} [{tag}]({src}) {r['pattern']}")
        return "\n".join(lines)

    if t.startswith("/刪除規則 "):
        raw = t[len("/刪除規則 ") :].strip()
        try:
            rid = int(raw)
        except ValueError:
            return "用法：/刪除規則 <數字編號>（用 /規則 看編號）"
        if memory.delete_filter_rule(group_id, rid):
            return f"已刪除規則 #{rid}"
        return f"找不到規則 #{rid}"

    if t == "/清除規則":
        n = memory.clear_filter_rules(group_id)
        return f"已清除 {n} 條過濾規則"

    # ── Layer 3：週期性自我檢討 ──────────────────────────────────────
    if t == "/檢討" or t == "/檢討 7":
        report, _ = review.run_weekly_review(group_id, days=7)
        return report

    if t.startswith("/檢討 "):
        raw = t[len("/檢討 ") :].strip()
        try:
            days = int(raw)
        except ValueError:
            return "用法：/檢討 <天數>  例如：/檢討 14"
        if days <= 0 or days > 30:
            return "天數請在 1~30 之間。"
        report, _ = review.run_weekly_review(group_id, days=days)
        return report

    if t == "/採用":
        drafts = memory.list_rule_drafts(group_id)
        if not drafts:
            return "目前沒有待採用的建議。先跑 /檢討 產生一份。"
        lines = ["目前的建議："]
        for d in drafts:
            tag = "不要回" if d["kind"] == "skip" else "要回"
            lines.append(f"{d['draft_id']}. [{tag}] {d['pattern']}")
            if d.get("reason"):
                lines.append(f"   理由：{d['reason']}")
        lines.append("")
        lines.append("用法：/採用 1 2  或  /採用 全部  或  /採用 無")
        return "\n".join(lines)

    if t.startswith("/採用 "):
        spec = t[len("/採用 ") :].strip()
        _, msg = review.adopt_drafts(group_id, spec)
        return msg

    # ── Layer 2：糾正剛剛的 bot 回覆 → 自動抽象成規則 ────────────────
    if t.startswith("/閉嘴"):
        reason = t[len("/閉嘴") :].lstrip()
        if not reason:
            return "用法：/閉嘴 <為什麼不應該回>\n例：/閉嘴 這種只是早安問候，不用回"
        return _handle_layer2_correction(group_id, reason)

    # ── 家族行事曆 ──────────────────────────────────────────────
    if t in ("/待辦", "/提醒事項", "/提醒清單"):
        return _build_todo_status_reply(group_id, t)

    if t in ("/行事曆", "/活動", "/聚餐"):
        return _format_calendar(group_id)

    if t.startswith("/取消活動 "):
        kw = t[len("/取消活動 ") :].strip()
        if not kw:
            return "用法：/取消活動 <活動關鍵字>"
        return _cancel_calendar_event(group_id, kw)

    return None


def _format_calendar(group_id: str) -> str:
    import calendar_db

    events = calendar_db.list_upcoming(group_id, days=30)
    if not events:
        return "📅 未來 30 天沒有家族活動。"
    lines = ["📅 **家族行事曆（未來 30 天）**"]
    for e in events:
        time_part = f" {e['event_time']}" if e["event_time"] else ""
        loc_part = f" @ {e['location']}" if e["location"] else ""
        try:
            import json as _j

            parts = _j.loads(e["participants"] or "[]")
        except Exception:
            parts = []
        ppl = "、".join(parts) if parts else ""
        ppl_part = f"（{ppl}）" if ppl else ""
        lines.append(
            f"• {e['event_date']}{time_part} {e['title']}{loc_part}{ppl_part}"
        )
    return "\n".join(lines)


def _cancel_calendar_event(group_id: str, keyword: str) -> str:
    import calendar_db

    target = calendar_db.find_active_event(group_id, keyword=keyword)
    if not target:
        return f"找不到含「{keyword}」的活動。用 /行事曆 看清單。"
    if calendar_db.cancel_event(target["event_id"]):
        return f"已取消：{target['event_date']} {target['title']}"
    return f"取消失敗（活動可能已被取消）：{target['title']}"


_HELP_TEXT = (
    "可用指令：\n"
    "【飲食】\n"
    "  /今晚煮什麼             用家裡現有食材推薦菜色\n"
    "  /該買什麼               待買食材清單\n"
    "  /家裡有什麼             目前記錄到的食材\n"
    "【民調】\n"
    "  /民調 <問題>            開一個群組民調，會 @all\n"
    "  /民調                   看目前民調統計\n"
    "  /催民調                 @all 提醒還沒回覆的人\n"
    "  /關閉民調               關閉目前民調\n"
    "【記憶】\n"
    "  /看記憶                 看長期事實\n"
    "  /記住 <內容>            手動加一條事實\n"
    "  /忘記 <關鍵字>          刪除含關鍵字的事實\n"
    "  /清除記憶               全砍\n"
    "【主動過濾】\n"
    "  /規則                   看過濾規則\n"
    "  /不要回 <特徵>          以後訊息裡有這個就不回\n"
    "  /以後要查 <特徵>        以後訊息裡有這個就主動查證\n"
    "  /刪除規則 <編號>        刪掉特定規則\n"
    "  /清除規則               全砍\n"
    "  /閉嘴 <理由>            針對剛剛那則 bot 回覆糾正,我會自動學一條規則\n"
    "【週期性自我檢討】\n"
    "  /檢討                   立刻跑一次過去 7 天的檢討(可接天數)\n"
    "  /採用                   列出待採用的建議\n"
    "  /採用 1 2 / 全部 / 無   把建議升級成正式規則\n"
    "【家族行事曆】\n"
    "  /待辦                   列出待辦與提醒事項\n"
    "  /行事曆                 列出未來 30 天的家族活動\n"
    "  /取消活動 <關鍵字>      取消含關鍵字的活動\n"
    "【其他】\n"
    "  /group_id               顯示本群 ID\n"
    "  /help                   看這張說明"
)


def _handle_layer2_correction(group_id: str, reason: str) -> str:
    """使用者覺得剛剛 bot 回覆不該出現 → 呼叫 Gemini 抽象一條 skip 規則。"""
    last = memory.get_last_bot_reply(group_id)
    if last is None:
        return "找不到最近的 bot 回覆可以糾正。"
    _, bot_reply = last
    trigger_text = _guess_last_trigger_text(group_id)

    pattern = gemini_client.generate_filter_rule(
        bot_reply=bot_reply,
        user_reason=reason,
        trigger_text=trigger_text,
    )
    if not pattern:
        return (
            "自動生成規則失敗，請改用 /不要回 <特徵> 手動加。\n"
            f"(你剛才說：{reason[:80]})"
        )
    rid = memory.add_filter_rule(group_id, "skip", pattern, source="learned")
    return (
        f"了解。我從這次糾正學到一條規則：\n"
        f"#{rid} [不要回] {pattern}\n"
        f"以後類似訊息就不會主動回了。覺得不對請用 /刪除規則 {rid}"
    )


def _guess_last_trigger_text(group_id: str) -> str:
    """找出最近一次 bot 回覆前，它看到的 user 訊息（當作 trigger 傳給規則產生器）。"""
    recent = memory.get_recent_raw_messages(group_id, limit=20)  # 舊→新
    last_bot_idx = None
    for i in range(len(recent) - 1, -1, -1):
        if recent[i][1] == "__bot__":
            last_bot_idx = i
            break
    if last_bot_idx is None:
        return ""
    before = [recent[i][2] for i in range(last_bot_idx) if recent[i][1] != "__bot__"]
    return "\n".join(before[-5:])


# ── LINE SDK helpers ──────────────────────────────────────────────────────────


def _is_mentioned(message: TextMessageContent) -> bool:
    """檢查這則訊息是否 @mention 了本 bot。
    LINE 的 mention 結構：message.mention.mentionees[i].is_self == True 代表 mention 到我。"""
    mention = getattr(message, "mention", None)
    if mention is None:
        return False
    mentionees = getattr(mention, "mentionees", None) or []
    for m in mentionees:
        if getattr(m, "is_self", False):
            return True
    return False


# 桌機 LINE 打不到 @bot 的後備觸發前綴（全大小寫都接受）
_ASK_PREFIXES = ("/ai ", "/ai", "/問 ", "/問", "/ask ", "/ask", "/AI ", "/AI")

# 桌機 LINE 有時候 @mention 不帶 mention 結構，只是純文字 @名稱。
# 列出 bot 名稱 + 通用 @AI 當 fallback。
# 桌機 LINE 會打出全形 ＠（U+FF20），所以半形全形都要接。
_TEXT_MENTION_PREFIXES = (
    "@咪寶 ",
    "@咪寶",
    "＠咪寶 ",
    "＠咪寶",
    "@ai ",
    "@ai",
    "＠ai ",
    "＠ai",
    "@AI ",
    "@AI",
    "＠AI ",
    "＠AI",
)

# 直接叫名字也算觸發（長輩不用 @，直接說「咪寶...」）
_BOT_NAME_KEYWORDS = ("咪寶",)


def _extract_gemini_trigger(text: str, message: TextMessageContent) -> str | None:
    """判斷這則訊息是否要丟給 Gemini；若是，回傳乾淨的問題文字。

    四種觸發方式（回 None 代表無視）：
    1. 手機 LINE：@ 本 bot，會有 mention 結構 → 去掉 mention 後剩下的字
    2. /ai、/問、/ask 前綴 → 去掉前綴後剩下的字
    3. 桌機 LINE fallback：純文字 @AI 開頭（沒有 mention 結構）→ 去掉前綴
    4. 訊息裡出現 bot 名字（咪寶）→ 去掉名字後剩下的字
    """
    t = text.strip()
    for prefix in _ASK_PREFIXES:
        if t == prefix.strip():
            return ""
        if t.startswith(prefix):
            return t[len(prefix) :].strip()
    if _is_mentioned(message):
        return _strip_mentions(message).strip()
    # fallback：桌機 LINE @mention 不帶結構，只有純文字 @AI
    for prefix in _TEXT_MENTION_PREFIXES:
        if t == prefix.strip():
            return ""
        if t.lower().startswith(prefix.lower()):
            return t[len(prefix) :].strip()
    # 名字偵測：訊息裡出現 bot 名字就觸發，把名字挖掉後剩下的當問題
    for name in _BOT_NAME_KEYWORDS:
        if name in t:
            clean = t.replace(name, "", 1).strip("，,、。！!？? \t")
            return clean
    return None


def _strip_mentions(message: TextMessageContent) -> str:
    """把訊息裡所有 @mention 的子字串挖掉，只留真正的問題。"""
    text = message.text or ""
    mention = getattr(message, "mention", None)
    if mention is None:
        return text
    mentionees = getattr(mention, "mentionees", None) or []
    # 從後往前刪，避免 index 位移
    ranges = sorted(
        [(m.index, m.index + m.length) for m in mentionees],
        key=lambda x: x[0],
        reverse=True,
    )
    for start, end in ranges:
        text = text[:start] + text[end:]
    return text


def _md_to_line(text: str) -> str:
    """把 Gemini 回傳的 Markdown 語法轉成 LINE 能直接閱讀的純文字。"""
    lines = text.splitlines()
    out = []
    in_code_block = False
    for line in lines:
        # code block fence
        if line.strip().startswith("```"):
            in_code_block = not in_code_block
            continue
        if in_code_block:
            out.append(line)
            continue
        # 水平分隔線
        if re.match(r"^\s*[-*_]{3,}\s*$", line):
            out.append("")
            continue
        # headers → 加底線感
        m = re.match(r"^#{1,3}\s+(.+)", line)
        if m:
            out.append(f"▌ {m.group(1)}")
            continue
        # blockquote
        line = re.sub(r"^>\s*", "", line)
        # bullet * / - → •（只處理行首）
        line = re.sub(r"^(\s*)[*-]\s+", r"\1• ", line)
        # bold **text** / __text__
        line = re.sub(r"\*\*(.+?)\*\*", r"\1", line)
        line = re.sub(r"__(.+?)__", r"\1", line)
        # italic *text* / _text_（小心不要吃掉 URL 或數學符號）
        line = re.sub(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", r"\1", line)
        line = re.sub(r"(?<!_)_(?!_)(.+?)(?<!_)_(?!_)", r"\1", line)
        # inline code
        line = re.sub(r"`(.+?)`", r"\1", line)
        # [text](url) → text（url）
        line = re.sub(r"\[([^\]]+)\]\((https?://[^\)]+)\)", r"\1（\2）", line)
        out.append(line)
    return "\n".join(out)


def _peek_text_pending_for_drain(
    group_id: str, max_count: int, deadline_sec: float
) -> list[tuple[str, str]]:
    """Peek up to max_count text pending items, render LLM reply (local fallback OK),
    return list of (rendered_text, original_msg_id) WITHOUT removing from pending.

    Caller 必須在 reply_message 成功後 call _commit_pending_removal(group_id, msg_ids)
    才會真的把 entries 從 pending 拿掉 (peek-then-confirm pattern, §3 GP1 C5)。

    Honors deadline_sec — 中途逾時 break。LLM 失敗的 entry skip (不阻塞後續)。
    純 text；file/image pending 不在此快速路徑處理。
    """
    pending = _load_pending_explicit()
    items = pending.get(group_id, [])
    text_items = [
        it for it in items
        if it.get("type") == "text" and (it.get("text") or "").strip()
    ]
    if not text_items:
        return []

    facts = memory.top_facts(group_id)
    context = memory.get_context(group_id)
    pnotes = _get_persona_notes(group_id)

    started = time.time()
    results: list[tuple[str, str]] = []
    for item in text_items[:max_count]:
        if time.time() - started > deadline_sec:
            logger.info(
                "peek drain deadline %.1fs hit, stopping at %d/%d items group=%s",
                deadline_sec, len(results), max_count, group_id,
            )
            break
        original = item["text"].strip()
        try:
            reply_text = _llm_chat(original, context, facts, pnotes)
        except Exception as e:
            logger.warning("peek drain LLM failed for one item: %s", str(e)[:120])
            continue
        if not reply_text:
            continue
        orig_preview = original[:300] + ("…" if len(original) > 300 else "")
        rendered = _prepare_outbound_text(
            reply_text, source="pending_peek"
        )[:1800]   # 收緊到 1800 (§3 codex Q3 / GP1)
        formatted = (
            f"📬 補回之前漏掉的訊息\n\n原文：\n{orig_preview}\n\n回應：\n{rendered}"
        )
        results.append((formatted[:4900], item.get("message_id", "")))
    return results


def _try_piggyback_drain_with_reply_token(
    reply_token: str | None, group_id: str
) -> None:
    """Gemini 爆時利用 incoming webhook 的 reply_token，把最多 2 條 pending 透過
    本機 LLM fallback 生回覆，bundle 進**單一** reply_message API call (LINE 免費，
    不耗月配額)。peek-then-confirm pattern：reply 失敗 pending 完整保留。

    Critical guards (per §3 codex + GP1 review):
      - drain lock 跟 _drain_pending_for_group 共用 (避免 race + lost-update)
      - deadline 6s (webhook 是 async def 但 _handle_event 是 sync call，太久會
        block uvicorn event loop)
      - max_count=2 (兩條 local LLM call ≈ 2-10s 內可接受 webhook latency)
      - peek-then-confirm (reply API fail → pending 不動，下次 webhook 再試)
      - 中文 char cap 1800 + total reply body 保守上限

    Silent on: 缺 reply_token / bot_muted / 無 pending / drain lock 拿不到 /
    所有 LLM 失敗 / reply_message API 失敗。
    """
    if not _pending_reply_enabled():
        return
    if not reply_token or not group_id:
        return
    if settings.bot_muted:
        return
    if not _load_pending_explicit().get(group_id):
        return

    slot = _try_acquire_drain_slot(group_id)
    if slot is None:
        logger.info(
            "quota-exhausted piggyback: another drain in progress group=%s",
            group_id,
        )
        return
    try:
        import pending_store as _ps
        _ps.ensure_message_ids(group_id)
        rendered = _peek_text_pending_for_drain(
            group_id, max_count=2, deadline_sec=6.0
        )
        if not rendered:
            return

        messages: list = []
        for text, _msg_id in rendered:
            if _is_system_status_outbound(text):
                logger.info(
                    "quota piggyback suppressed system-status text group=%s preview=%r",
                    group_id, text[:120],
                )
                continue
            messages.append(TextMessage(text=text))
        if not messages:
            return

        try:
            with ApiClient(_get_line_config()) as api_client:
                MessagingApi(api_client).reply_message(
                    ReplyMessageRequest(
                        reply_token=reply_token, messages=messages,
                    )
                )
        except Exception as e:
            logger.warning(
                "quota-exhausted piggyback reply failed (pending preserved): %s",
                str(e)[:200],
            )
            return

        # commit：reply 成功才從 pending 移除
        committed_ids = [msg_id for _, msg_id in rendered if msg_id]
        _commit_pending_removal(group_id, committed_ids)

        # log bot's own replies
        import uuid as _uuid
        for m in messages:
            txt = getattr(m, "text", "")[:500]
            uniq = f"piggyback_{_uuid.uuid4().hex[:12]}"
            memory.log_raw_message(group_id, uniq, "__bot__", txt)
        logger.info(
            "quota-exhausted piggyback: drained %d via reply_token group=%s",
            len(rendered), group_id,
        )
    finally:
        slot.release()


def _peek_pending_for_piggyback(
    group_id: str, skip_ids: set[str] | None = None
) -> tuple[str, list[str]] | None:
    """pending 有訊息時，從佇列頭取 1 則處理生回覆，格式化成 piggyback 訊息。

    優先序：
      1. text pending（每次 1 條，慢消化）
      2. 沒 text → file pending（PDF/Office，走 local fallback，每次 1 個）
      3. 沒 file → image pending（每次 1 張，走 media_pipeline.analyze_image
         local vision；50s thread timeout 跟 _handle_image_message 對齊）

    成功回 (格式化字串, message_ids)；失敗回 None，pending 不動。
    caller 必須等 LINE reply 成功後才 commit removal。
    """
    if not _pending_reply_enabled():
        return None
    _drop_stale_pending(group_id)  # D1 TTL: age>7d 進 DLQ，避免 PDF stuck 卡 slot
    pending = _load_pending_explicit()
    items = pending.get(group_id, [])
    if skip_ids:
        items = [it for it in items if it.get("message_id") not in skip_ids]
    if not items:
        return None

    facts = memory.top_facts(group_id)
    context = memory.get_context(group_id)
    pnotes = _get_persona_notes(group_id)

    # ── 1. text pending（user 偏好慢消化）─────────────────────────────────
    _text_items = [
        it
        for it in items
        if it.get("type") == "text" and (it.get("text") or "").strip()
    ]
    batch = _text_items[:1]
    if batch:
        original = "\n".join(it["text"].strip() for it in batch)
        try:
            reply_text = _llm_chat(original, context, facts, pnotes)
        except Exception:
            return None
        if not reply_text:
            return None
        orig_preview = original[:300] + ("…" if len(original) > 300 else "")
        reply_preview = _prepare_outbound_text(
            reply_text, source="legacy_pending_piggyback"
        )
        return (
            f"📬 補回之前漏掉的訊息\n\n原文：\n{orig_preview}\n\n回應：\n{reply_preview}",
            [it.get("message_id") for it in batch if it.get("message_id")],
        )

    # ── 2. file pending（PDF/Office 走 local fallback）──────────────────
    _file_items = [
        it
        for it in items
        if it.get("type") == "file"
        and it.get("media_path")
        and os.path.exists(it["media_path"])
    ]
    file_batch = _file_items[:1]
    if file_batch:
        file_item = file_batch[0]
        file_name = file_item.get("file_name", "unknown")
        media_path = file_item["media_path"]
        try:
            with open(media_path, "rb") as f:
                data = f.read()
        except Exception as e:
            logger.warning("drain pending file read failed: %s", e)
            return None

        reply_text = _drain_pending_file(data, file_name, group_id, context, facts, pnotes)
        if not reply_text or not reply_text.strip():
            return None

        reply_preview = _prepare_outbound_text(
            reply_text, source="legacy_pending_file_piggyback"
        )[:1500]
        return (
            f"📬 補回之前漏掉的檔案 [{file_name}]\n\n回應：\n{reply_preview}",
            [file_item.get("message_id")] if file_item.get("message_id") else [],
        )

    # ── 3. image pending（local 圖像辨識 via media_pipeline.analyze_image）─
    _image_items = [
        it
        for it in items
        if it.get("type") == "image"
        and it.get("media_path")
        and os.path.exists(it["media_path"])
    ]
    img_batch = _image_items[:1]
    if not img_batch:
        return None

    img_item = img_batch[0]
    media_path = img_item["media_path"]
    try:
        with open(media_path, "rb") as f:
            img_data = f.read()
    except Exception as e:
        logger.warning("drain pending image read failed: %s", e)
        return None

    # 跟 _handle_image_message 同 50s thread timeout (reply_token TTL ~1min)
    import threading as _threading
    holder: dict[str, str | None] = {"reply": None}

    def _run():
        try:
            import media_pipeline
            holder["reply"] = media_pipeline.analyze_image(img_data, group_id=group_id)
        except Exception as e:
            logger.warning("piggyback image analyze failed: %s", e)

    t = _threading.Thread(target=_run, daemon=True)
    t.start()
    t.join(timeout=50)
    if t.is_alive():
        logger.info("piggyback image analyze timeout (>50s), leaving pending for next retry")
        return None
    reply_text = holder.get("reply")
    if not reply_text or not reply_text.strip():
        return None

    reply_preview = _prepare_outbound_text(
        reply_text, source="legacy_pending_image_piggyback"
    )[:1500]
    return (
        f"📷 補回之前漏掉的圖片\n\n圖片分析：\n{reply_preview}",
        [img_item.get("message_id")] if img_item.get("message_id") else [],
    )


def _pop_pending_for_piggyback(group_id: str) -> str | None:
    """Legacy pop wrapper. New reply path uses peek-then-confirm."""
    if not _pending_reply_enabled():
        return None
    peeked = _peek_pending_for_piggyback(group_id)
    if not peeked:
        return None
    text, msg_ids = peeked
    _commit_pending_removal(group_id, msg_ids)
    return text


def _drain_pending_file(
    data: bytes,
    file_name: str,
    group_id: str,
    context: list,
    facts: list,
    pnotes: list | None,
) -> str | None:
    """drain pending file → 回 reply_text。對齊 _handle_file_message 的 quota 爆 fallback：

    - image → media_pipeline.analyze_image (mlx-vlm + OCR)
    - PDF → pypdf 抽文字 → _llm_chat (fallback chain)；scanned PDF rasterize 走 vision
    - Office (docx/xlsx/pptx) → _extract_office_text → _llm_chat
    - text/* → decode → _llm_chat

    任何失敗回 None，caller 不移 pending（下次 retry）。
    """
    mime_type = _guess_mime_type(file_name)
    ext = file_name.rsplit(".", 1)[-1].lower() if "." in file_name else ""

    # image
    if mime_type.startswith("image/"):
        try:
            from media_pipeline import analyze_image
            return analyze_image(
                data,
                user_prompt=(
                    "請用繁體中文分析這張圖的內容、主題、可見文字，"
                    "第一句必須是具體判斷或結論，"
                    "不要以「使用者」「我看到」「咪寶」「這張圖」等空話開頭。"
                    f"\n\n[檔名：{file_name}]"
                ),
                group_id=group_id,
            )
        except Exception as e:
            logger.warning("drain image fallback failed: %s", e)
            return None

    # PDF
    if mime_type == "application/pdf":
        try:
            from pypdf import PdfReader
            import io
            reader = PdfReader(io.BytesIO(data))
            total_pages = len(reader.pages)
            content = "\n".join(
                (p.extract_text() or "") for p in reader.pages[:30]
            )
            if content.strip():
                content = content[:_TEXT_CHAR_LIMIT]
                page_note = (
                    f"（PDF 共 {total_pages} 頁，只看前 30 頁）"
                    if total_pages > 30 else ""
                )
                prompt_text = (
                    "請根據以下 PDF 內容做分析回應，用繁體中文，"
                    "第一句必須是具體判斷或結論，"
                    "不要以「使用者」「我看到」「咪寶」「這份檔案」等空話開頭。\n\n"
                    f"--- PDF 內容開始 ---\n{content}\n--- PDF 內容結束 ---\n\n"
                    f"[檔名：{file_name}]{page_note}"
                )
                return _llm_chat(prompt_text, context, facts, pnotes)
            # scanned PDF → rasterize 第一頁
            if len(data) > 50 * 1024 * 1024:
                logger.info("drain pdf rasterize skip: %s 太大", file_name)
                return None
            try:
                import fitz
                doc = fitz.open(stream=data, filetype="pdf")
                fitz_pages = len(doc)
                img_bytes = None
                try:
                    if fitz_pages > 0:
                        pix = doc[0].get_pixmap(dpi=150)
                        try:
                            img_bytes = pix.tobytes("png")
                        finally:
                            pix = None
                finally:
                    doc.close()
                if img_bytes:
                    from media_pipeline import analyze_image
                    desc = analyze_image(
                        img_bytes,
                        user_prompt=(
                            f"這是 scanned PDF「{file_name}」的第一頁。"
                            "請分析內容，第一句必須是具體判斷，"
                            "不要以「使用者」「我看到」「咪寶」開頭。"
                        ),
                        group_id=group_id,
                    )
                    if desc and fitz_pages > 1:
                        return f"{desc}\n\n（scanned PDF 共 {fitz_pages} 頁，只分析第一頁）"
                    return desc
            except Exception as e:
                logger.warning("drain pdf rasterize failed: %s", e)
            return None
        except Exception as e:
            logger.warning("drain pdf fallback failed: %s", e)
            return None

    # Office
    if ext in ("docx", "xlsx", "xls", "pptx"):
        content = _extract_office_text(data, file_name)
        if not content or not content.strip():
            return None
        content = content[:_TEXT_CHAR_LIMIT]
        prompt_text = (
            "請根據以下檔案內容做分析回應，用繁體中文，"
            "第一句必須是具體判斷或結論。\n\n"
            f"--- 內容開始 ---\n{content}\n--- 內容結束 ---\n\n"
            f"[檔名：{file_name}]"
        )
        return _llm_chat(prompt_text, context, facts, pnotes)

    # text/*
    if mime_type.startswith("text/"):
        try:
            content = data.decode("utf-8", errors="replace")
        except Exception:
            return None
        if not content.strip():
            return None
        content = content[:_TEXT_CHAR_LIMIT]
        prompt_text = (
            "請根據以下檔案內容做分析回應，用繁體中文，"
            "第一句必須是具體判斷或結論。\n\n"
            f"--- 內容開始 ---\n{content}\n--- 內容結束 ---\n\n"
            f"[檔名：{file_name}]"
        )
        return _llm_chat(prompt_text, context, facts, pnotes)

    return None


def _reply(
    reply_token: str,
    text: str,
    group_id: str | None = None,
    *,
    allow_push_fallback: bool = True,
) -> None:
    """
    回覆 LINE 訊息。若帶 group_id,成功後會把 bot 的回覆也存進 raw_messages,
    這樣使用者引用 bot 的回覆問後續問題時,能查得到原文。

    若 reply_token 已過期（例如 redelivery）且有 group_id,
    自動 fallback 到 push_message 補送。

    settings.bot_muted=True 時整個函式 short-circuit:
    不 reply、不 push、只把原本要送的 text 寫進 log 方便除錯。
    """
    if not text or not text.strip():
        return
    # Markdown → LINE 純文字，並在 LINE API 呼叫前做最後 factual-safety gate。
    text = _prepare_outbound_text(text, source="reply")
    primary_suppressed = not bool(text and text.strip())
    # LINE 單則訊息上限 5000 字；在截斷前先預留 footer 空間
    footer = _get_quota_footer()
    if not primary_suppressed:
        text = text[: 4900 - len(footer)] + footer
    if not primary_suppressed and _is_system_status_outbound(text):
        logger.info(
            "suppressed system-status LINE reply group=%s preview=%r",
            group_id, text[:120],
        )
        text = ""
        primary_suppressed = True

    # ── Mute 守門 ─────────────────────────────────────────────────────────────
    # 修 bug 期間預設靜音。webhook 照收、classifier/chat 照跑、log 照寫，只是不送 LINE。
    if settings.bot_muted:
        logger.info(
            "[MUTED] would_reply group=%s len=%d preview=%r",
            group_id,
            len(text),
            text[:120],
        )
        return

    # LINE reply_message 上限 5 則。pending reply piggyback 已取消；due
    # reminders 可趁正常 reply 一起送，作為 LINE push quota 429 時的補救路徑。
    # legacy pending branch 僅在 _PENDING_REPLY_ENABLED=True 的測試/rollback 場景會啟用。
    messages_to_send: list = [] if primary_suppressed else [TextMessage(text=text)]
    pending_commit_ids: list[str] = []
    pending_reminders: list[tuple[str, int]] = []  # (event_id, offset) — reply 成功才 mark
    pending_reminder_pushes: list[tuple[int, str]] = []  # (reminder_id, stage) — reply 成功才 mark
    pending_slot = None
    try:
        if group_id:
            # Step 1: legacy pending reply branch。拿同一把 drain lock，且只
            # peek；reply 成功後才 commit removal。
            if (
                messages_to_send
                and _pending_reply_enabled()
                and _load_pending_explicit().get(group_id)
            ):
                pending_slot = _try_acquire_drain_slot(group_id)
                if pending_slot is None:
                    logger.info(
                        "piggyback skip: another drain in progress group=%s",
                        group_id,
                    )
                else:
                    import pending_store as _ps
                    _ps.ensure_message_ids(group_id)
                    if not _quota_exhausted():
                        rendered = _peek_text_pending_for_drain(
                            group_id, max_count=4, deadline_sec=6.0
                        )
                        for pig_text, msg_id in rendered:
                            if _is_system_status_outbound(pig_text):
                                logger.info(
                                    "piggyback suppressed system-status text group=%s preview=%r",
                                    group_id, pig_text[:120],
                                )
                                continue
                            messages_to_send.append(TextMessage(text=pig_text[:5000]))
                            if msg_id:
                                pending_commit_ids.append(msg_id)
                            logger.info(
                                "piggyback peeked text group=%s pig_len=%d",
                                group_id, len(pig_text),
                            )
                            if len(messages_to_send) >= 5:
                                break
                    else:
                        logger.info("piggyback skip: gemini exhausted group=%s", group_id)
                    if len(messages_to_send) == 1:
                        pending_count = len(_load_pending_explicit().get(group_id, []))
                        if pending_count > 0:
                            logger.info(
                                "piggyback skip: pending=%d but no text entry rendered group=%s",
                                pending_count, group_id,
                            )
            if _reminder_reply_piggyback_enabled():
                # Step 2: 剩餘 slot 給 due reminders (legacy quota fallback)
                try:
                    import calendar_db
                    import event_reminder as _er
                    for offset in calendar_db.REMINDER_OFFSETS:
                        if len(messages_to_send) >= 5:
                            break
                        due = calendar_db.list_due_for_reminder(
                            group_id, days_ahead=offset
                        )
                        for e in due:
                            if len(messages_to_send) >= 5:
                                break
                            spec = _er.build_reminder_message_spec(
                                e, offset, allow_mention=True
                            )
                            if spec is None:
                                continue
                            reminder_message = _er.sdk_message_from_spec(spec)
                            if reminder_message is None:
                                continue
                            messages_to_send.append(reminder_message)
                            pending_reminders.append((e["event_id"], offset))
                except Exception as e:
                    logger.warning("reminder piggyback skip: %s", e)

                # Step 3: 再把自然語言 reminder_push 的 due reminder 塞進剩餘 slot。
                try:
                    if len(messages_to_send) < 5:
                        import reminder_push as _rp
                        remaining = 5 - len(messages_to_send)
                        for item in _rp.due_reminders_for_reply(
                            group_id, limit=remaining
                        ):
                            messages_to_send.append(
                                item.get("message")
                                or TextMessage(text=item["text"][:5000])
                            )
                            pending_reminder_pushes.append(
                                (item["reminder_id"], item["stage"])
                            )
                except Exception as e:
                    logger.warning("reminder_push piggyback skip: %s", e)

        resp = None
        if not messages_to_send:
            logger.info("reply suppressed with no piggyback messages group=%s", group_id)
            return
        try:
            with ApiClient(_get_line_config()) as api_client:
                resp = MessagingApi(api_client).reply_message(
                    ReplyMessageRequest(
                        reply_token=reply_token,
                        messages=messages_to_send,
                    )
                )
        except Exception as e:
            logger.warning("reply failed: %s", str(e)[:300])
            if group_id and not _is_definite_reply_token_error(e):
                logger.warning(
                    "reply failure ambiguous; skip fallback push to avoid duplicate group=%s",
                    group_id,
                )
                return
            # reply_token 明確過期 / invalid / 已用過 → fallback 到 push_message
            if group_id:
                if primary_suppressed or not allow_push_fallback or _is_market_quote_outbound(text):
                    logger.info(
                        "reply token invalid; skip fallback push group=%s quote=%s primary_suppressed=%s",
                        group_id,
                        _is_market_quote_outbound(text),
                        primary_suppressed,
                    )
                    return
                try:
                    with ApiClient(_get_line_config()) as api_client:
                        MessagingApi(api_client).push_message(
                            PushMessageRequest(
                                to=group_id,
                                messages=[TextMessage(text=text)],
                            )
                        )
                    logger.info("fallback push_message sent to group=%s", group_id)
                    memory.log_raw_message(
                        group_id, f"push_{int(time.time() * 1000)}", "__bot__", text
                    )
                except Exception as push_err:
                    logger.warning("fallback push also failed: %s", str(push_err)[:300])
            return

        if pending_commit_ids and group_id:
            try:
                removed = _commit_pending_removal(group_id, pending_commit_ids)
                if removed < len(pending_commit_ids):
                    logger.error(
                        "piggyback reply succeeded but committed %d/%d pending group=%s",
                        removed, len(pending_commit_ids), group_id,
                    )
                else:
                    logger.info("piggyback committed %d pending group=%s", removed, group_id)
            except Exception as e:
                logger.exception(
                    "piggyback reply succeeded but pending commit failed group=%s: %s",
                    group_id, e,
                )

        # reply 成功 → mark reminders (peek-then-confirm pattern)
        if pending_reminders:
            try:
                import calendar_db as _cdb
                for ev_id, off in pending_reminders:
                    _cdb.mark_reminded(ev_id, off, group_id)
                logger.info(
                    "reminder piggyback marked: %d reminders group=%s",
                    len(pending_reminders), group_id,
                )
            except Exception as e:
                logger.warning("reminder mark_reminded failed: %s", e)

        # reply 成功 → mark reminder_push reminders (peek-then-confirm pattern)
        if pending_reminder_pushes:
            try:
                import reminder_push as _rp
                marked = _rp.mark_reminders_pushed(pending_reminder_pushes)
                if marked < len(pending_reminder_pushes):
                    logger.error(
                        "reminder_push piggyback marked %d/%d reminders group=%s",
                        marked, len(pending_reminder_pushes), group_id,
                    )
                else:
                    logger.info(
                        "reminder_push piggyback marked: %d reminders group=%s",
                        marked, group_id,
                    )
            except Exception as e:
                logger.warning("reminder_push mark_reminders_pushed failed: %s", e)

        # 把 bot 自己的回覆也記進 raw_messages,供之後 quote-lookup
        if group_id is None:
            return
        sent_messages = getattr(resp, "sent_messages", None) or []
        for idx, sm in enumerate(sent_messages):
            sm_id = getattr(sm, "id", None)
            if sm_id:
                sent_text = text
                if idx < len(messages_to_send):
                    sent_text = getattr(messages_to_send[idx], "text", text)
                memory.log_raw_message(group_id, sm_id, "__bot__", sent_text)
    finally:
        if pending_slot is not None:
            pending_slot.release()


def _download_content(message_id: str) -> bytes:
    """從 LINE 下載 image/video/audio/file 訊息的原始 bytes。"""
    with ApiClient(_get_line_config()) as api_client:
        return bytes(MessagingApiBlob(api_client).get_message_content(message_id))


def _guess_mime_type(file_name: str) -> str:
    mt, _ = mimetypes.guess_type(file_name)
    return mt or "application/octet-stream"
